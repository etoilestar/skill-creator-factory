"""Creator router — file-by-file Skill generation endpoints.

These endpoints decouple the file-creation phase from the main
/api/chat/creator conversation endpoint:

- POST /api/creator/analyze-blueprint  — extract file list from blueprint (no LLM)
- POST /api/creator/init-skill          — create Skill directory structure
- POST /api/creator/generate-file       — SSE: stream single-file content from LLM
- POST /api/creator/write-file          — write generated content to disk
- POST /api/creator/validate-skill      — validate SKILL.md format
- POST /api/creator/package-skill       — package Skill directory into .skill archive
"""

import ast
import base64
import csv
import io
import json
from dataclasses import dataclass, field, replace
import logging
import re
import shlex
import subprocess
import tempfile
import yaml
from pathlib import Path
from typing import Any, Optional

from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field

from ..config import settings
from ..services.blueprint_parser import BlueprintPlan, parse_blueprint
from ..services.skill_plan import SkillPlanEntry, build_skill_plan_entry, capabilities_for_role, command_template_for_entry, default_io_for_role, file_role_classifier, file_type_for_path, language_for_path, runtime_for_language
from ..services.llm_proxy import complete_chat_once, stream_chat
from ..services.model_router import VALIDATOR_TASK, route_creator_file_model, route_model
from ..services.skill_executor import _build_script_runtime_env, run_action
from ..services.artifact_validator import validate_stdout_file_outputs, FileOutputValidationError
from .chat_utils import _get_skill_venv_python, _scan_and_install_python_deps

logger = logging.getLogger(__name__)


def _log_creator_model_usage(
    *,
    phase: str,
    file_path: str,
    route=None,
    model: str | None = None,
    skill_name: str = "",
    attempt: int | None = None,
    actual_model: str | None = None,
    provider: dict | None = None,
    extra: str = "",
) -> None:
    """Emit compact model routing/ack logs for docker logs debugging."""
    expected_model = model or (getattr(route, "model", "") if route is not None else "")
    task = getattr(route, "task", "") if route is not None else ""
    requested_model = getattr(route, "requested_model", None) if route is not None else None
    reason = getattr(route, "reason", "") if route is not None else ""
    matched = None
    if route is not None and actual_model:
        try:
            matched = route.ack(actual_model=actual_model).get("matched")
        except Exception:  # pragma: no cover - defensive logging only
            matched = None
    logger.info(
        "[Creator][model] phase=%s skill=%s file=%s attempt=%s task=%s model=%s requested_model=%s actual_model=%s matched=%s reason=%s provider=%s%s",
        phase,
        skill_name,
        file_path,
        "" if attempt is None else attempt,
        task,
        expected_model,
        requested_model or "",
        actual_model or "",
        "" if matched is None else matched,
        reason,
        provider or {},
        f" extra={extra}" if extra else "",
    )


_SKILL_MD_MARKDOWN_EXECUTION_GUIDE = """

宿主 Markdown 执行说明（写入生成的 SKILL.md 正文时必须保持常见 Markdown 形态）：
- SKILL.md 是普通 Markdown 说明书，只描述做什么、何时使用资源，以及 assistant 在运行时应如何表达动作；不要引入自定义协议章节（例如 `Runtime Contract` JSON）。
- 对纯文本即可完成的任务，明确写“直接回答”，不要要求运行脚本。
- 如果确实需要运行 scripts/ 下的脚本，使用市面常见的 Markdown fenced code block 给出命令示例/模板，例如：
  执行命令：
  ```bash
  python scripts/<script-name> '{"topic":"{{topic}}","keywords":"{{keywords}}"}'
  ```
- 命令示例必须与脚本真实接口一致：脚本读 JSON argv 时，示例就传 JSON；脚本读 stdin 时，正文就说明 stdin 内容。禁止让运行时主模型根据脚本名临时猜 CLI flags。
- 参数映射用普通 Markdown 列表说明，例如 `topic` 从用户输入提取、`keywords` 从用户输入提取、可选参数给出默认值；不要使用单独的 JSON contract。
- 只有 assistant 在 Sandbox 当轮回复中输出的 fenced code block 才会被宿主解析和执行；SKILL.md 中的 block 是运行说明/示例，不会在加载时自动执行。
- 如果需要写文件，用普通 Markdown 说明 assistant 应输出 `写入文件：<path>` 或 `保存到：<path>`，并把完整文件内容放在紧随其后的 fenced code block。
- assistant 不得假装脚本已经执行；必须等待宿主返回 stdout/stderr/observation 后，再基于 observation 生成最终回答。
- 禁止在 SKILL.md 中只写“立即调用 `scripts/...`”这种隐式执行描述；应写成“运行时 assistant 输出以下命令块交由宿主执行”，并给出具体命令示例。
- 如果用户要求使用平台内置模型、图像模型或多模态模型，不要写外部 API key、关键词数据库或假 API；应说明由宿主配置的模型完成相关步骤。任何脚本都必须是有实际功能的实现：要么执行确定性的真实计算/转换/文件处理，要么在需要开放式生成、语义理解、视觉/图像能力时使用宿主已配置的模型能力；模型与认证相关参数由平台运行时注入；生成脚本可按需读取 `IMAGE_MODEL`、`IMAGE_BASE_URL`、`IMAGE_SIZE`、`IMAGE_API_KEY` / `LLM_API_KEY` / `OPENAI_API_KEY` 等环境变量，但不要硬编码这些值，也不需要额外校验它们是否存在。
- 如果需要生成图片，SKILL.md 只描述“使用平台稳定扩散图片生成能力”即可；不要把中文 prompt 翻译、TEXT_MODEL 调用、接口字段解析等平台细节写入创建出来的 Skill 正文。平台运行时会静默完成中文 topic 到英文 Stable Diffusion prompt 的转换。
"""

router = APIRouter(prefix="/api/creator", tags=["creator"])

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# Directories allowed as parents when writing non-SKILL.md files.
_ALLOWED_FOLDERS: frozenset[str] = frozenset({"scripts", "references", "assets"})

# Trailing conversation turns to include in file-generation prompts.
_MAX_HISTORY_TURNS = 6

# Generated files can repair themselves by sending validator/static/trial-run
# failures back to the same routed model before returning content to the frontend.
_MAX_FILE_REPAIR_ATTEMPTS = 10
_SCRIPT_TRIAL_TIMEOUT_SECONDS = 30

# Human-readable language labels indexed by file extension.
_LANG_LABELS: dict[str, str] = {
    ".py":       "Python",
    ".js":       "JavaScript",
    ".mjs":      "JavaScript",
    ".cjs":      "JavaScript",
    ".ts":       "TypeScript",
    ".sh":       "Bash",
    ".bash":     "Bash",
    ".rb":       "Ruby",
    ".go":       "Go",
    ".md":       "Markdown",
    ".yaml":     "YAML",
    ".yml":      "YAML",
    ".json":     "JSON",
    ".toml":     "TOML",
    ".txt":      "Text",
    ".jinja":    "Jinja2 模板",
    ".jinja2":   "Jinja2 模板",
    ".template": "模板文件",
    ".tmpl":     "模板文件",
}


# ---------------------------------------------------------------------------
# Request / Response schemas
# ---------------------------------------------------------------------------

class AnalyzeBlueprintRequest(BaseModel):
    messages: list[dict]
    model: Optional[str] = None


class FileSpecOut(BaseModel):
    path: str
    purpose: str
    required: bool
    can_skip: bool
    file_type: Optional[str] = None
    role: Optional[str] = None
    inputs: list[str] = Field(default_factory=list)
    outputs: list[str] = Field(default_factory=list)
    dependencies: list[str] = Field(default_factory=list)
    required_capabilities: list[str] = Field(default_factory=list)
    forbidden_capabilities: list[str] = Field(default_factory=list)
    reference_files: list[str] = Field(default_factory=list)
    skill_local_references: list[str] = Field(default_factory=list)
    creator_internal_references: list[str] = Field(default_factory=list)
    language: str = "text"
    runtime: str = "none"
    entrypoint: str = ""
    command_template: str = ""
    references: list[str] = Field(default_factory=list)
    low_confidence: bool = False
    confidence: float = 0.0
    reason: str = ""
    heuristic_signals: list[str] = Field(default_factory=list)


class AnalyzeBlueprintResponse(BaseModel):
    skill_name: str
    files: list[FileSpecOut]
    warnings: list[str]


class InitSkillRequest(BaseModel):
    skill_name: str


class InitSkillResponse(BaseModel):
    success: bool
    path: Optional[str] = None
    message: str


class GenerateFileRequest(BaseModel):
    skill_name: str
    file_path: str
    purpose: str
    blueprint_text: str
    conversation_history: list[dict]
    model: Optional[str] = None
    role: Optional[str] = None
    skill_plan_entry: Optional[dict[str, Any]] = None


class WriteFileRequest(BaseModel):
    skill_name: str
    file_path: str
    content: str
    role: Optional[str] = None
    skill_plan_entry: Optional[dict[str, Any]] = None


class WriteFileResponse(BaseModel):
    success: bool
    path: Optional[str] = None
    bytes: int = 0
    message: str


class SkillActionRequest(BaseModel):
    skill_name: str


class SkillActionResponse(BaseModel):
    success: bool
    path: Optional[str] = None
    message: str


class ListFilesRequest(BaseModel):
    skill_name: str


class FileInfo(BaseModel):
    path: str
    is_directory: bool
    size: int = 0


class ListFilesResponse(BaseModel):
    success: bool
    files: list[FileInfo]
    message: str


class InitFromBlueprintRequest(BaseModel):
    skill_name: str
    files: list[FileSpecOut]


class InitFromBlueprintResponse(BaseModel):
    success: bool
    path: Optional[str] = None
    files_created: int = 0
    message: str


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _contains_path_wildcard(file_path: str) -> bool:
    return any(ch in file_path for ch in "*?[]{}")


def _validate_file_path(file_path: str) -> None:
    """Raise HTTP 400 if file_path is outside allowed locations."""
    p = Path(file_path)
    if p.is_absolute() or ".." in p.parts or _contains_path_wildcard(file_path):
        raise HTTPException(
            status_code=400,
            detail=(
                f"非法文件路径: {file_path}。"
                "Creator 只能逐个生成具体文件，不能生成通配符路径。"
            ),
        )

    if file_path == "SKILL.md":
        return

    if not p.parts or p.parts[0] not in _ALLOWED_FOLDERS:
        raise HTTPException(
            status_code=400,
            detail=(
                f"文件路径 '{file_path}' 不合法。"
                f"只允许 SKILL.md 或 scripts/*、references/*、assets/* 下的文件。"
            ),
        )

    filename = p.name
    if not filename or filename.startswith(".") or "\x00" in filename or len(filename) > 255:
        raise HTTPException(status_code=400, detail=f"文件名非法: {filename!r}")


def _validate_skill_name(skill_name: str) -> str:
    """Strip, validate, and return the skill_name or raise HTTP 400."""
    name = skill_name.strip()
    if not name:
        raise HTTPException(status_code=400, detail="skill_name 不能为空。")
    if not re.fullmatch(r"[a-z0-9][a-z0-9-]*", name):
        raise HTTPException(
            status_code=400,
            detail="skill_name 只能包含小写字母、数字和连字符，且必须以字母或数字开头。",
        )
    return name


def _skill_plan_entry_for_file(
    *,
    file_path: str,
    purpose: str = "",
    blueprint_text: str = "",
    role: str | None = None,
    skill_plan_entry: dict[str, Any] | None = None,
) -> SkillPlanEntry:
    """Return the per-file SkillPlan contract used by Creator.

    If the UI passes an explicit role from /analyze-blueprint, keep it.
    Otherwise classify from the concrete file path/purpose, using blueprint text
    only as auxiliary context.
    """
    if skill_plan_entry and skill_plan_entry.get("path") == file_path:
        explicit_role = str(skill_plan_entry.get("role") or role or "").strip()
        allowed_roles = {"text_generator", "image_generator", "composite_generator", "pdf_builder", "docx_builder", "pptx_builder", "html_asset_builder", "asset_builder", "generic_script", "skill_overview", "reference", "asset"}
        if explicit_role in allowed_roles:
            required_capabilities = list(
                skill_plan_entry.get("required_capabilities") or capabilities_for_role(explicit_role)[0]
            )
            if (
                {"text_generation", "image_generation"}.issubset(set(required_capabilities))
                and file_type_for_path(file_path) == "script"
                and explicit_role not in {"pdf_builder", "docx_builder", "pptx_builder", "html_asset_builder", "asset_builder"}
            ):
                explicit_role = "composite_generator"
            inputs = list(skill_plan_entry.get("inputs") or default_io_for_role(explicit_role)[0])
            outputs = list(skill_plan_entry.get("outputs") or default_io_for_role(explicit_role)[1])
            forbidden_capabilities = list(
                skill_plan_entry.get("forbidden_capabilities") or capabilities_for_role(explicit_role)[1]
            )
            forbidden_capabilities = [capability for capability in forbidden_capabilities if capability not in required_capabilities]
            file_type = file_type_for_path(file_path)
            language = str(skill_plan_entry.get("language") or language_for_path(file_path))
            runtime = str(skill_plan_entry.get("runtime") or runtime_for_language(language, file_type))
            return SkillPlanEntry(
                path=file_path,
                file_type=file_type,
                role=explicit_role,  # type: ignore[arg-type]
                purpose=str(skill_plan_entry.get("purpose") or purpose),
                inputs=inputs,
                outputs=outputs,
                dependencies=list(skill_plan_entry.get("dependencies") or []),
                required_capabilities=required_capabilities,
                optional_capabilities=list(skill_plan_entry.get("optional_capabilities") or []),
                allowed_capabilities=list(skill_plan_entry.get("allowed_capabilities") or []),
                forbidden_capabilities=forbidden_capabilities,
                reference_files=[ref for ref in list(skill_plan_entry.get("reference_files") or []) if str(ref).startswith(("references/", "assets/", "scripts/"))],
                skill_local_references=[ref for ref in list(skill_plan_entry.get("skill_local_references") or skill_plan_entry.get("reference_files") or []) if str(ref).startswith(("references/", "assets/", "scripts/"))],
                creator_internal_references=list(skill_plan_entry.get("creator_internal_references") or []),
                language=language,  # type: ignore[arg-type]
                runtime=runtime,  # type: ignore[arg-type]
                entrypoint=str(skill_plan_entry.get("entrypoint") or (file_path if file_type == "script" else "")),
                command_template=str(skill_plan_entry.get("command_template") or (command_template_for_entry(file_path, runtime, inputs) if file_type == "script" else "")),
                required=bool(skill_plan_entry.get("required", True)),
                can_skip=bool(skill_plan_entry.get("can_skip", False)),
                confidence=float(skill_plan_entry.get("confidence") or 1.0),
                reason=str(skill_plan_entry.get("reason") or "explicit SkillPlan entry from UI"),
                heuristic_signals=list(skill_plan_entry.get("heuristic_signals") or []),
            )

    entry = build_skill_plan_entry(
        file_path=file_path,
        purpose=purpose,
        blueprint_summary=blueprint_text[:4000],
    )
    if role and role != entry.role:
        classification = file_role_classifier(
            file_path=file_path,
            purpose=purpose,
            blueprint_summary=blueprint_text[:4000],
        )
        if role in {"text_generator", "image_generator", "composite_generator", "pdf_builder", "docx_builder", "pptx_builder", "html_asset_builder", "asset_builder", "generic_script", "skill_overview", "reference", "asset"}:
            inputs, outputs = default_io_for_role(role)
            required_capabilities, forbidden_capabilities = capabilities_for_role(role)
            forbidden_capabilities = [capability for capability in forbidden_capabilities if capability not in required_capabilities]
            return SkillPlanEntry(
                path=entry.path,
                file_type=entry.file_type,
                role=role,  # type: ignore[arg-type]
                purpose=entry.purpose,
                inputs=inputs,
                outputs=outputs,
                dependencies=entry.dependencies,
                required_capabilities=required_capabilities,
                optional_capabilities=entry.optional_capabilities,
                allowed_capabilities=entry.allowed_capabilities,
                forbidden_capabilities=forbidden_capabilities,
                reference_files=entry.reference_files,
                skill_local_references=entry.skill_local_references,
                creator_internal_references=entry.creator_internal_references,
                language=entry.language,
                runtime=entry.runtime,
                entrypoint=entry.entrypoint,
                command_template=command_template_for_entry(entry.path, entry.runtime, inputs) if entry.file_type == "script" else "",
                required=entry.required,
                can_skip=entry.can_skip,
                confidence=classification.confidence,
                reason=f"explicit role from SkillPlan/UI: {role}",
                heuristic_signals=entry.heuristic_signals,
            )
    return entry


def _extract_first_fenced_block(content: str) -> str | None:
    """Return the first fenced block body from content, or None."""
    lines = content.splitlines(keepends=True)

    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.lstrip()
        match = re.match(r"(`{3,}|~{3,})([^\n`]*)\n?$", stripped.rstrip("\n"))
        if not match:
            i += 1
            continue

        fence = match.group(1)
        fence_char = fence[0]
        fence_len = len(fence)
        code_lines: list[str] = []
        i += 1

        while i < len(lines):
            close_line = lines[i]
            close_stripped = close_line.lstrip()
            close_match = re.match(
                rf"{re.escape(fence_char)}{{{fence_len},}}\s*$",
                close_stripped.rstrip("\n"),
            )
            if close_match:
                return "".join(code_lines).strip()
            code_lines.append(close_line)
            i += 1

        return "".join(code_lines).strip()

    return None


def _extract_target_file_from_bundle(content: str, file_path: str) -> str | None:
    """Extract the requested file when a model returns a multi-file bundle."""
    escaped_path = re.escape(file_path)
    heading_re = re.compile(
        rf"(?im)^\s*#{{1,6}}\s*(?:[^\n`]*?)`?{escaped_path}`?\s*$"
    )

    for match in heading_re.finditer(content):
        section = content[match.end():]
        block = _extract_first_fenced_block(section)
        if block is not None:
            return block

    return None


_MULTI_FILE_MARKER_RE = re.compile(
    r"(?im)^\s*#{1,6}\s*(?:[^\n`]*)(?:SKILL\.md|scripts/|references/|assets/)"
)


_SCRIPT_FAKE_IMPLEMENTATION_RE = re.compile(
    r"placeholder|TODO|your_api_key|api\.example\.com|example\.com|模拟|占位|假装|"
    r"实际使用时|实际开发中|仅为演示|演示目的|空的占位图|纯色图片|ASCII插图|ascii_art|fake",
    re.IGNORECASE,
)
_SKILL_CUSTOM_RUNTIME_CONTRACT_RE = re.compile(r"(?im)^\s*#{1,6}\s*Runtime\s+Contract\s*$")
_HOST_MODEL_CAPABILITY_RE = re.compile(
    r"宿主.{0,12}模型|内置.{0,12}模型|配置.{0,12}模型|文本模型|图像模型|视觉模型|"
    r"多模态|大语言模型|LLM|AI生成|模型生成|调用模型|TEXT_MODEL|IMAGE_MODEL|VISION_MODEL",
    re.IGNORECASE,
)
_CONFIGURED_MODEL_CALL_RE = re.compile(
    r"LLM_BASE_URL|TEXT_MODEL|IMAGE_MODEL|VISION_MODEL|/v1/chat/completions|"
    r"chat/completions|complete_chat_once|stream_chat|openai|"
    r"generate_text_with_llm|generate_stable_diffusion_image|backend\.services\.skill_runtime",
    re.IGNORECASE,
)
_CREATOR_FLOW_LEAK_RE = re.compile(
    r"点击\s*(?:\*\*)?[‘'\"“”]?开始创建[’'\"“”]?(?:\*\*)?|开始生成文件|文件清单预览|确认无误后|"
    r"你也可以在创建后继续编辑内容|确认项列表|系统将自动创建|自动创建以下文件|"
    r"创建文件面板|文件创建面板|若当前无误|已预置|所有路径与命名与蓝图一致|"
    r"不包含任何隐藏逻辑或隐式执行|输出格式符合 Markdown 标准，支持宿主解析",
    re.IGNORECASE,
)
_SKILL_FILE_PATH_RE = re.compile(r"(?<![\w./-])((?:scripts|references|assets)/[A-Za-z0-9_./-]+|SKILL\.md)(?![\w./-])")
_KERNEL_RESOURCE_LEAK_RE = re.compile(
    r"(?<![\w./-])(kernel/references/[A-Za-z0-9_./-]+)(?![\w./-])",
    re.IGNORECASE,
)

_IMAGE_MODEL_USAGE_RE = re.compile(r"IMAGE_MODEL|IMAGE_BASE_URL|/v1/images/generations|images/generations|generate_stable_diffusion_image", re.IGNORECASE)
_DIRECT_IMAGE_API_RE = re.compile(r"IMAGE_BASE_URL|/v1/images/generations|images/generations", re.IGNORECASE)
_PLATFORM_IMAGE_HELPER_RE = re.compile(r"generate_stable_diffusion_image", re.IGNORECASE)
_IMAGE_URL_ONLY_RE = re.compile(r'\[0\]\s*\.get\(\s*[\'"]url[\'"]|\[\s*[\'"]url[\'"]\s*\]', re.IGNORECASE)
_DATA_URI_RE = re.compile(r"data:image/[^;]+;base64", re.IGNORECASE)
_REFERENCE_PLACEHOLDER_RE = re.compile(r"placeholder|TODO|待补充|将要生成|仅为示例|空壳|占位", re.IGNORECASE)


def _reject_custom_skill_md_protocol(content: str) -> None:
    """Reject non-standard runtime protocol sections in generated SKILL.md."""
    if _SKILL_CUSTOM_RUNTIME_CONTRACT_RE.search(content):
        raise ValueError(
            "SKILL.md 不应包含自定义 Runtime Contract JSON 协议；"
            "请使用普通 Markdown 说明和 ```bash 命令示例描述运行时动作。"
        )


def _extract_declared_skill_paths(text: str) -> list[str]:
    """Return normalized Skill package file paths mentioned by a blueprint."""
    seen: set[str] = set()
    paths: list[str] = []
    for raw in _SKILL_FILE_PATH_RE.findall(text or ""):
        path = raw.strip().rstrip("`，,。；;:)）]").replace("\\", "/")
        if path not in seen:
            seen.add(path)
            paths.append(path)
    return paths


def _paths_requiring_skill_md_mentions(blueprint_text: str, *, prefix: str) -> list[str]:
    return [path for path in _extract_declared_skill_paths(blueprint_text) if path.startswith(prefix)]


def _reject_creator_flow_leak(content: str) -> None:
    """Reject Creator UI/workflow text copied into generated SKILL.md."""
    if _CREATOR_FLOW_LEAK_RE.search(content):
        raise ValueError(
            "SKILL.md 包含 Creator 界面流程/确认清单文本（例如“点击开始创建”“确认项列表”“系统将自动创建”），"
            "这是平台创建流程泄露，不属于 Skill 使用说明。请删除这些流程文本，只保留 Skill 的使用说明、资源引用和可执行命令示例。"
        )


@dataclass(frozen=True)
class ContractCheckResult:
    id: str
    passed: bool
    target: str
    message: str
    expected: str
    minimal_edit: str
    matched_paths: list[str] = field(default_factory=list)


class ContractValidationError(ValueError):
    """Validation error carrying structured contract check results."""

    def __init__(self, message: str, results: list[ContractCheckResult]):
        super().__init__(message)
        self.results = results


def _infer_script_input_keys_from_blueprint(script_path: str, blueprint_text: str) -> list[str]:
    """Infer stable JSON argv keys from the blueprint for a script path."""
    lowered = (blueprint_text or "").lower()
    key_candidates = ["topic", "prompt", "text", "keywords"]
    keys = [key for key in key_candidates if key in lowered]
    if not keys:
        keys = ["topic"]
    return keys[:3]


def _script_command_template(script_path: str, blueprint_text: str, entry: SkillPlanEntry | None = None) -> str:
    """Render the command template from SkillPlanEntry, the sole execution contract."""
    if entry is not None and entry.command_template:
        return entry.command_template
    keys = list(entry.inputs) if entry is not None else _infer_script_input_keys_from_blueprint(script_path, blueprint_text)
    if not keys:
        keys = ["payload"]
    payload = json.dumps({key: f"{{{{{key}}}}}" for key in keys}, ensure_ascii=False)
    runtime = entry.runtime if entry is not None else runtime_for_language(language_for_path(script_path), file_type_for_path(script_path))
    if runtime == "node":
        return f"node {script_path} '{payload}'"
    if runtime == "bash":
        return f"bash {script_path} '{payload}'"
    if runtime == "shell":
        return f"sh {script_path} '{payload}'"
    if runtime == "python":
        return f"python {script_path} '{payload}'"
    return f"{script_path} '{payload}'"



def _command_signature(command: str, script_path: str) -> dict[str, Any] | None:
    """Normalize a command block into runner/script/JSON-argv placeholder contract."""
    try:
        parts = shlex.split((command or "").strip())
    except ValueError:
        return None
    expected_script = script_path.replace("\\", "/")
    for idx, part in enumerate(parts):
        normalized_script = part.replace("\\", "/")
        if normalized_script == expected_script or normalized_script.endswith("/" + expected_script):
            if idx + 1 >= len(parts):
                return {
                    "runner": Path(parts[idx - 1]).name if idx > 0 else "",
                    "script_path": expected_script,
                    "keys": set(),
                    "placeholders": {},
                }
            try:
                payload = json.loads(parts[idx + 1])
            except json.JSONDecodeError:
                return None
            if not isinstance(payload, dict):
                return None
            placeholders: dict[str, str] = {}
            for key, value in payload.items():
                if isinstance(value, str):
                    match = re.fullmatch(r"\{\{\s*([A-Za-z_][\w-]*)\s*\}\}", value.strip())
                    placeholders[str(key)] = match.group(1) if match else value.strip()
                else:
                    placeholders[str(key)] = ""
            return {
                "runner": Path(parts[idx - 1]).name if idx > 0 else "",
                "script_path": expected_script,
                "keys": set(str(key) for key in payload.keys()),
                "placeholders": placeholders,
            }
    return None


def _command_template_equivalent(command: str, script_path: str, entry: SkillPlanEntry) -> bool:
    """Compare command blocks by normalized execution contract, not raw bytes."""
    command_sig = _command_signature(command, script_path)
    template_sig = _command_signature(_script_command_template(script_path, "", entry), script_path)
    if not command_sig or not template_sig:
        return False
    if command_sig["runner"] != template_sig["runner"]:
        return False
    if command_sig["script_path"] != template_sig["script_path"]:
        return False
    if command_sig["keys"] != template_sig["keys"]:
        return False
    return command_sig["placeholders"] == template_sig["placeholders"]

def _command_payload_keys(command: str, script_path: str) -> set[str] | None:
    """Return JSON argv keys passed to script_path, or None if unparsable/non-JSON."""
    try:
        parts = shlex.split(command)
    except ValueError:
        return None
    for idx, part in enumerate(parts):
        normalized = part.replace("\\", "/")
        if normalized == script_path or normalized.endswith("/" + script_path):
            if idx + 1 >= len(parts):
                return set()
            try:
                payload = json.loads(parts[idx + 1])
            except json.JSONDecodeError:
                return None
            if not isinstance(payload, dict):
                return None
            return {str(key) for key in payload.keys()}
    return None



def _command_runtime_matches(command: str, script_path: str, entry: SkillPlanEntry) -> bool:
    try:
        parts = shlex.split(command)
    except ValueError:
        return False
    for idx, part in enumerate(parts):
        normalized = part.replace("\\", "/")
        if normalized == script_path or normalized.endswith("/" + script_path):
            runner = parts[idx - 1] if idx > 0 else ""
            if entry.runtime == "python":
                return Path(runner).name.startswith("python")
            if entry.runtime == "node":
                return Path(runner).name == "node"
            if entry.runtime == "bash":
                return Path(runner).name in {"bash", "sh"}
            if entry.runtime == "shell":
                return Path(runner).name in {"sh", "bash"}
            return True
    return False

def _check_command_block_contract(script_path: str, commands: list[str], entry: SkillPlanEntry) -> list[ContractCheckResult]:
    """Validate Markdown command examples against the SkillPlan input contract."""
    required = set(entry.inputs or ["payload"])
    results: list[ContractCheckResult] = []
    for idx, command in enumerate(commands, start=1):
        command = command.strip()
        keys = _command_payload_keys(command, script_path)
        target = f"{script_path}#command-{idx}"
        template_matches = _command_template_equivalent(command, script_path, entry)
        results.append(ContractCheckResult(
            id="command_block.command_template.equivalent",
            passed=template_matches,
            target=target,
            message=(
                "命令块与 SkillPlan.command_template 规范化等价。" if template_matches
                else "命令块必须复用 SkillPlan.command_template 的执行合同：runner、script path、JSON argv keys 和 placeholder 映射必须一致。"
            ),
            expected=f"command_template: {_script_command_template(script_path, '', entry)}",
            minimal_edit=f"将命令调整为与模板等价，推荐整行替换为：{_script_command_template(script_path, '', entry)}",
        ))
        runtime_matches = _command_runtime_matches(command, script_path, entry)
        results.append(ContractCheckResult(
            id="command_block.runtime.matches_skillplan",
            passed=runtime_matches,
            target=target,
            message=(
                "命令块 runner 与 SkillPlan runtime 一致。" if runtime_matches
                else f"命令块必须按 SkillPlan.runtime={entry.runtime} 调用 {script_path}。"
            ),
            expected=f"按 runtime 调用：{_script_command_template(script_path, '', entry)}",
            minimal_edit=f"将命令改为：{_script_command_template(script_path, '', entry)}",
        ))
        results.append(ContractCheckResult(
            id="command_block.json_argv.parseable",
            passed=keys is not None,
            target=target,
            message=(
                "命令块使用可解析 JSON argv。" if keys is not None
                else f"{script_path} 命令块必须在脚本路径后传入一个 JSON object argv。"
            ),
            expected=f"命令形如：{_script_command_template(script_path, '', entry)}",
            minimal_edit=f"将命令改为：{_script_command_template(script_path, '', entry)}",
        ))
        if keys is None:
            continue
        missing = sorted(required - keys)
        extra = sorted(keys - required)
        results.append(ContractCheckResult(
            id="command_block.skillplan_inputs.exact",
            passed=not missing and not extra,
            target=target,
            message=(
                "命令块 JSON keys 与 SkillPlan inputs 完全一致。" if not missing and not extra
                else f"命令块 JSON keys 与 SkillPlan inputs 不一致；missing={missing} extra={extra}。"
            ),
            expected=f"JSON argv keys 必须且只能是：{', '.join(sorted(required))}",
            minimal_edit=f"将 JSON argv 调整为：{json.dumps({key: '{{' + key + '}}' for key in sorted(required)}, ensure_ascii=False)}",
        ))
    return results


def _build_skill_md_contract_text(blueprint_text: str) -> str:
    """Build the explicit SKILL.md contract injected before generation/repair."""
    script_paths = _paths_requiring_skill_md_mentions(blueprint_text, prefix="scripts/")
    reference_paths = _paths_requiring_skill_md_mentions(blueprint_text, prefix="references/")

    lines = [
        "必须满足以下 SKILL.md 合同，逐项覆盖：",
        "A. frontmatter:",
        "- 必须以 --- 开始。",
        "- 必须包含 name 和 description。",
        "- 必须用 --- 关闭 frontmatter。",
        "B. 复合任务编排:",
        "- SKILL.md 必须作为总流程/编排说明，描述执行顺序、数据流、每步预期输出，以及何时读取 references/assets。",
        "- 对复合任务，SKILL.md 只做流程总览，不把子任务详细规则写满；详细规范放入对应 references/*.md。",
        "C. scripts 命令块:",
    ]
    if script_paths:
        for script_path in script_paths:
            entry = _skill_plan_entry_for_file(file_path=script_path, blueprint_text=blueprint_text)
            keys = ", ".join(entry.inputs or ["payload"])
            lines.extend([
                "- SKILL.md 可以直接包含普通 Markdown ```bash fenced code block；若不直接包含，必须明确说明由对应 reference 定义命令或执行步骤。",
                f"- 直接命令 block 内必须出现精确路径：{script_path}",
                f"- 命令必须使用 JSON argv，JSON keys 必须且只能来自 SkillPlan inputs：{keys}。",
                f"- 推荐命令模板：{_script_command_template(script_path, blueprint_text, entry)}",
            ])
    else:
        lines.append("- 蓝图没有 scripts/，不要强行写脚本命令。")

    lines.append("D. references 引用:")
    if reference_paths:
        for reference_path in reference_paths:
            lines.append(f"- 必须在正文中出现并说明用途：{reference_path}")
    else:
        lines.append("- 蓝图没有 references/，不要强行编造参考资料。")

    lines.extend([
        "E. 禁止项:",
        "- 不要包含 Runtime Contract JSON。",
        "- 不要包含 Creator 创建流程、确认清单、点击开始创建、系统将自动创建等平台流程文案。",
    ])
    return "\n".join(lines)



def _declared_skill_paths_from_blueprint(blueprint_text: str) -> set[str]:
    paths: set[str] = set()
    for prefix in ("scripts/", "references/", "assets/"):
        paths.update(_paths_requiring_skill_md_mentions(blueprint_text, prefix=prefix))
    return paths


def _skill_local_paths_in_markdown(content: str) -> set[str]:
    return {match.group(1).strip() for match in _SKILL_FILE_PATH_RE.finditer(content or "")}


def _kernel_resource_leak_paths(content: str) -> list[str]:
    """Return explicit kernel/references paths mentioned in final Skill text."""
    seen: set[str] = set()
    paths: list[str] = []
    for match in _KERNEL_RESOURCE_LEAK_RE.finditer(content or ""):
        path = match.group(1).rstrip("`，,。；;:)）]")
        if path not in seen:
            seen.add(path)
            paths.append(path)
    return paths


def _normalize_similarity_text(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "")).strip().lower()


def _kernel_reference_content_copy_paths(content: str) -> list[str]:
    """Detect large verbatim copying from kernel references without banning same names."""
    import difflib

    candidate = _normalize_similarity_text(content)
    if len(candidate) < 500:
        return []
    matches: list[str] = []
    kernel_refs = settings.kernel_path / "references"
    if not kernel_refs.is_dir():
        return []
    for ref in sorted(kernel_refs.glob("*.md")):
        try:
            kernel_text = _normalize_similarity_text(ref.read_text(encoding="utf-8"))
        except OSError:
            continue
        if len(kernel_text) < 500:
            continue
        matcher = difflib.SequenceMatcher(None, candidate, kernel_text, autojunk=False)
        longest = max((block.size for block in matcher.get_matching_blocks()), default=0)
        # A contiguous 500+ char copy is almost certainly a leaked reference;
        # for shorter kernel refs, also catch near-whole-file copies.
        copied_ratio = longest / max(1, min(len(candidate), len(kernel_text)))
        if longest >= 500 or (longest >= 300 and copied_ratio >= 0.60):
            matches.append(f"kernel/references/{ref.name}")
    return matches


def _existing_skill_local_paths_for_skill(skill_name: str) -> set[str]:
    skill_dir = settings.skills_path / skill_name
    paths: set[str] = set()
    if not skill_dir.exists():
        return paths
    for folder in ("scripts", "references", "assets"):
        folder_path = skill_dir / folder
        if not folder_path.is_dir():
            continue
        for child in folder_path.rglob("*"):
            if child.is_file():
                paths.add(child.relative_to(skill_dir).as_posix())
    return paths


def _check_skill_md_contract(content: str, blueprint_text: str) -> list[ContractCheckResult]:
    """Return structured SKILL.md contract checks for generation and repair."""
    stripped = content.strip()
    results: list[ContractCheckResult] = []

    has_frontmatter = bool(re.match(r"^---\nname: [^\n]+\ndescription: [^\n]+\n---\n", stripped))
    results.append(ContractCheckResult(
        id="skill_md.frontmatter",
        passed=has_frontmatter,
        target="SKILL.md",
        message=(
            "SKILL.md frontmatter 合格。" if has_frontmatter
            else "SKILL.md 必须以 YAML frontmatter 开始，且包含 name 和 description：--- / name: ... / description: ... / ---。"
        ),
        expected="文件以 --- / name: ... / description: ... / --- 开头。",
        minimal_edit="修正文件开头的 YAML frontmatter，不要改动正文其它已通过内容。",
    ))

    has_runtime_contract = bool(_SKILL_CUSTOM_RUNTIME_CONTRACT_RE.search(content))
    results.append(ContractCheckResult(
        id="skill_md.forbidden_runtime_contract",
        passed=not has_runtime_contract,
        target="SKILL.md",
        message=(
            "未包含自定义 Runtime Contract JSON 协议。" if not has_runtime_contract
            else "SKILL.md 不应包含自定义 Runtime Contract JSON 协议；请使用普通 Markdown 说明和 ```bash 命令示例描述运行时动作。"
        ),
        expected="不要包含 Runtime Contract JSON。",
        minimal_edit="删除 Runtime Contract JSON/协议小节，改为普通 Markdown 说明。",
    ))

    has_creator_flow = bool(_CREATOR_FLOW_LEAK_RE.search(content))
    results.append(ContractCheckResult(
        id="skill_md.forbidden_creator_flow",
        passed=not has_creator_flow,
        target="SKILL.md",
        message=(
            "未包含 Creator 界面流程文案。" if not has_creator_flow
            else "SKILL.md 包含 Creator 界面流程/确认清单文本（例如“点击开始创建”“确认项列表”“系统将自动创建”），这是平台创建流程泄露，不属于 Skill 使用说明。"
        ),
        expected="不要包含 Creator 创建流程、确认清单、点击开始创建等平台流程文案。",
        minimal_edit="删除 Creator UI/确认清单/点击开始创建相关文案，只保留 Skill 使用说明。",
    ))

    mentioned_skill_paths = _skill_local_paths_in_markdown(content)
    declared_skill_paths = _declared_skill_paths_from_blueprint(blueprint_text)
    missing_local_paths = sorted(path for path in mentioned_skill_paths if path not in declared_skill_paths)
    kernel_leak_paths = _kernel_resource_leak_paths(content)
    results.append(ContractCheckResult(
        id="skill_md.resource.no_kernel_leak",
        passed=not kernel_leak_paths,
        target="SKILL.md",
        message=(
            "SKILL.md 未引用 Creator 内部 kernel resources。"
            if not kernel_leak_paths
            else "SKILL.md 显式引用了 Creator 内部 kernel resources：" + ", ".join(kernel_leak_paths)
        ),
        expected="最终业务 SKILL.md 只能引用本轮生成或当前 skill 目录真实存在的 resources；只有显式 kernel/references/... 路径会被判定为 kernel 路径泄露。",
        minimal_edit="删除 kernel/references/... 内部 Creator 资源引用；如果业务 Skill 需要同名 reference，请创建并引用 skill-local references/*.md。",
        matched_paths=kernel_leak_paths,
    ))
    kernel_copy_paths = _kernel_reference_content_copy_paths(content)
    results.append(ContractCheckResult(
        id="skill_md.resource.no_kernel_content_copy",
        passed=not kernel_copy_paths,
        target="SKILL.md",
        message=(
            "SKILL.md 未大段复制 Creator kernel reference 内容。"
            if not kernel_copy_paths
            else "SKILL.md 大段复制了 Creator kernel reference 内容：" + ", ".join(kernel_copy_paths)
        ),
        expected="最终业务 SKILL.md 不得大段复制 kernel/references 中的 Creator 内部说明；同名本地资源可引用，但内容应是业务 Skill 自有说明。",
        minimal_edit="删除复制的 kernel 内部说明，改写为面向该业务 Skill 的简洁使用说明和本地资源引用。",
        matched_paths=kernel_copy_paths,
    ))
    results.append(ContractCheckResult(
        id="skill_md.resource.local_declared",
        passed=not missing_local_paths,
        target="SKILL.md",
        message=(
            "SKILL.md 只引用本轮声明/已存在的 skill-local resources。"
            if not missing_local_paths
            else "SKILL.md 引用了未在本轮生成或当前 Skill 中不存在的资源：" + ", ".join(missing_local_paths)
        ),
        expected="SKILL.md 中的 scripts/references/assets 路径必须是当前 Skill 内真实存在或本轮生成的资源；裸 references/*.md 只做本地资源存在性校验，不因文件名与 kernel reference 同名而失败。",
        minimal_edit="删除未声明/不存在的 references/assets/scripts 引用；如果确实需要该资源，请先在蓝图中加入并生成对应文件。",
        matched_paths=missing_local_paths,
    ))

    for script_path in _paths_requiring_skill_md_mentions(blueprint_text, prefix="scripts/"):
        commands = _extract_script_command_templates(content, script_path)
        passed = bool(commands)
        entry = _skill_plan_entry_for_file(file_path=script_path, blueprint_text=blueprint_text)
        if entry.role == "generic_script" and entry.inputs == ["payload"]:
            inferred_inputs = _infer_script_input_keys_from_blueprint(script_path, blueprint_text)
            if commands:
                command_keys: set[str] = set()
                for command in commands:
                    keys = _command_payload_keys(command, script_path)
                    if keys:
                        command_keys.update(keys)
                if command_keys:
                    inferred_inputs = sorted(command_keys)
            entry = replace(entry, inputs=inferred_inputs, command_template=commands[0].strip() if commands else command_template_for_entry(script_path, entry.runtime, inferred_inputs))
        template = _script_command_template(script_path, blueprint_text, entry)
        reference_delegates_execution = bool(reference_path_candidates := _paths_requiring_skill_md_mentions(blueprint_text, prefix="references/")) and any(ref in content for ref in reference_path_candidates) and bool(re.search(r"命令|执行|步骤|command|execute|steps", content, re.I))
        results.append(ContractCheckResult(
            id="skill_md.script_command.exists",
            passed=passed or reference_delegates_execution,
            target=script_path,
            message=(
                f"SKILL.md 已包含调用 {script_path} 的可执行 Markdown 命令块。" if passed
                else (f"SKILL.md 未直接写 {script_path} 命令块，但已显式委托 references 定义执行步骤。" if reference_delegates_execution else f"SKILL.md 缺少调用 {script_path} 的可执行 Markdown 命令块，且未显式说明由 references 定义命令/执行步骤。")
            ),
            expected=f"SKILL.md 直接包含命令块，或明确说明由 references/*.md 定义执行命令；直接命令推荐：{template}",
            minimal_edit=f"在执行/运行脚本小节加入命令块：```bash\n{template}\n```，或明确写明读取 reference 中的命令模板。",
        ))
        if commands:
            results.extend(_check_command_block_contract(script_path, commands, entry))

    for reference_path in _paths_requiring_skill_md_mentions(blueprint_text, prefix="references/"):
        passed = reference_path in content
        results.append(ContractCheckResult(
            id="skill_md.reference.mentioned",
            passed=passed,
            target=reference_path,
            message=(
                f"SKILL.md 已引用参考资料 {reference_path}。" if passed
                else f"SKILL.md 缺少对参考资料 {reference_path} 的引用。请在“参考资料/资源”小节用普通 Markdown 明确说明何时读取该 reference。"
            ),
            expected=f"正文中逐字出现 {reference_path} 并说明用途/何时读取。",
            minimal_edit=f"在参考资料/资源小节加入 `{reference_path}` 及其用途说明。",
        ))

    return results


def _format_contract_checks(results: list[ContractCheckResult], *, passed: bool) -> str:
    selected = [result for result in results if result.passed is passed]
    if not selected:
        return "- 无"
    lines: list[str] = []
    for result in selected:
        matched = f"\n  matched_paths: {', '.join(result.matched_paths)}" if result.matched_paths else ""
        lines.append(
            f"- {result.id} target={result.target}: {result.message}\n"
            f"  expected: {result.expected}\n"
            f"  minimal_edit: {result.minimal_edit}"
            f"{matched}"
        )
    return "\n".join(lines)


def _format_contract_failures(results: list[ContractCheckResult]) -> str:
    failed = [result for result in results if not result.passed]
    if not failed:
        return ""
    return (
        "SKILL.md contract 未通过：\n"
        + _format_contract_checks(failed, passed=False)
    )


def _validate_skill_md_contract(content: str, blueprint_text: str) -> None:
    """Validate generated SKILL.md against blueprint-declared resources."""
    results = _check_skill_md_contract(content, blueprint_text)
    failed = [result for result in results if not result.passed]
    if failed:
        raise ContractValidationError(_format_contract_failures(results), results)





def _strip_orphan_trailing_fence(content: str) -> str:
    """Remove isolated Markdown fence markers at file boundaries.

    This is intentionally narrower than generic fence stripping: it deletes only
    standalone trailing ```/~~~ lines left by model output, plus an optional
    standalone opening fence when no matching closing fence remains.
    """
    lines = content.strip().splitlines()
    changed = False
    while lines and re.fullmatch(r"\s*(`{3,}|~{3,})\s*", lines[-1]):
        lines.pop()
        changed = True
    if lines and re.fullmatch(r"\s*(`{3,}|~{3,})[A-Za-z0-9_-]*\s*", lines[0]):
        body = "\n".join(lines[1:])
        if "```" not in body and "~~~" not in body:
            lines = lines[1:]
            changed = True
    return ("\n".join(lines).strip() if changed else content.strip())


def _build_script_file_contract_text(
    file_path: str,
    blueprint_text: str,
    *,
    purpose: str = "",
    role: str | None = None,
    skill_plan_entry: dict[str, Any] | None = None,
) -> str:
    entry = _skill_plan_entry_for_file(
        file_path=file_path,
        purpose=purpose,
        blueprint_text=blueprint_text,
        role=role,
        skill_plan_entry=skill_plan_entry,
    )
    keys = ", ".join(entry.inputs or ["payload"])
    lines = [
        f"必须满足以下脚本文件合同：{file_path}",
        f"SkillPlan role: {entry.role}",
        f"file_type: {entry.file_type}",
        f"language: {entry.language}",
        f"runtime: {entry.runtime}",
        f"entrypoint: {entry.entrypoint or file_path}",
        f"command_template: {_script_command_template(file_path, blueprint_text, entry)}",
        f"inputs: {', '.join(entry.inputs) or 'payload'}",
        f"outputs: {', '.join(entry.outputs)}",
        f"required_capabilities: {', '.join(entry.required_capabilities) or 'none'}",
        f"forbidden_capabilities: {', '.join(entry.forbidden_capabilities) or 'none'}",
        "A. 输出形态:",
        "- 只输出单个脚本源码本身，不要 Markdown fence、说明文字、写入文件标签或多文件包。",
        "- Python 脚本必须能通过 ast.parse 语法检查。",
        "B. 参数接口:",
        "- 默认使用 JSON argv 接口：按脚本 runtime 读取 JSON argv 并解析。",
        f"- 必须实际使用用户输入 keys：{keys}。",
        f"- 与 SKILL.md/reference 命令模板保持一致；推荐命令：{_script_command_template(file_path, blueprint_text, entry)}",
        "C. 角色输出合同:",
    ]
    if entry.role == "text_generator":
        lines.extend([
            "- stdout 必须输出 JSON object，至少包含一个非空字段；字段名由 SKILL.md 数据流决定，不强制 text/story_text。",
            "- 能力边界由 required_capabilities/forbidden_capabilities 决定；若 required_capabilities 额外包含 image_generation，则必须调用图片 helper。",
        ])
    elif entry.role == "composite_generator":
        lines.extend([
            "- stdout 必须输出 JSON object，至少包含一个非空字段；字段名由 SKILL.md 数据流决定。",
            "- 必须同时调用 generate_text_with_llm 与 generate_stable_diffusion_image；不得用固定 f-string/template-only 文本或占位图片替代。",
            "- 保证 stdout 字段能被 SKILL.md 后续命令的 {{变量}} 消费，形成闭环。",
        ])
    elif entry.role == "image_generator":
        lines.extend([
            "- stdout 必须输出 JSON object，至少包含一个非空字段；字段名由 SKILL.md 数据流决定，不强制 image_paths/images。",
            "- 必须调用平台 Stable Diffusion helper，不要直接调用 /v1/images/generations。",
            "- 不要在脚本里写中文 prompt 翻译逻辑；平台 helper 会处理 Stable Diffusion prompt 适配。",
            "- 必须保留 result = generate_stable_diffusion_image(desc) 的调用骨架；禁止 image_path = generate_stable_diffusion_image(...)。",
        ])
    elif entry.role == "pdf_builder":
        lines.extend([
            "- stdout 必须输出 JSON object，并在任意业务字段中返回真实生成文件路径。",
            "- forbidden_capabilities 生效：PDF 构建脚本不得生成图片或调用 generate_stable_diffusion_image。",
            "- 只能消费已有 story_text/text/image_paths/previous_stdout/template/assets 并构建 PDF/文件。",
        ])
    elif entry.role == "docx_builder":
        lines.extend([
            "- stdout 必须输出 JSON object，且 docx_path 或 file_paths 至少一个指向真实生成文件。",
            "- 只能消费上一步 stdout JSON、story_text/text、image_paths 或模板；不得生成新文本或图片。",
        ])
    elif entry.role == "pptx_builder":
        lines.extend([
            "- stdout 必须输出 JSON object，且 pptx_path 或 file_paths 至少一个指向真实生成文件。",
            "- 只能消费上一步 stdout JSON、story_text/text、image_paths 或模板；不得生成新文本或图片。",
        ])
    else:
        lines.extend([
            "- stdout 必须输出 JSON object，不要混入调试说明。",
            "- 根据角色职责返回 text、file_paths 或其它明确结果字段。",
        ])
    lines.extend([
        "E. 禁止项:",
        "- 禁止 placeholder/mock/fake API/固定模板冒充真实能力。",
        "- 只能实现本 role 的职责；复合任务由 SKILL.md 编排多个 scripts/references/assets 完成。",
    ])
    return "\n".join(lines)


def _build_reference_file_contract_text(file_path: str, purpose: str, blueprint_text: str) -> str:
    script_paths = _paths_requiring_skill_md_mentions(blueprint_text, prefix="scripts/")
    script_lines: list[str] = []
    for script_path in script_paths:
        entry = _skill_plan_entry_for_file(file_path=script_path, blueprint_text=blueprint_text)
        if file_path in entry.reference_files or not entry.reference_files:
            script_lines.extend([
                f"- 本 reference 默认只写规范、风格、正例、反例和质量标准；不要重新定义 {script_path} 的 role/inputs/outputs/capabilities/command_template。",
                f"- 如果 SKILL.md 已包含 {script_path} 的可执行命令块，本 reference 不要再写命令块。",
                f"- 若 reference 必须提供命令块，只能复用 SkillPlan.command_template 的等价执行合同：```bash\n{_script_command_template(script_path, blueprint_text, entry)}\n```",
                f"- 禁止把上述正确命令、JSON keys（{', '.join(entry.inputs or ['payload'])}）或 role={entry.role} 写成反例。",
            ])
    if not script_lines:
        script_lines.append("- 本 reference 对应一个独立子任务/模块；必须写清 inputs、outputs、执行步骤、约束和示例。")
    return "\n".join([
        f"必须满足以下参考资料文件合同：{file_path}",
        "A. 输出形态:",
        "- 只输出该 reference 的 Markdown 文档内容，不要写入文件标签、Creator 流程说明或多文件包。",
        "- 可以包含普通 Markdown 标题/列表/示例；如确实需要代码示例，可以包含文档内部 fenced block。",
        "B. 内容职责:",
        f"- 职责说明：{purpose or '根据蓝图提供可操作参考资料'}",
        "- 每个 reference 只对应一个子任务/模块，不要把整个 Skill 包打包到一个 reference。",
        "- 默认只写写作规范、风格要求、正例、反例、质量标准；如果 SKILL.md 已有可执行命令块，reference 不要重复写命令块。",
        "- 不要重新定义 role / inputs / outputs / capabilities / command_template。",
        *script_lines,
        "- 内容必须是有实际指导价值的参考资料，不是对‘将要生成参考资料’的再描述。",
        "- 必须包含任务规范/步骤、示例、反例、约束/禁止项等章节，且正文长度足以指导子任务。",
        "C. 禁止项:",
        "- 不要包含 Creator 创建流程、确认清单、点击开始创建等平台流程文案。",
        "- 不要包含其它 SKILL.md/scripts/assets/references 文件的打包内容。",
        "- 不要包含 placeholder/TODO/待补充等占位文本。",
    ])


def _build_asset_file_contract_text(file_path: str, purpose: str) -> str:
    return "\n".join([
        f"必须满足以下 asset 文件合同：{file_path}",
        "A. 输出形态:",
        "- 只输出当前 asset 文件内容，不要写入文件标签、说明文字或多文件包。",
        "- 文件必须非空；JSON 资源必须可被 json.loads 解析。",
        "B. 内容职责:",
        f"- 职责说明：{purpose or '根据蓝图提供模板或静态资源'}",
        "C. 禁止项:",
        "- asset 是模板或静态资源，不得包含运行时代码、图片生成调用或 Creator 创建流程文案。",
    ])


def _build_generated_file_contract_text(
    file_path: str,
    blueprint_text: str,
    purpose: str = "",
    *,
    role: str | None = None,
    skill_plan_entry: dict[str, Any] | None = None,
) -> str:
    if file_path == "SKILL.md":
        return _build_skill_md_contract_text(blueprint_text)
    if file_path.startswith("scripts/"):
        return _build_script_file_contract_text(file_path, blueprint_text, purpose=purpose, role=role, skill_plan_entry=skill_plan_entry)
    if file_path.startswith("references/"):
        return _build_reference_file_contract_text(file_path, purpose, blueprint_text)
    if file_path.startswith("assets/"):
        return _build_asset_file_contract_text(file_path, purpose)
    return ""


def _check_reference_file_contract(file_path: str, content: str, purpose: str = "") -> list[ContractCheckResult]:
    stripped = content.strip()
    results = [
        ContractCheckResult(
            id="reference.not_empty",
            passed=bool(stripped),
            target=file_path,
            message=("参考资料内容非空。" if stripped else f"{file_path} 参考资料内容为空。"),
            expected="输出该 reference 的 Markdown 文档内容。",
            minimal_edit="补充有实际指导价值的 Markdown 参考资料正文。",
        ),
        ContractCheckResult(
            id="reference.no_creator_flow",
            passed=not bool(_CREATOR_FLOW_LEAK_RE.search(content)),
            target=file_path,
            message=(
                "未包含 Creator 创建流程文案。"
                if not _CREATOR_FLOW_LEAK_RE.search(content)
                else f"{file_path} 包含 Creator 创建流程/确认清单/点击开始创建等平台流程文案。"
            ),
            expected="不要包含 Creator 创建流程、确认清单、点击开始创建等平台流程文案。",
            minimal_edit="删除平台创建流程文案，只保留参考资料正文。",
        ),
        ContractCheckResult(
            id="reference.single_file",
            passed=not bool(_MULTI_FILE_MARKER_RE.search(content)) and not bool(re.search(r"(?im)^\s*写入文件[:：]", content)),
            target=file_path,
            message=(
                "参考资料是单文件内容。"
                if not _MULTI_FILE_MARKER_RE.search(content) and not re.search(r"(?im)^\s*写入文件[:：]", content)
                else f"{file_path} 包含多文件包、其它文件路径标题或写入文件标签。"
            ),
            expected="只输出当前 reference 文件内容，不要包含 SKILL.md/scripts/assets/references 的多文件包。",
            minimal_edit="删除其它文件内容、路径标题和写入文件标签，只保留当前参考资料正文。",
        ),
    ]
    min_chars = 120
    has_min_length = len(stripped) >= min_chars
    results.append(ContractCheckResult(
        id="reference.min_quality_length",
        passed=has_min_length,
        target=file_path,
        message=(
            "参考资料长度满足最低质量要求。"
            if has_min_length
            else f"{file_path} 内容过短，无法作为子任务参考资料。"
        ),
        expected=f"至少 {min_chars} 个字符，包含可执行的任务规则、示例和约束。",
        minimal_edit="扩充为任务专用参考资料，加入步骤、格式要求、示例、反例和约束。",
    ))
    required_sections = {
        "rules": bool(re.search(r"(?im)^#{1,3}.*(规范|规则|步骤|流程|要求|Rules|Steps)", stripped)),
        "examples": bool(re.search(r"(?im)^#{1,3}.*(示例|例子|Examples?)", stripped)),
        "anti_examples": bool(re.search(r"(?im)^#{1,3}.*(反例|错误示例|Anti[- ]?examples?)", stripped)),
        "constraints": bool(re.search(r"(?im)^#{1,3}.*(约束|限制|禁止|Constraints?)", stripped)),
    }
    sections_ok = all(required_sections.values())
    missing_sections = [name for name, present in required_sections.items() if not present]
    results.append(ContractCheckResult(
        id="reference.required_sections",
        passed=sections_ok,
        target=file_path,
        message=(
            "参考资料包含规范/示例/反例/约束章节。"
            if sections_ok
            else f"{file_path} 缺少必要章节：{', '.join(missing_sections)}。"
        ),
        expected="包含规范/步骤、示例、反例、约束/禁止项章节。",
        minimal_edit="补齐 Markdown 标题章节：## 规范、## 示例、## 反例、## 约束。",
    ))

    role_sections = {
        "io": bool(re.search(r"(?im)^#{1,3}.*(输入|输出|Inputs?|Outputs?)", stripped)),
        "execution": bool(re.search(r"(?im)^#{1,3}.*(执行|命令|步骤|Execution|Commands?)", stripped)) or "```bash" in stripped,
        "role_constraints": bool(re.search(r"(?im)^#{1,3}.*(角色|能力|禁止|Role|Capabilities?)", stripped)),
    }
    role_section_required = bool(re.search(r"子任务|脚本|script|SkillPlan|text_generator|image_generator|composite_generator|pdf_builder|docx_builder|pptx_builder|html_asset_builder|asset_builder|generic_script|命令模板|执行参考", purpose, re.I))
    role_sections_ok = (not role_section_required) or all(role_sections.values())
    missing_role_sections = [name for name, present in role_sections.items() if not present]
    results.append(ContractCheckResult(
        id="reference.subtask_contract_sections",
        passed=role_sections_ok,
        target=file_path,
        message=(
            "参考资料包含 inputs/outputs、执行步骤、角色能力边界，或该 reference 是非执行型指导文档。"
            if role_sections_ok
            else f"{file_path} 缺少子任务合同章节：{', '.join(missing_role_sections)}。"
        ),
        expected="执行型 reference 包含 inputs/outputs、执行步骤或命令模板、角色能力/禁止能力边界；非执行型 reference 可只保留规则/示例/约束。",
        minimal_edit="增加 ## Inputs / Outputs、## 执行步骤、## 角色与能力边界 等章节。",
    ))

    reference_commands = _reference_script_commands(stripped)
    script_paths_to_check = {script_path for script_path, _ in reference_commands}
    script_paths_to_check.update(_paths_requiring_skill_md_mentions(purpose, prefix="scripts/"))
    for script_path in sorted(script_paths_to_check):
        entry = _skill_plan_entry_for_file(file_path=script_path, purpose=purpose, blueprint_text=purpose)
        results.extend(_check_reference_skillplan_redefinitions(file_path, stripped, entry))
    for script_path, command in reference_commands:
        entry = _skill_plan_entry_for_file(file_path=script_path, purpose=purpose, blueprint_text=purpose)
        skill_md_commands = _extract_script_command_templates(purpose, script_path)
        results.append(ContractCheckResult(
            id="reference.command_block.not_duplicate_skill_md",
            passed=not skill_md_commands,
            target=f"{file_path}#{script_path}",
            message=(
                "SKILL.md 未提供该脚本命令，reference 可承担必要命令块。"
                if not skill_md_commands
                else "SKILL.md 已包含该脚本可执行命令块；reference 默认应只写规范/风格/示例/质量标准，不要重复写命令块。"
            ),
            expected="如果 SKILL.md 已有命令块，reference 不再写命令块；只有 SKILL.md 委托 reference 定义命令时才写。",
            minimal_edit="删除 reference 中重复的命令块，保留写作规范、风格要求、正例、反例和质量标准。",
        ))
        results.extend(_check_command_block_contract(script_path, [command], entry))
        required_capabilities = list(entry.required_capabilities or [])
        missing_capability_mentions = [capability for capability in required_capabilities if capability not in stripped]
        results.append(ContractCheckResult(
            id="reference.required_capabilities.mentioned",
            passed=not missing_capability_mentions,
            target=f"{file_path}#{script_path}",
            message=(
                "reference 命令块说明了脚本 required_capabilities。"
                if not missing_capability_mentions
                else f"reference 缺少这些 required_capabilities 说明：{', '.join(missing_capability_mentions)}。"
            ),
            expected=f"执行型 reference 必须说明 role={entry.role} 以及 required_capabilities={', '.join(required_capabilities) or 'none'}。",
            minimal_edit="在角色与能力边界小节补充 role 和 required_capabilities，并说明输入/输出如何跨脚本传递。",
        ))

    has_placeholder = bool(_REFERENCE_PLACEHOLDER_RE.search(stripped))
    results.append(ContractCheckResult(
        id="reference.no_placeholder_phrases",
        passed=not has_placeholder,
        target=file_path,
        message=(
            "参考资料未包含占位短语。"
            if not has_placeholder
            else f"{file_path} 包含 placeholder/TODO/待补充等占位短语。"
        ),
        expected="不要使用 placeholder、TODO、待补充、将要生成等占位表达。",
        minimal_edit="删除占位短语并替换为实际任务规则和示例。",
    ))
    return results



def _reference_script_commands(content: str) -> list[tuple[str, str]]:
    commands: list[tuple[str, str]] = []
    for match in re.finditer(r"```(?:bash|sh|shell)?\s*\n([\s\S]*?)\n```", content, flags=re.IGNORECASE):
        command = match.group(1).strip()
        for script_match in re.finditer(r"scripts/[A-Za-z0-9_./-]+", command):
            commands.append((script_match.group(0), command))
    return commands


def _declared_list_in_text(field_name: str, content: str) -> list[str] | None:
    pattern = re.compile(rf"(?:^|\b){re.escape(field_name)}\s*[：:=]\s*\[?([^\]\n;]+)\]?", re.I | re.M)
    match = pattern.search(content or "")
    if not match:
        return None
    return [re.sub(r"[^A-Za-z0-9_./-]", "", item.strip().strip("'\"")) for item in re.split(r"[,，、]\s*", match.group(1)) if item.strip()]


def _declared_role_in_text(content: str) -> str | None:
    match = re.search(r"(?:^|\b)role\s*[：:=]\s*(text_generator|image_generator|composite_generator|pdf_builder|docx_builder|pptx_builder|html_asset_builder|asset_builder|generic_script)", content or "", re.I | re.M)
    return match.group(1) if match else None


def _anti_example_sections(content: str) -> str:
    chunks: list[str] = []
    matches = list(re.finditer(r"(?im)^#{1,3}.*(?:反例|错误示例|Anti[- ]?examples?).*$", content or ""))
    for idx, match in enumerate(matches):
        start = match.end()
        next_heading = re.search(r"(?m)^#{1,3}\s+", content[start:])
        end = start + next_heading.start() if next_heading else len(content)
        chunks.append(content[start:end])
    return "\n".join(chunks)


def _check_reference_skillplan_redefinitions(file_path: str, content: str, entry: SkillPlanEntry) -> list[ContractCheckResult]:
    """Ensure references do not invent a second script interface contract."""
    results: list[ContractCheckResult] = []
    declared_role = _declared_role_in_text(content)
    role_ok = declared_role is None or declared_role == entry.role
    results.append(ContractCheckResult(
        id="reference.role.matches_skillplan",
        passed=role_ok,
        target=f"{file_path}#{entry.path}",
        message=("reference 未重新定义冲突 role。" if role_ok else f"reference 重新定义 role={declared_role}，与 SkillPlan.role={entry.role} 冲突。"),
        expected=f"reference 默认不要定义 role；如提及只能是 role={entry.role}。",
        minimal_edit="删除 reference 中的 role/能力合同定义，改写为写作规范、风格要求、示例和质量标准。",
    ))
    for field_name, expected_values in (
        ("inputs", entry.inputs or ["payload"]),
        ("outputs", entry.outputs),
        ("required_capabilities", entry.required_capabilities),
        ("forbidden_capabilities", entry.forbidden_capabilities),
    ):
        declared = _declared_list_in_text(field_name, content)
        ok = declared is None or declared == expected_values
        results.append(ContractCheckResult(
            id=f"reference.{field_name}.matches_skillplan",
            passed=ok,
            target=f"{file_path}#{entry.path}",
            message=(f"reference 未重新定义冲突 {field_name}。" if ok else f"reference {field_name}={declared} 与 SkillPlan {field_name}={expected_values} 冲突。"),
            expected=f"reference 默认不要定义 {field_name}；如提及必须逐字等于 SkillPlan: {expected_values}。",
            minimal_edit=f"删除或修正 {field_name} 小节，避免产生第二套接口合同。",
        ))
    anti = _anti_example_sections(content)
    correct_command_in_anti = bool(anti and entry.command_template and entry.command_template in anti)
    correct_keys_in_anti = False
    for command in re.findall(r"```(?:bash|sh|shell)?\s*\n([\s\S]*?)\n```", anti, flags=re.I):
        keys = _command_payload_keys(command.strip(), entry.path)
        if keys == set(entry.inputs or ["payload"]):
            correct_keys_in_anti = True
    results.append(ContractCheckResult(
        id="reference.anti_example.not_skillplan_command",
        passed=not correct_command_in_anti and not correct_keys_in_anti,
        target=f"{file_path}#anti-examples",
        message=("reference 未把 SkillPlan 正确命令/JSON keys 写成反例。" if not correct_command_in_anti and not correct_keys_in_anti else "reference 把 SkillPlan.command_template 或正确 JSON keys 写入反例，导致合同冲突。"),
        expected="反例只能展示 extra key、缺失 key、错误 runner 或非 JSON argv；不得否定 SkillPlan.command_template。",
        minimal_edit="从反例中移除正确命令，改为错误示例例如 extra 参数或 payload 包装。",
    ))
    return results

def _validate_reference_file_contract(file_path: str, content: str, purpose: str = "") -> None:
    results = _check_reference_file_contract(file_path, content, purpose)
    if any(not result.passed for result in results):
        raise ContractValidationError(_format_contract_failures(results).replace("SKILL.md contract", f"{file_path} contract"), results)



def _asset_extension_check(file_path: str, stripped: str) -> tuple[bool, str, str]:
    ext = Path(file_path).suffix.lower()
    if not stripped:
        return True, "空内容由 asset.not_empty 检查处理。", "当前 asset 文件内容非空。"
    if ext == ".json":
        try:
            json.loads(stripped)
        except json.JSONDecodeError as exc:
            return False, f"{file_path} 不是合法 JSON: {exc.msg}", "JSON asset 必须可被 json.loads 解析。"
        return True, "JSON asset 可解析。", "JSON asset 必须可被 json.loads 解析。"
    if ext in {".yaml", ".yml"}:
        try:
            yaml.safe_load(stripped)
        except yaml.YAMLError as exc:
            return False, f"{file_path} 不是合法 YAML: {exc}", "YAML asset 必须可被 yaml.safe_load 解析。"
        return True, "YAML asset 可解析。", "YAML asset 必须可被 yaml.safe_load 解析。"
    if ext == ".csv":
        rows = list(csv.reader(io.StringIO(stripped)))
        header = rows[0] if rows else []
        if len(rows) < 3 or not header or any(not cell.strip() for cell in header):
            return False, f"{file_path} CSV 必须包含非空表头和至少 2 行数据。", "CSV asset 必须包含 header 和至少 2 行数据。"
        return True, "CSV asset 包含表头和至少 2 行数据。", "CSV asset 必须包含 header 和至少 2 行数据。"
    if ext in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
        data = stripped.encode("latin1", errors="ignore")
        if re.fullmatch(r"[A-Za-z0-9+/=\s]+", stripped) and len(stripped) > 24:
            try:
                data = base64.b64decode(stripped, validate=True)
            except ValueError:
                data = stripped.encode("latin1", errors="ignore")
        magic_ok = (
            data.startswith(b"\x89PNG\r\n\x1a\n")
            or data.startswith(b"\xff\xd8\xff")
            or data.startswith(b"GIF87a")
            or data.startswith(b"GIF89a")
            or data.startswith(b"RIFF") and b"WEBP" in data[:16]
        )
        dims = _image_dimensions(data)
        size_ok = dims is not None and dims[0] >= 64 and dims[1] >= 64
        image_ok = magic_ok and size_ok
        return (
            image_ok,
            "image asset 头部和尺寸合法。" if image_ok else f"{file_path} 不是有效图片或尺寸小于 64x64。",
            "图片 asset 必须是有效图片字节或 base64，且尺寸 >= 64x64。",
        )
    if ext == ".pdf":
        pdf_ok = stripped.startswith("%PDF-") and "%%EOF" in stripped and len(stripped.encode("latin1", errors="ignore")) > 100
        return (
            pdf_ok,
            "PDF asset 结构合法。" if pdf_ok else f"{file_path} 必须以 %PDF- 开头、包含 %%EOF 且大于 100 bytes。",
            "PDF asset 必须是有效、非空 PDF 内容。",
        )
    if ext in {".md", ".txt"}:
        quality_ok = len(stripped) >= 40 and not _REFERENCE_PLACEHOLDER_RE.search(stripped)
        return (
            quality_ok,
            "Markdown/text asset 满足最低质量要求。" if quality_ok else f"{file_path} 文本资源过短或包含占位短语。",
            "Markdown/text asset 至少 40 个字符且不能包含占位短语。",
        )
    return True, "asset 格式可解析。", "当前 asset 文件内容必须符合其扩展名对应格式。"



def _image_dimensions(data: bytes) -> tuple[int, int] | None:
    if data.startswith(b"\x89PNG\r\n\x1a\n") and len(data) >= 24:
        return int.from_bytes(data[16:20], "big"), int.from_bytes(data[20:24], "big")
    if data.startswith((b"GIF87a", b"GIF89a")) and len(data) >= 10:
        return int.from_bytes(data[6:8], "little"), int.from_bytes(data[8:10], "little")
    if data.startswith(b"\xff\xd8"):
        idx = 2
        while idx + 9 < len(data):
            if data[idx] != 0xFF:
                idx += 1
                continue
            marker = data[idx + 1]
            idx += 2
            if marker in {0xD8, 0xD9}:
                continue
            if idx + 2 > len(data):
                break
            length = int.from_bytes(data[idx:idx + 2], "big")
            if length < 2:
                break
            if marker in {0xC0, 0xC1, 0xC2, 0xC3, 0xC5, 0xC6, 0xC7, 0xC9, 0xCA, 0xCB, 0xCD, 0xCE, 0xCF} and idx + 7 <= len(data):
                height = int.from_bytes(data[idx + 3:idx + 5], "big")
                width = int.from_bytes(data[idx + 5:idx + 7], "big")
                return width, height
            idx += length
    return None


def _validate_pdf_trial_outputs(payload: dict[str, Any], *, skill_dir: Path | None, args: list[str], stdout: str) -> None:
    candidates: list[str] = []
    pdf_path = payload.get("pdf_path")
    if isinstance(pdf_path, str) and pdf_path.strip():
        candidates.append(pdf_path.strip())
    file_paths = payload.get("file_paths")
    if isinstance(file_paths, list):
        candidates.extend(p.strip() for p in file_paths if isinstance(p, str) and p.strip())

    checked: list[Path] = []
    for candidate in candidates:
        path = Path(candidate)
        if not path.is_absolute() and skill_dir is not None:
            path = (skill_dir / "scripts" / path).resolve()
        checked.append(path)
        if path.is_file():
            data = path.read_bytes()
            if len(data) > 100 and data.startswith(b"%PDF-") and b"%%EOF" in data[-2048:]:
                return
    raise ValueError(
        "pdf_builder 试运行必须生成真实可用 PDF：文件存在、大小 > 100 bytes、以 %PDF- 开头且包含 %%EOF。"
        f" argv={args!r} stdout={stdout[-4000:]} checked={[str(p) for p in checked]}"
    )

def _check_asset_file_contract(file_path: str, content: str) -> list[ContractCheckResult]:
    stripped = content.strip()
    has_runtime_code = bool(_PLATFORM_IMAGE_HELPER_RE.search(stripped))
    parse_ok, parse_message, parse_expected = _asset_extension_check(file_path, stripped)
    return [
        ContractCheckResult(
            id="asset.not_empty",
            passed=bool(stripped),
            target=file_path,
            message=("asset 内容非空。" if stripped else f"{file_path} asset 内容为空。"),
            expected="输出当前 asset 的模板或静态资源内容。",
            minimal_edit="补充真实模板/静态资源内容，不要输出空壳。",
        ),
        ContractCheckResult(
            id="asset.parseable",
            passed=parse_ok,
            target=file_path,
            message=parse_message,
            expected=parse_expected,
            minimal_edit="按文件扩展名修正格式：JSON/YAML/CSV/image/PDF/Markdown 文本必须可解析且非空。",
        ),
        ContractCheckResult(
            id="asset.no_runtime_capability",
            passed=not has_runtime_code,
            target=file_path,
            message=(
                "asset 未包含运行时图片生成能力。"
                if not has_runtime_code
                else f"{file_path} 是 asset，但包含图片生成 helper/运行时代码。"
            ),
            expected="asset 只能是模板或静态资源，不得执行 image_generation 等能力。",
            minimal_edit="删除运行时代码或将该职责拆分为 scripts/ 文件。",
        ),
    ]


def _validate_asset_file_contract(file_path: str, content: str) -> None:
    results = _check_asset_file_contract(file_path, content)
    if any(not result.passed for result in results):
        raise ContractValidationError(_format_contract_failures(results).replace("SKILL.md contract", f"{file_path} contract"), results)



def _script_satisfies_required_capability(content: str, capability: str) -> bool:
    capability = capability.lower()
    if capability == "text_generation":
        return bool(re.search(r"generate_text_with_llm|LLM_BASE_URL|TEXT_MODEL|chat/completions|complete_chat_once|stream_chat", content, re.IGNORECASE))
    if capability == "image_generation":
        # Direct image API usage is rejected later with a more specific error
        # (for example VISION_MODEL misuse).  Treat it as an attempted image
        # capability here so the repair loop sees the precise role/API failure
        # instead of a generic "missing required_capabilities" message.
        return bool(_PLATFORM_IMAGE_HELPER_RE.search(content) or _DIRECT_IMAGE_API_RE.search(content))
    if capability == "pdf_generation":
        return bool(re.search(r"FPDF|reportlab|PdfWriter|pdf\.output|canvas\.Canvas|build_pdf|%PDF-|pdf_path|\.pdf[\"\']", content, re.IGNORECASE))
    if capability == "docx_generation":
        return bool(re.search(r"Document\(|python-docx|word/document.xml|ZipFile\(|build_docx", content, re.IGNORECASE))
    if capability == "pptx_generation":
        return bool(re.search(r"Presentation\(|python-pptx|ppt/presentation.xml|ZipFile\(|build_pptx", content, re.IGNORECASE))
    if capability in {"html_generation", "html_asset_generation"}:
        return bool(re.search(r"write_text|open\s*\(|<html|<!DOCTYPE html|assets/generated|build_html", content, re.IGNORECASE))
    if capability == "file_output":
        return bool(re.search(r"write_text|write_bytes|open\s*\(|fs\.writeFile|pdf\.output|prs\.save|ZipFile\(|build_(?:pdf|docx|pptx|html)", content, re.IGNORECASE))
    return True



_ARTIFACT_OUTPUT_KEYS = {"pdf_path", "docx_path", "pptx_path", "html_path", "file_paths"}
_ARTIFACT_CAPABILITIES = {"pdf_generation", "docx_generation", "pptx_generation", "html_generation", "html_asset_generation", "file_output"}


def _script_has_real_file_creation_logic(content: str, *, outputs: list[str], capabilities: list[str]) -> bool:
    declared_artifacts = _ARTIFACT_OUTPUT_KEYS & set(outputs or [])
    required_artifacts = _ARTIFACT_CAPABILITIES & set(capabilities or [])
    if not declared_artifacts and not required_artifacts:
        return True
    has_writer = bool(re.search(
        r"write_text|write_bytes|open\s*\([^)]*,\s*[rbu'\"]*w|pdf\.output|prs\.save|ZipFile\s*\(|shutil\.copy|Path\([^)]*\)\.write_|build_(?:pdf|docx|pptx|html)",
        content,
        re.IGNORECASE,
    ))
    returns_only_paths = bool(re.search(
        r"return\s*\{[^}]*['\"](?:pdf_path|docx_path|pptx_path|html_path|file_paths)['\"][^}]*\}",
        content,
        re.IGNORECASE | re.DOTALL,
    )) and not has_writer
    return has_writer and not returns_only_paths

_MODEL_CAPABILITIES = {"text_generation", "image_generation"}
_DETERMINISTIC_BUILDER_ROLES = {"pdf_builder", "docx_builder", "pptx_builder", "html_asset_builder", "asset_builder"}


def _effective_required_capabilities_for_script(plan_entry: SkillPlanEntry) -> list[str]:
    """Return capabilities that this script source must visibly exercise.

    File builders/exporters are deterministic by default.  If a global
    SKILL.md/blueprint model declaration was accidentally copied into a
    builder's required_capabilities, do not turn that into a requirement for
    ``build_pdf.py`` (or sibling exporters) to call LLM/IMAGE_MODEL.  Model
    scripts keep their text/image requirements through their own generator
    roles, while builders are validated for real artifact creation.
    """
    capabilities = list(plan_entry.required_capabilities or [])
    if plan_entry.role in _DETERMINISTIC_BUILDER_ROLES:
        capabilities = [capability for capability in capabilities if capability not in _MODEL_CAPABILITIES]
    return capabilities


def _script_required_capability_failures(content: str, capabilities: list[str]) -> list[str]:
    return [capability for capability in capabilities if not _script_satisfies_required_capability(content, capability)]

def _check_script_file_contract(file_path: str, content: str, role: str | None = None, skill_plan_entry: dict[str, Any] | None = None) -> list[ContractCheckResult]:
    plan_entry = _skill_plan_entry_for_file(file_path=file_path, role=role, skill_plan_entry=skill_plan_entry)
    strict_interface = skill_plan_entry is not None
    stripped = content.strip()
    effective_required_capabilities = _effective_required_capabilities_for_script(plan_entry)
    has_markdown_or_bundle = "```" in stripped or "~~~" in stripped or bool(_MULTI_FILE_MARKER_RE.search(stripped))
    raw_ok = bool(stripped) and not has_markdown_or_bundle
    results = [
        ContractCheckResult(
            id="script.raw_source.single_file",
            passed=raw_ok,
            target=file_path,
            message=(
                "脚本是单个裸源码文件。"
                if raw_ok
                else f"{file_path} 生成内容包含 Markdown 代码块或多文件包，不是单个脚本源码。请重新生成该文件。"
            ),
            expected="只输出单个脚本源码本身，不要 Markdown fence、说明文字、写入文件标签或多文件包。",
            minimal_edit="从上一次内容中只保留目标脚本源码；删除所有 ``` fence、文件路径标题、写入文件标签和说明文字。",
        )
    ]
    if not raw_ok:
        return results

    syntax_ok = True
    syntax_message = f"{plan_entry.language} 源码基础校验通过。"
    syntax_expected = "脚本源码必须符合 language/runtime 的基础语法与入口约定。"
    if plan_entry.language == "python":
        try:
            ast.parse(stripped)
        except SyntaxError as exc:
            syntax_ok = False
            syntax_message = f"{file_path} 生成内容不是合法 Python 源码: {exc.msg}"
        syntax_expected = "Python 脚本必须能通过 ast.parse 语法检查。"
    elif plan_entry.runtime == "node":
        syntax_ok = "process.argv" in stripped and "console.log" in stripped
        syntax_message = "Node/JS 脚本包含 process.argv 和 stdout JSON 输出。" if syntax_ok else f"{file_path} Node/JS 脚本必须使用 process.argv 读取 JSON argv 并 console.log 输出 JSON。"
        syntax_expected = "Node/JS 脚本必须使用 process.argv[2] + JSON.parse，并通过 console.log(JSON.stringify(...)) 输出 JSON。"
    elif plan_entry.runtime in {"bash", "shell"}:
        syntax_ok = "$1" in stripped or "${1" in stripped
        syntax_message = "Shell/Bash 脚本读取 $1 JSON argv。" if syntax_ok else f"{file_path} Shell/Bash 脚本必须读取 $1 JSON argv。"
        syntax_expected = "Shell/Bash 脚本必须读取 $1 JSON argv，并向 stdout 输出 JSON 或写入声明的文件产物。"
    results.append(
        ContractCheckResult(
            id="script.source.syntax",
            passed=syntax_ok,
            target=file_path,
            message=syntax_message,
            expected=syntax_expected,
            minimal_edit="修正源码语法/入口错误，同时保持 stdout JSON 和参数接口不变。",
        )
    )

    if strict_interface:
        reads_json = _script_reads_json_argv(stripped, plan_entry.runtime)
        results.append(
            ContractCheckResult(
                id="script.json_argv.runtime",
                passed=reads_json,
                target=file_path,
                message=(
                    f"脚本按 {plan_entry.runtime} runtime 读取 JSON argv。"
                    if reads_json
                    else f"{file_path} 必须按 {plan_entry.runtime} runtime 读取 JSON argv。"
                ),
                expected="Python: sys.argv[1]+json.loads；Node: process.argv[2]+JSON.parse；Bash: $1 JSON。",
                minimal_edit="补充 runtime 对应 JSON argv 解析入口。",
            )
        )
        uses_keys, missing_keys = _script_uses_input_keys(stripped, list(plan_entry.inputs or ["payload"]))
        results.append(
            ContractCheckResult(
                id="script.skillplan_inputs.used",
                passed=uses_keys,
                target=file_path,
                message=(
                    "脚本源码使用了所有 SkillPlan inputs。"
                    if uses_keys
                    else f"脚本源码未使用这些 SkillPlan inputs：{', '.join(missing_keys)}。"
                ),
                expected=f"源码必须实际读取/使用 inputs：{', '.join(plan_entry.inputs or ['payload'])}。",
                minimal_edit="在参数解析或业务逻辑中读取并使用缺失的 payload key。",
            )
        )
        has_entry = _script_has_main_entry(stripped, plan_entry.runtime)
        results.append(
            ContractCheckResult(
                id="script.runtime.entrypoint",
                passed=has_entry,
                target=file_path,
                message=("脚本包含 runtime 入口与 stdout 输出。" if has_entry else f"{file_path} 缺少 {plan_entry.runtime} 入口或 stdout 输出。"),
                expected="脚本包含对应 runtime 的入口函数/语句，并向 stdout 输出 JSON。",
                minimal_edit="补齐 main/入口调用和 JSON stdout 输出。",
            )
        )
        missing_capabilities = _script_required_capability_failures(stripped, effective_required_capabilities)
        results.append(
            ContractCheckResult(
                id="script.required_capabilities.called",
                passed=not missing_capabilities,
                target=file_path,
                message=(
                    "脚本调用了 role.required_capabilities 对应的平台/文件能力。"
                    if not missing_capabilities
                    else f"脚本未调用这些 required_capabilities 对应接口：{', '.join(missing_capabilities)}。"
                ),
                expected="text_generation 调用 generate_text_with_llm/平台 LLM；image_generation 调用 generate_stable_diffusion_image；pdf_generation 使用 reportlab/fpdf/PDF 构建；file_output 写入声明文件。",
                minimal_edit="按 role + runtime 注入对应平台 helper 或真实文件/PDF 构建逻辑，不要返回固定占位文本。",
            )
        )

    missing_capabilities = _script_required_capability_failures(stripped, effective_required_capabilities)
    results.append(
        ContractCheckResult(
            id="script.required_capabilities.called",
            passed=not missing_capabilities,
            target=file_path,
            message=(
                "脚本调用了 role.required_capabilities 对应的平台/文件能力。"
                if not missing_capabilities
                else f"脚本没有调用这些 required_capabilities 对应接口：{', '.join(missing_capabilities)}。"
            ),
            expected="text_generation 调用 generate_text_with_llm/平台 LLM；image_generation 调用 generate_stable_diffusion_image；pdf_generation 使用 reportlab/fpdf/PDF 构建；file_output 写入声明文件。",
            minimal_edit="按 role + runtime 注入对应平台 helper 或真实文件/PDF 构建逻辑，不要返回固定占位文本。",
        )
    )

    missing_capabilities = _script_required_capability_failures(stripped, effective_required_capabilities)
    results.append(
        ContractCheckResult(
            id="script.required_capabilities.called",
            passed=not missing_capabilities,
            target=file_path,
            message=(
                "脚本调用了 role.required_capabilities 对应的平台/文件能力。"
                if not missing_capabilities
                else f"脚本没有调用这些 required_capabilities 对应接口：{', '.join(missing_capabilities)}。"
            ),
            expected="text_generation 调用 generate_text_with_llm/平台 LLM；image_generation 调用 generate_stable_diffusion_image；pdf_generation 使用 reportlab/fpdf/PDF 构建；file_output 写入声明文件。",
            minimal_edit="按 role + runtime 注入对应平台 helper 或真实文件/PDF 构建逻辑，不要返回固定占位文本。",
        )
    )

    missing_capabilities = _script_required_capability_failures(stripped, effective_required_capabilities)
    results.append(
        ContractCheckResult(
            id="script.required_capabilities.called",
            passed=not missing_capabilities,
            target=file_path,
            message=(
                "脚本调用了 role.required_capabilities 对应的平台/文件能力。"
                if not missing_capabilities
                else f"脚本没有调用这些 required_capabilities 对应接口：{', '.join(missing_capabilities)}。"
            ),
            expected="text_generation 调用 generate_text_with_llm/平台 LLM；image_generation 调用 generate_stable_diffusion_image；pdf_generation 使用 reportlab/fpdf/PDF 构建；file_output 写入声明文件。",
            minimal_edit="按 role + runtime 注入对应平台 helper 或真实文件/PDF 构建逻辑，不要返回固定占位文本。",
        )
    )

    missing_capabilities = _script_required_capability_failures(stripped, effective_required_capabilities)
    results.append(
        ContractCheckResult(
            id="script.required_capabilities.called",
            passed=not missing_capabilities,
            target=file_path,
            message=(
                "脚本调用了 role.required_capabilities 对应的平台/文件能力。"
                if not missing_capabilities
                else f"脚本没有调用这些 required_capabilities 对应接口：{', '.join(missing_capabilities)}。"
            ),
            expected="text_generation 调用 generate_text_with_llm/平台 LLM；image_generation 调用 generate_stable_diffusion_image；pdf_generation 使用 reportlab/fpdf/PDF 构建；file_output 写入声明文件。",
            minimal_edit="按 role + runtime 注入对应平台 helper 或真实文件/PDF 构建逻辑，不要返回固定占位文本。",
        )
    )

    enforce_artifact_outputs = strict_interface or bool(_ARTIFACT_CAPABILITIES & set(effective_required_capabilities))
    has_real_file_output = _script_has_real_file_creation_logic(
        stripped,
        outputs=list(plan_entry.outputs or []) if enforce_artifact_outputs else [],
        capabilities=effective_required_capabilities if enforce_artifact_outputs else [],
    )
    results.append(
        ContractCheckResult(
            id="script.file_outputs.real_creation_logic",
            passed=has_real_file_output,
            target=file_path,
            message=(
                "脚本声明文件输出时包含真实文件创建逻辑。"
                if has_real_file_output
                else f"{file_path} 声明 pdf/docx/pptx/html/file_paths 输出或文件生成能力，但只返回路径字符串或缺少真实写文件逻辑。"
            ),
            expected="声明 pdf_path/docx_path/pptx_path/html_path/file_paths 或文件生成 required_capabilities 时，必须真实创建对应文件，禁止只返回路径占位。",
            minimal_edit="拆出 build_pdf/build_docx/build_pptx/build_html 等函数，在其中写入真实文件并只返回已创建文件路径。",
        )
    )

    has_fake = bool(_SCRIPT_FAKE_IMPLEMENTATION_RE.search(stripped))
    results.append(
        ContractCheckResult(
            id="script.no_fake_implementation",
            passed=not has_fake,
            target=file_path,
            message=(
                "脚本未包含占位/模拟/假 API 实现。"
                if not has_fake
                else f"{file_path} 包含占位/模拟/假 API 实现。Creator 生成的脚本必须具备真实可执行功能。"
            ),
            expected="不得使用 placeholder/mock/fake API/固定模板冒充真实能力。",
            minimal_edit="替换占位或模拟逻辑，实现真实可执行算法或调用平台配置模型/helper。",
        )
    )

    uses_image_helper = bool(_PLATFORM_IMAGE_HELPER_RE.search(stripped))
    if plan_entry.role in {"pdf_builder", "docx_builder", "pptx_builder", "html_asset_builder", "asset_builder"}:
        results.append(
            ContractCheckResult(
                id="script.capability.forbidden_image_generation",
                passed=not uses_image_helper,
                target=file_path,
                message=(
                    "文件构建脚本未调用图片生成 helper。"
                    if not uses_image_helper
                    else f"{file_path} 是文件构建脚本，但调用了图片生成 helper。"
                ),
                expected="文件构建脚本只能消费已有数据并真实创建文件；不得调用 generate_stable_diffusion_image。",
                minimal_edit="移除图片 helper 调用，只保留文件构建逻辑。",
            )
        )
    if "image_generation" in (plan_entry.forbidden_capabilities or []):
        results.append(
            ContractCheckResult(
                id="script.capability.forbidden_image_generation",
                passed=not uses_image_helper,
                target=file_path,
                message=(
                    "脚本未调用 forbidden_capabilities 中禁止的 image_generation。"
                    if not uses_image_helper
                    else f"{file_path} 的 SkillPlan forbidden_capabilities 包含 image_generation，但脚本调用了图片生成 helper。"
                ),
                expected="只有 required_capabilities 包含 image_generation 且未被 forbidden_capabilities 禁止时，脚本才可调用 generate_stable_diffusion_image。",
                minimal_edit="若脚本确实需要图片能力，请修正 SkillPlan required_capabilities/forbidden_capabilities；否则移除图片 helper 调用。",
            )
        )

    uses_text_helper = bool(re.search(r"generate_text_with_llm|LLM_BASE_URL|TEXT_MODEL|chat/completions|complete_chat_once|stream_chat", stripped, re.IGNORECASE))
    if "text_generation" in (plan_entry.forbidden_capabilities or []):
        results.append(
            ContractCheckResult(
                id="script.capability.forbidden_text_generation",
                passed=not uses_text_helper,
                target=file_path,
                message=(
                    "脚本未调用 forbidden_capabilities 中禁止的 text_generation。"
                    if not uses_text_helper
                    else f"{file_path} 的 SkillPlan forbidden_capabilities 包含 text_generation，但脚本调用了文本生成 helper/LLM。"
                ),
                expected="只有 required_capabilities 包含 text_generation 且未被 forbidden_capabilities 禁止时，脚本才可调用 generate_text_with_llm 或平台文本模型。",
                minimal_edit="若脚本确实需要文本生成能力，请修正 SkillPlan required_capabilities/forbidden_capabilities；否则移除文本模型调用。",
            )
        )

    writes_pdf = bool(re.search(r"\.pdf[\"']|pdf_path|FPDF|reportlab|PdfWriter|write_bytes\s*\(\s*b?[\"']%PDF-", stripped, re.IGNORECASE))
    if "pdf_generation" in (plan_entry.forbidden_capabilities or []):
        results.append(
            ContractCheckResult(
                id="script.capability.forbidden_pdf_generation",
                passed=not writes_pdf,
                target=file_path,
                message=(
                    "脚本未调用 forbidden_capabilities 中禁止的 pdf_generation。"
                    if not writes_pdf
                    else f"{file_path} 的 SkillPlan forbidden_capabilities 包含 pdf_generation，但源码包含 PDF 生成/输出逻辑。"
                ),
                expected="只有 required_capabilities 包含 pdf_generation 且未被 forbidden_capabilities 禁止时，脚本才可构建 PDF。",
                minimal_edit="若脚本确实需要 PDF 能力，请修正 SkillPlan required_capabilities/forbidden_capabilities；否则移除 PDF 生成逻辑。",
            )
        )

    if plan_entry.role == "image_generator":
        pdf_only = ("pdf_path" in stripped or "file_paths" in stripped) and "image_paths" not in stripped and "images" not in stripped
        results.append(
            ContractCheckResult(
                id="script.role.image_forbidden_pdf_only_outputs",
                passed=not pdf_only and not writes_pdf,
                target=file_path,
                message=(
                    "image_generator 未输出或生成 PDF-only 结果。"
                    if not pdf_only and not writes_pdf
                    else f"{file_path} 是 image_generator，但源码包含 PDF-only 输出或 PDF 生成逻辑。"
                ),
                expected="image_generator 必须输出 image_paths/images，不得写 PDF 或只输出 pdf_path/file_paths；需要文本+图片时请使用 composite_generator。",
                minimal_edit="返回 image_paths/images 并删除 PDF 写入逻辑；若要同时生成文本和图片，请将 role/required_capabilities 改为 composite_generator + text_generation/image_generation。",
            )
        )
    return results


def _validate_script_file_source_contract(file_path: str, content: str, role: str | None = None, skill_plan_entry: dict[str, Any] | None = None) -> None:
    # Accept otherwise-valid raw source with a dangling orphan fence marker at
    # the boundary.  Full fenced/bundled responses are still rejected by the
    # lower-level checker unless the sanitize path extracted a single code block.
    candidate = _strip_orphan_trailing_fence(content)
    results = _check_script_file_contract(file_path, candidate, role=role, skill_plan_entry=skill_plan_entry)
    if any(not result.passed for result in results):
        raise ContractValidationError(_format_contract_failures(results).replace("SKILL.md contract", f"{file_path} contract"), results)

def _validate_skill_md_against_existing_files(skill_name: str, content: str) -> None:
    """Validate SKILL.md against files already initialized in the Skill dir."""
    skill_dir = settings.skills_path / skill_name
    if not skill_dir.exists():
        return
    declared_paths = sorted(_existing_skill_local_paths_for_skill(skill_name))
    if declared_paths:
        _validate_skill_md_contract(content, "\n".join(declared_paths))
    else:
        _reject_creator_flow_leak(content)


def _clean_blueprint_for_file_prompt(blueprint_text: str) -> str:
    """Remove Creator UI confirmation text from blueprint context before generation."""
    cleaned_lines: list[str] = []
    in_confirmation_block = False
    for line in (blueprint_text or "").splitlines():
        stripped = line.strip()
        if _CREATOR_FLOW_LEAK_RE.search(stripped):
            in_confirmation_block = True
            continue
        if in_confirmation_block:
            if stripped.startswith("```") or stripped.startswith("- [") or stripped.startswith(">"):
                continue
            if not stripped:
                in_confirmation_block = False
                continue
        cleaned_lines.append(line)
    return "\n".join(cleaned_lines).strip() or blueprint_text


def _reject_fake_script_implementation(file_path: str, content: str) -> None:
    """Reject placeholder/mock scripts that pretend to implement capabilities."""
    if _SCRIPT_FAKE_IMPLEMENTATION_RE.search(content):
        raise ValueError(
            f"{file_path} 包含占位/模拟/假 API 实现。"
            "Creator 生成的脚本必须具备真实可执行功能；"
            "如需图像或多模态能力，应通过宿主配置的模型/服务完成，不能写 placeholder 文件或假装调用 API。"
        )


def _requires_configured_model_call(*, plan_entry: SkillPlanEntry | None) -> bool:
    """Return whether the current script contract requires host model use.

    Model-call requirements are scoped to this script's SkillPlanEntry.
    Whole-SKILL.md wording about LLM/image models can describe earlier or later
    steps, but must not force deterministic exporter/builder scripts to call a
    model unless their own required_capabilities declare text/image generation.
    """
    if plan_entry is None:
        return False
    return bool({"text_generation", "image_generation"} & set(_effective_required_capabilities_for_script(plan_entry)))


def _script_uses_configured_model(content: str) -> bool:
    """Detect whether script calls the configured host LLM/VL endpoint."""
    return bool(_CONFIGURED_MODEL_CALL_RE.search(content))


def _validate_configured_model_usage_static(*, file_path: str, content: str, skill_md: str, plan_entry: SkillPlanEntry | None = None) -> None:
    """Reject scripts whose own SkillPlanEntry requires host-model behavior but do not call models."""
    if _DIRECT_IMAGE_API_RE.search(content) and "VISION_MODEL" in content:
        raise ValueError(
            f"{file_path} 将 VISION_MODEL 与图片生成接口混用。"
            "生成图片必须使用平台 Stable Diffusion 图片运行时或 IMAGE_MODEL；"
            "VISION_MODEL 只用于看图理解/OCR/多模态问答。"
        )

    if _DIRECT_IMAGE_API_RE.search(content) and not _PLATFORM_IMAGE_HELPER_RE.search(content):
        raise ValueError(
            f"{file_path} 直接调用图片生成接口。"
            "Creator 生成的图片脚本必须调用 backend.services.skill_runtime.generate_stable_diffusion_image，"
            "由平台侧静默完成中文 topic 翻译、Stable Diffusion IMAGE_MODEL 选择、b64_json 解析和图片落盘。"
        )

    if _IMAGE_MODEL_USAGE_RE.search(content) and _IMAGE_URL_ONLY_RE.search(content) and "b64_json" not in content:
        raise ValueError(
            f"{file_path} 假设图片接口只返回 url。"
            "平台图片运行时默认使用 b64_json，并会落盘为文件路径；请调用平台图片运行时 helper。"
        )

    if _DATA_URI_RE.search(content):
        raise ValueError(
            f"{file_path} 输出 base64 data URI。"
            "图片结果必须由平台运行时写入 OUTPUT_DIR，并在 stdout JSON 中返回 image_paths/images。"
        )

    if re.search(r"(?m)^\s*image_path\s*=\s*generate_stable_diffusion_image\s*\(", content):
        raise ValueError(
            f"{file_path} 将 helper 返回 dict 直接赋给 image_path。"
            "图片脚本必须先保存 result = generate_stable_diffusion_image(desc)，"
            "再执行 image_paths.append(result.get(\"image_path\")) 和 images.append(result)。"
        )

    effective_required_capabilities = _effective_required_capabilities_for_script(plan_entry) if plan_entry else []
    if plan_entry and plan_entry.role in {"pdf_builder", "docx_builder", "pptx_builder", "html_asset_builder", "asset_builder"} and not ({"text_generation", "image_generation"} & set(effective_required_capabilities)):
        return
    skill_md_declares_model = bool(re.search(r"宿主|内置|配置模型|LLM|大语言|文本模型|图像模型|vision|TEXT_MODEL|IMAGE_MODEL", skill_md or "", re.IGNORECASE))
    if not _requires_configured_model_call(plan_entry=plan_entry):
        # Keep backward compatibility for legacy generic scripts whose SKILL.md
        # has no local SkillPlan block but clearly says the script is model-backed.
        # Deterministic builders/exporters returned above, so global model prose
        # still cannot force build_pdf.py to call LLM/IMAGE_MODEL.
        if not (plan_entry and plan_entry.role == "generic_script" and skill_md_declares_model):
            return
    if _script_uses_configured_model(content):
        return
    raise ValueError(
        f"{file_path} 的当前脚本职责/SkillPlan.required_capabilities 声明需要使用宿主/内置/配置模型，但脚本没有调用这些模型。"
        "脚本不能用固定模板、随机词表或 ASCII 图替代模型能力；"
        "请通过 LLM_BASE_URL + TEXT_MODEL 调用文本模型，需要图像/视觉能力时使用 IMAGE_MODEL/VISION_MODEL。"
    )

def _validate_generated_file_content(file_path: str, content: str, role: str | None = None, skill_plan_entry: dict[str, Any] | None = None) -> None:
    """Reject content that is clearly not the requested single file."""
    if file_path == "SKILL.md":
        _reject_custom_skill_md_protocol(content)
        return

    if file_path.startswith("scripts/"):
        _validate_script_file_source_contract(file_path, content, role=role, skill_plan_entry=skill_plan_entry)
        return

    if file_path.startswith("references/"):
        _validate_reference_file_contract(file_path, content)
        return

    if file_path.startswith("assets/"):
        _validate_asset_file_contract(file_path, content)
        return


def _extract_script_command_templates(skill_md: str, script_path: str) -> list[str]:
    """Return shell command templates in SKILL.md that invoke script_path."""
    commands: list[str] = []
    for match in re.finditer(r"```(?:bash|sh|shell)?\s*\n([\s\S]*?)\n```", skill_md, flags=re.IGNORECASE):
        command = match.group(1).strip()
        if script_path in command:
            commands.append(command)
    return commands


def _command_uses_json_argv(command: str) -> bool:
    return "{" in command and "}" in command


def _script_reads_json_argv(content: str, runtime: str = "python") -> bool:
    if runtime == "node":
        return "JSON.parse" in content and "process.argv" in content
    if runtime in {"bash", "shell"}:
        return "$1" in content or "${1" in content or "jq" in content
    return "json.loads" in content and "sys.argv" in content


def _script_uses_input_keys(content: str, keys: list[str]) -> tuple[bool, list[str]]:
    missing = [key for key in keys if key not in content]
    return not missing, missing


def _script_has_main_entry(content: str, runtime: str) -> bool:
    if runtime == "python":
        return "def main" in content and "__main__" in content
    if runtime == "node":
        return "process.argv" in content and "console.log" in content
    if runtime in {"bash", "shell"}:
        return ("$1" in content or "${1" in content) and ("echo" in content or "printf" in content or "print(json.dumps" in content)
    return True


def _validate_script_contract_static(*, file_path: str, content: str, skill_md: str) -> None:
    """Validate script source against existing SKILL.md command examples."""
    _reject_fake_script_implementation(file_path, content)
    plan_entry = _skill_plan_entry_for_file(file_path=file_path, blueprint_text=skill_md)
    _validate_configured_model_usage_static(file_path=file_path, content=content, skill_md=skill_md, plan_entry=plan_entry)
    commands = _extract_script_command_templates(skill_md, file_path)
    if not commands:
        return

    # Backward compatibility for existing SKILL.md files that predate explicit
    # SkillPlan inputs: keep generic_script conservative capabilities, but treat
    # the JSON argv command itself as the declared interface for static script
    # validation.  New Creator generation still emits command blocks from
    # SkillPlan inputs via _check_skill_md_contract.
    if plan_entry.role == "generic_script" and plan_entry.inputs == ["payload"]:
        command_keys: set[str] = set()
        for command in commands:
            keys = _command_payload_keys(command, file_path)
            if keys:
                command_keys.update(keys)
        if command_keys:
            plan_entry = replace(plan_entry, inputs=sorted(command_keys), command_template=commands[0].strip())
    command_results: list[ContractCheckResult] = []
    command_results.extend(_check_command_block_contract(file_path, commands, plan_entry))
    failed_command_results = [result for result in command_results if not result.passed]
    if failed_command_results:
        raise ValueError(_format_contract_checks(failed_command_results, passed=False))

    json_argv_commands = [cmd for cmd in commands if _command_uses_json_argv(cmd)]
    if json_argv_commands and not _script_reads_json_argv(content, plan_entry.runtime):
        raise ValueError(
            f"{file_path} 的 SKILL.md Markdown 命令示例传入 JSON 参数，但脚本没有按脚本 runtime 读取 JSON argv 并解析（Python 应读取 sys.argv[1] 并 json.loads）；"
            "禁止保存与命令示例不一致的脚本。"
        )

    for cmd in commands:
        for key in re.findall(r"{{\s*([a-zA-Z_][\w-]*)\s*}}", cmd):
            if key not in content:
                raise ValueError(
                    f"{file_path} 的 Markdown 命令示例包含参数 {{{{{key}}}}}，但脚本源码未使用该参数。"
                )


def _validate_script_against_existing_skill_contract(skill_name: str, file_path: str, content: str) -> None:
    """Refuse saving scripts that do not match the current SKILL.md contract."""
    if not file_path.startswith("scripts/"):
        return
    skill_md_path = settings.skills_path / skill_name / "SKILL.md"
    if not skill_md_path.is_file():
        return
    skill_md = skill_md_path.read_text(encoding="utf-8")
    _validate_script_contract_static(file_path=file_path, content=content, skill_md=skill_md)


def _sample_value_for_placeholder(key: str) -> str:
    """Return realistic trial input; image/diffusion prompts are already English."""
    lowered = key.lower()
    if any(token in lowered for token in ("prompt", "diffusion", "image", "picture", "photo", "scene")):
        return "a cinematic watercolor cat under a warm sunset"
    if any(token in lowered for token in ("text", "content", "input", "query")):
        return "测试输入 text sample"
    if any(token in lowered for token in ("topic", "theme", "subject")):
        return "system time"
    return f"sample {key}"


def _render_trial_command_args(command: str, script_path: str) -> list[str] | None:
    """Extract argv for script_path from a SKILL.md command template."""
    rendered = re.sub(
        r"{{\s*([a-zA-Z_][\w-]*)\s*}}",
        lambda m: _sample_value_for_placeholder(m.group(1)),
        command,
    )
    try:
        parts = shlex.split(rendered)
    except ValueError:
        return None

    for idx, part in enumerate(parts):
        normalized = part.replace("\\", "/")
        if normalized == script_path or normalized.endswith("/" + script_path):
            return parts[idx + 1:]
    return None


def _dedupe_trial_arg_sets(arg_sets: list[list[str]]) -> list[list[str]]:
    seen: set[tuple[str, ...]] = set()
    deduped: list[list[str]] = []
    for args in arg_sets:
        key = tuple(args)
        if key in seen:
            continue
        seen.add(key)
        deduped.append(args)
    return deduped


def _json_argv_text_optional_variants(args: list[str]) -> list[list[str]]:
    """Add trial cases proving optional text can be omitted or empty."""
    if len(args) != 1:
        return []
    try:
        payload = json.loads(args[0])
    except json.JSONDecodeError:
        return []
    if not isinstance(payload, dict) or "text" not in payload:
        return []

    variants: list[list[str]] = []
    without_text = dict(payload)
    without_text.pop("text", None)
    if without_text:
        variants.append([json.dumps(without_text, ensure_ascii=False)])

    empty_text = dict(payload)
    empty_text["text"] = ""
    variants.append([json.dumps(empty_text, ensure_ascii=False)])
    return variants


def _trial_args_for_script(skill_md: str, file_path: str, content: str) -> list[list[str]]:
    commands = _extract_script_command_templates(skill_md, file_path)
    arg_sets = [args for cmd in commands if (args := _render_trial_command_args(cmd, file_path)) is not None]
    if not arg_sets and _script_reads_json_argv(content, runtime_for_language(language_for_path(file_path), file_type_for_path(file_path))):
        arg_sets = [[json.dumps({
            "prompt": _sample_value_for_placeholder("prompt"),
            "text": _sample_value_for_placeholder("text"),
            "topic": _sample_value_for_placeholder("topic"),
        }, ensure_ascii=False)]]
    if not arg_sets:
        arg_sets = [[]]

    expanded: list[list[str]] = []
    for args in arg_sets:
        expanded.append(args)
        expanded.extend(_json_argv_text_optional_variants(args))
    return _dedupe_trial_arg_sets(expanded)


def _format_trial_failure(*, args: list[str], returncode: int, stdout: str, stderr: str) -> str:
    return (
        "脚本试运行失败：\n"
        f"argv={args!r}\n"
        f"exit_code={returncode}\n"
        f"stdout={stdout[-4000:]}\n"
        f"stderr={stderr[-4000:]}"
    )


def _validate_image_payload_shape(payload: dict[str, Any]) -> bool:
    image_path = payload.get("image_path")
    if isinstance(image_path, str) and image_path.strip():
        return True

    image_paths = payload.get("image_paths")
    if isinstance(image_paths, list) and image_paths and all(isinstance(p, str) and p.strip() for p in image_paths):
        return True

    images = payload.get("images")
    if (
        isinstance(images, list)
        and images
        and all(isinstance(item, dict) and isinstance(item.get("image_path"), str) and item.get("image_path").strip() for item in images)
    ):
        return True

    return False



def _html_output_candidates(payload: dict[str, Any]) -> list[str]:
    candidates: list[str] = []
    html_path = payload.get("html_path")
    if isinstance(html_path, str) and html_path.strip():
        candidates.append(html_path.strip())
    asset_paths = payload.get("asset_paths")
    if isinstance(asset_paths, list):
        candidates.extend(path.strip() for path in asset_paths if isinstance(path, str) and path.strip())
    return candidates


def _validate_html_trial_outputs(payload: dict[str, Any], *, skill_dir: Path | None, args: list[str], stdout: str) -> None:
    checked: list[Path] = []
    for candidate in _html_output_candidates(payload):
        candidate_path = Path(candidate)
        if not candidate_path.is_absolute() and skill_dir is not None:
            if candidate.startswith("assets/"):
                candidate_path = (skill_dir / candidate_path).resolve()
            else:
                candidate_path = (skill_dir / "scripts" / candidate_path).resolve()
        checked.append(candidate_path)
        if skill_dir is not None:
            generated_root = (skill_dir / "assets" / "generated").resolve()
            try:
                candidate_path.relative_to(generated_root)
            except ValueError:
                raise ValueError(
                    "html_asset_builder 输出路径必须位于当前 Skill 的 assets/generated/ 下："
                    f"argv={args!r} stdout={stdout[-4000:]} checked={candidate_path}"
                )
        if candidate_path.is_file() and candidate_path.suffix.lower() in {".html", ".htm"}:
            text = candidate_path.read_text(encoding="utf-8", errors="replace").lower()
            if "<html" in text or "<!doctype html" in text:
                return
    raise ValueError(
        "html_asset_builder 试运行必须输出 html_path 或 asset_paths，并生成真实 HTML 文件："
        f"argv={args!r} stdout={stdout[-4000:]} checked={[str(p) for p in checked]}"
    )

def _validate_file_payload_shape(payload: dict[str, Any]) -> bool:
    pdf_path = payload.get("pdf_path")
    if isinstance(pdf_path, str) and pdf_path.strip():
        return True

    file_paths = payload.get("file_paths")
    if isinstance(file_paths, list) and file_paths and all(isinstance(p, str) and p.strip() for p in file_paths):
        return True

    return False



_LEGACY_OUTPUT_ALIASES: dict[str, tuple[str, ...]] = {
    "text": ("text_content", "story_text", "content"),
    "image_paths": ("image_path", "images"),
    "images": ("image_path", "image_paths"),
    "pdf_path": ("file_path", "file_paths"),
    "docx_path": ("file_path", "file_paths"),
    "pptx_path": ("file_path", "file_paths"),
    "html_path": ("file_path", "asset_paths"),
    "file_paths": ("file_path", "pdf_path", "docx_path", "pptx_path"),
}


def _payload_has_declared_output(payload: dict[str, Any], output_key: str) -> bool:
    if output_key in payload:
        return True
    return any(alias in payload for alias in _LEGACY_OUTPUT_ALIASES.get(output_key, ()))


def _payload_output_value(payload: dict[str, Any], output_key: str) -> Any:
    if output_key in payload:
        return payload.get(output_key)
    for alias in _LEGACY_OUTPUT_ALIASES.get(output_key, ()):
        if alias in payload:
            return payload.get(alias)
    return None

def _json_value_non_empty(value: Any) -> bool:
    if value is None:
        return False
    if isinstance(value, str):
        return bool(value.strip())
    if isinstance(value, (list, tuple, set)):
        return any(_json_value_non_empty(item) for item in value)
    if isinstance(value, dict):
        return any(_json_value_non_empty(item) for item in value.values())
    return True


def _validate_trial_stdout_json(*, stdout: str, content: str, args: list[str], role: str | None = None, skill_dir: Path | None = None, skill_plan_entry: dict[str, Any] | None = None) -> None:
    """Validate trial stdout with dynamic, field-name-agnostic rules.

    SkillPlan.outputs is a blueprint hint, not the sole runtime contract.  The
    hard requirements here are: stdout is a JSON object, it has at least one
    non-empty value, it does not report an error, and any file-looking values it
    declares point at real files.
    """
    stripped = (stdout or "").strip()
    if not stripped:
        raise ValueError(f"脚本试运行 stdout 为空：argv={args!r}")
    try:
        payload = json.loads(stripped)
    except json.JSONDecodeError as exc:
        raise ValueError(f"脚本试运行 stdout 不是合法 JSON object：argv={args!r} stdout={stripped[-4000:]}") from exc
    if not isinstance(payload, dict):
        raise ValueError(f"脚本试运行 stdout 必须是 JSON object：argv={args!r} stdout={stripped[-4000:]}")
    if "error" in payload:
        raise ValueError(f"脚本试运行 stdout JSON 不得包含 error 字段：argv={args!r} stdout={stripped[-4000:]}")
    if not any(_json_value_non_empty(value) for value in payload.values()):
        raise ValueError(f"脚本试运行 stdout JSON 至少需要一个非空字段：argv={args!r} stdout={stripped[-4000:]}")

    try:
        if skill_dir is not None:
            validate_stdout_file_outputs(stripped, skill_dir=skill_dir, cwd=skill_dir / "scripts")
    except FileOutputValidationError as exc:
        raise ValueError(str(exc)) from exc



def _trial_run_generated_script_with_plan(
    skill_name: str,
    file_path: str,
    content: str,
    *,
    role: str | None = None,
    skill_plan_entry: dict[str, Any] | None = None,
) -> None:
    try:
        _trial_run_generated_script(skill_name, file_path, content, role, skill_plan_entry)
    except TypeError as exc:
        # Some tests monkeypatch the trial runner with the historical
        # 4-argument signature; keep production role-aware calls while allowing
        # those narrow fakes to exercise the repair loop.
        if "positional" not in str(exc) and "unexpected keyword" not in str(exc):
            raise
        _trial_run_generated_script(skill_name, file_path, content, role)


def _trial_run_generated_script(
    skill_name: str,
    file_path: str,
    content: str,
    role: str | None = None,
    skill_plan_entry: dict[str, Any] | None = None,
) -> None:
    """Run a generated Python script before accepting it from Creator.

    Python scripts are executed in a temporary per-skill virtual environment.
    Before each trial run, imports are statically scanned and missing common
    third-party packages are installed into that venv, matching sandbox runtime
    behavior and allowing generation-test-repair-test loops to focus on real
    script defects instead of missing packages.
    """
    if not file_path.startswith("scripts/") or Path(file_path).suffix.lower() != ".py":
        return

    skill_md_path = settings.skills_path / skill_name / "SKILL.md"
    skill_md = skill_md_path.read_text(encoding="utf-8") if skill_md_path.is_file() else ""
    _validate_script_contract_static(file_path=file_path, content=content, skill_md=skill_md)

    with tempfile.TemporaryDirectory(prefix="creator-script-trial-") as tmp:
        skill_dir = Path(tmp) / skill_name
        scripts_dir = skill_dir / "scripts"
        scripts_dir.mkdir(parents=True, exist_ok=True)
        (skill_dir / "SKILL.md").write_text(
            skill_md or f"---\nname: {skill_name}\ndescription: trial\n---\n",
            encoding="utf-8",
        )
        script_path = scripts_dir / Path(file_path).name
        script_path.write_text(content, encoding="utf-8")

        try:
            venv_python = _get_skill_venv_python(skill_dir)
            _scan_and_install_python_deps(script_path, venv_python)
        except RuntimeError as exc:
            raise ValueError(f"脚本试运行环境准备失败：{exc}") from exc

        for args in _trial_args_for_script(skill_md, file_path, content):
            try:
                proc = subprocess.run(
                    [str(venv_python), str(script_path), *args],
                    cwd=str(scripts_dir),
                    capture_output=True,
                    text=True,
                    timeout=_SCRIPT_TRIAL_TIMEOUT_SECONDS,
                    env={**_build_script_runtime_env(skill_dir), "SKILL_TRIAL_RUN": "1"},
                )
            except subprocess.TimeoutExpired as exc:
                raise ValueError(f"脚本试运行超时（超过 {_SCRIPT_TRIAL_TIMEOUT_SECONDS} 秒）：argv={args!r}") from exc
            if proc.returncode != 0:
                raise ValueError(_format_trial_failure(
                    args=args,
                    returncode=proc.returncode,
                    stdout=proc.stdout,
                    stderr=proc.stderr,
                ))
            inferred_entry = _skill_plan_entry_for_file(
                file_path=file_path,
                blueprint_text=skill_md,
                skill_plan_entry=skill_plan_entry,
            )
            inferred_role = role or inferred_entry.role
            _validate_trial_stdout_json(
                stdout=proc.stdout,
                content=content,
                args=args,
                role=inferred_role,
                skill_dir=skill_dir,
                skill_plan_entry=skill_plan_entry,
            )


async def _repair_generated_file_with_feedback(
    *,
    prompt_messages: list[dict],
    model: str,
    file_path: str,
    previous_content: str,
    validation_error: str,
    targeted_repair: str = "",
    contract_text: str = "",
    passed_checks_text: str = "",
    failed_checks_text: str = "",
    repair_mode: str = "minimal_edit",
    skill_plan_entry: dict[str, Any] | None = None,
) -> str:
    """Ask the routed generation model to fix one file using validator feedback.

    The repaired model output is intentionally not validated here. Validation
    happens at the top of the generate-file retry loop so format errors in the
    repaired response consume one retry attempt and can be sent back as feedback.
    Previous content is passed as user-quoted data instead of an assistant turn
    so Markdown-wrapped failures are not reinforced as the desired answer shape.
    """
    is_script = file_path.startswith("scripts/")
    plan_entry = _skill_plan_entry_for_file(file_path=file_path, skill_plan_entry=skill_plan_entry) if is_script else None
    repair_language = plan_entry.language if plan_entry is not None else language_for_path(file_path)
    repair_runtime = plan_entry.runtime if plan_entry is not None else runtime_for_language(repair_language, file_type_for_path(file_path))
    runtime_rule = {
        "python": "Python: parse sys.argv[1] as JSON，保留 main() 入口并 print JSON。",
        "node": "Node: parse process.argv[2] as JSON，并 console.log(JSON.stringify(...))。",
        "bash": "Bash: parse $1 as JSON，并向 stdout 输出 JSON 或声明的文件产物。",
        "shell": "Shell: parse $1 as JSON，并向 stdout 输出 JSON 或声明的文件产物。",
    }.get(repair_runtime, "只输出该文件类型的原始内容；不得包含 Markdown fence。")
    output_contract = (
        f"Rewrite as raw {repair_language} source. Remove any fenced code blocks or file labels. Do NOT include Markdown fences, explanations, file headers, or multi-file content. Align JSON argv keys with the existing SKILL.md command placeholders; do not change the blueprint or SKILL.md. {runtime_rule}"
        if is_script
        else "最终只返回 SKILL.md 文件正文；不要在文件外层套 Markdown 代码块，不要输出 Creator 创建流程、确认清单或 `点击开始创建` 文案。"
    )
    local_edit_scope = (
        "保留其它已经正确的导入、函数、参数解析、stdout JSON 协议和业务逻辑。"
        if is_script
        else "保留已经正确的 frontmatter、章节结构、脚本命令示例和 reference 引用。"
    )
    if is_script:
        role_rule = ""
        if plan_entry is not None and plan_entry.role == "composite_generator":
            caps = set(plan_entry.required_capabilities or [])
            role_rule = "role=composite_generator：表示多能力组合脚本；具体 helper 由 SkillPlan.required_capabilities 决定；stdout 字段名由现有 SKILL.md 后续变量引用/业务语义决定，不要把 composite 固定理解为 text+image。"
            if {"text_generation", "image_generation"} <= caps:
                role_rule += " 当前合同同时要求文本和图片能力，必须保留 generate_text_with_llm 与 generate_stable_diffusion_image。"
        elif plan_entry is not None and plan_entry.role == "image_generator":
            role_rule = "role=image_generator：必须保留并调用 generate_stable_diffusion_image；stdout 输出非空 JSON，并使用现有 SKILL.md/脚本链路会消费的字段名；禁止占位图片或删除真实 helper。"
        elif plan_entry is not None and plan_entry.role == "text_generator":
            role_rule = "role=text_generator：必须调用 generate_text_with_llm 或平台 LLM，禁止调用图片 helper 或输出固定 template-only 文本。"
        elif plan_entry is not None and plan_entry.role == "pdf_builder":
            role_rule = "role=pdf_builder：默认是纯文件合并/排版/PDF 构建脚本，只需真实构建文件并在 stdout JSON 返回实际存在路径；不要因为全局 SKILL.md 提到模型就调用 LLM/IMAGE_MODEL，除非当前脚本的有效 required_capabilities 明确要求模型。"
        elif plan_entry is not None and plan_entry.role == "generic_script":
            role_rule = "role=generic_script：只能调用 SkillPlan.required_capabilities 声明的能力；若现有 SkillPlan.required_capabilities 已要求文本+图片，只能修当前脚本实现以匹配，禁止修改蓝图或 SKILL.md。"
        extra_rules = (
            "Python / Node / Bash 必须按 SkillPlan.runtime 读取单个 JSON argv，并且 JSON argv keys 匹配现有 SKILL.md 命令占位符；"
            "禁止生成 topicstring / tonehumorous / stylepopular-science 这类把 key、类型或默认值拼接起来的字段；"
            f"{role_rule}"
            "不要直接调用 /v1/images/generations，不要用 VISION_MODEL 生成图片，不要写 placeholder/模拟图片。"
        )
    else:
        extra_rules = (
            "如果蓝图包含 scripts/，必须包含调用对应 scripts/ 路径的 ```bash fenced code block；"
            "如果蓝图包含 references/，必须在正文中明确引用对应 reference 路径；"
            "不得复制 Creator UI 流程、待确认清单、文件创建面板说明或系统自动创建文件提示。"
        )

    repair_messages = [*prompt_messages]
    previous_for_prompt = previous_content[-16000:]
    previous_label = "待编辑草稿"
    if is_script and repair_mode == "strict_contract_rewrite":
        previous_for_prompt = (_extract_probable_python_source(previous_content) if repair_language == "python" else _drop_common_non_code_lines(_extract_single_wrapping_fence(previous_content) or previous_content)) or ""
        previous_label = "可参考的源码候选（已移除 Markdown 外壳；若为空，请根据原始任务重新生成）"
    repair_messages.append({
        "role": "user",
        "content": (
            f"以下是上一次生成但未通过校验的 {file_path} 内容。它可能包含错误示范（例如 Markdown fence 或 Creator 流程泄露），"
            f"不要模仿错误格式，只把它当作{previous_label}：\n"
            "<previous_content>\n"
            f"{previous_for_prompt}\n"
            "</previous_content>"
        ),
    })
    edit_strategy = (
        "这是 scripts/ 首次失败后的严格重写：不要走 minimal_edit，不要修补错误外壳；"
        "请根据合同、蓝图和骨架重新输出完整可运行源码。"
        if is_script and repair_mode == "strict_contract_rewrite"
        else (
            "请优先做局部修改：只修改校验意见指出的最小错误片段，"
            f"{local_edit_scope}"
            "修改完成后可以整合输出。"
        )
    )
    repair_messages.append({
        "role": "user",
        "content": (
            f"上一次生成的 {file_path} 没有通过校验模型/静态校验/试运行。"
            f"{edit_strategy}"
            f"{output_contract}\n"
            f"{extra_rules}\n\n"
            f"错误信息：\n{validation_error}"
            + (f"\n\n完整 contract（最终输出必须满足全部条目）：\n{contract_text}" if contract_text else "")
            + (f"\n\n已通过检查（必须保留，不要重写或删除对应内容）：\n{passed_checks_text}" if passed_checks_text else "")
            + (f"\n\n未通过检查（本轮只修这些项）：\n{failed_checks_text}" if failed_checks_text else "")
            + (f"\n\n本轮修复模式：{repair_mode}" if repair_mode else "")
            + ("\n- minimal_edit：只做最小编辑；strict_contract_rewrite：上一轮仍未通过同一 contract，必须重写目标小节但保留已通过项。")
            + ("\n- scripts/ 修复示例：Rewrite as raw <language> source, remove any fenced code blocks or file labels, align JSON argv keys with SkillPlan inputs." if is_script else "")
            + ("\n- 如果这是 scripts/ 文件且进入 strict_contract_rewrite：不要继续修补 Markdown 包裹草稿；必须重新输出会被直接保存的单文件源码，第一行必须是当前 runtime 的源码字符，全文不得出现 ``` 或 ~~~。" if is_script and repair_mode == "strict_contract_rewrite" else "")
            + (f"\n\n后端根据确定性错误生成的必做修复步骤：\n{targeted_repair}" if targeted_repair else "")
        ),
    })
    logger.info(
        "[Creator][model] phase=repair.request file=%s model=%s repair_mode=%s messages=%d previous_chars=%d contract_chars=%d failed_checks_chars=%d",
        file_path,
        model,
        repair_mode,
        len(repair_messages),
        len(previous_content or ""),
        len(contract_text or ""),
        len(failed_checks_text or ""),
    )
    repaired = await complete_chat_once(repair_messages, model)
    if is_script:
        normalized = _normalize_generated_file_content(file_path, repaired)
        logger.info(
            "[Creator][model] phase=repair.response file=%s model=%s repair_mode=%s raw_chars=%d normalized_chars=%d normalized=%s",
            file_path,
            model,
            repair_mode,
            len(repaired or ""),
            len(normalized or ""),
            normalized != repaired,
        )
        return normalized
    logger.info(
        "[Creator][model] phase=repair.response file=%s model=%s repair_mode=%s raw_chars=%d",
        file_path,
        model,
        repair_mode,
        len(repaired or ""),
    )
    return repaired


def _parse_validator_json_object(text: str) -> dict | None:
    stripped = (text or "").strip()
    if stripped.startswith("```json"):
        stripped = stripped[7:].strip()
    if stripped.startswith("```"):
        stripped = stripped[3:].strip()
    if stripped.endswith("```"):
        stripped = stripped[:-3].strip()
    try:
        data = json.loads(stripped)
    except json.JSONDecodeError:
        return None
    return data if isinstance(data, dict) else None


def _deterministic_failed_check_ids(failed_checks_text: str) -> set[str]:
    ids: set[str] = set()
    for match in re.finditer(r"^-\s+([^\s]+)\s+target=", failed_checks_text or "", re.M):
        ids.add(match.group(1))
    return ids


def _filter_validator_failed_checks(failed_checks: list[Any], failed_checks_text: str) -> list[Any]:
    """Keep only model failed_checks backed by deterministic failed checks."""
    allowed_ids = _deterministic_failed_check_ids(failed_checks_text)
    if not allowed_ids:
        return []
    filtered: list[Any] = []
    for check in failed_checks:
        if not isinstance(check, dict):
            continue
        check_id = str(check.get("id") or "")
        if check_id in allowed_ids:
            filtered.append(check)
    return filtered


_MISSING_SKILL_SCRIPT_BLOCK_RE = re.compile(
    r"SKILL\.md 缺少调用 (?P<script>scripts/[A-Za-z0-9_./-]+) 的可执行 Markdown 命令块"
)
_MISSING_SKILL_REFERENCE_RE = re.compile(
    r"SKILL\.md 缺少对参考资料 (?P<reference>references/[A-Za-z0-9_./-]+) 的引用"
)


def _targeted_generated_file_repair_instructions(*, file_path: str, deterministic_error: str) -> str:
    """Return deterministic, actionable instructions for recurring validation failures."""
    if file_path == "SKILL.md":
        script_match = _MISSING_SKILL_SCRIPT_BLOCK_RE.search(deterministic_error)
        if script_match:
            script_path = script_match.group("script")
            return (
                f"必须在 SKILL.md 正文中加入一个真实的 bash fenced code block，且 block 内必须逐字包含 `{script_path}`。\n"
                "注意：‘不要在文件外层套 Markdown 代码块’不等于禁止 SKILL.md 正文内部的命令示例；"
                "这个内部 ```bash block 是校验必需内容。\n"
                "请使用平台占位符 `{{topic}}`，不要使用 shell 变量 `${topic}`。\n"
                "可直接插入如下命令示例，并围绕它补充参数说明：\n"
                "```bash\n"
                f"python {script_path} " "'{\"topic\":\"{{topic}}\"}'" "\n"
                "```"
            )

        reference_match = _MISSING_SKILL_REFERENCE_RE.search(deterministic_error)
        if reference_match:
            reference_path = reference_match.group("reference")
            return (
                f"必须在 SKILL.md 的参考资料/资源小节逐字引用 `{reference_path}`，"
                "并说明何时读取该 reference；不要只泛称‘参考资料’。"
            )

        if "Creator 界面流程" in deterministic_error:
            return (
                "删除 Creator 创建流程、确认清单、点击开始创建、系统将自动创建等 UI 文案；"
                "只保留 Skill 使用说明、资源引用、参数映射和运行时命令示例。"
            )

        if "skill_md.resource.local_declared" in deterministic_error or "未在本轮生成或当前 Skill 中不存在" in deterministic_error:
            return (
                "删除最终 SKILL.md 中未声明/不存在的 references/assets/scripts 引用，尤其是 Creator 内部 kernel references "
                "（references/workflows.md、references/output-patterns.md、references/best-practices.md 等）。"
                "这些 kernel 资料只能作为生成上下文，不能出现在 Skill 的资源列表或运行说明中；"
                "保留合法脚本命令块和合法 text/image helper 说明，不要删除 generate_text_with_llm 或 generate_stable_diffusion_image 相关合法调用。"
            )

    if file_path.startswith("scripts/"):
        if "Markdown 代码块或多文件包" in deterministic_error or "script.raw_source.single_file" in deterministic_error:
            return (
                "本轮必须把上一次内容改成单个裸脚本源码：删除所有 ``` fence、```python/```text 标签、"
                "文件路径标题、写入文件标签、解释性文字和多文件包内容；最终响应第一个字符应是脚本源码字符。"
            )
        if "forbidden_image_generation" in deterministic_error or "调用了图片生成 helper" in deterministic_error:
            return (
                "当前脚本的 SkillPlan forbidden_capabilities 禁止 image_generation，因此 validator 禁止调用 generate_stable_diffusion_image。"
                "蓝图和 SKILL.md 确定后不能由后台修复流程修改；只能修当前脚本，使其与既有 role、required_capabilities 和 SKILL.md 数据流一致。"
            )
        if "script.required_capabilities.called" in deterministic_error or "未调用这些 required_capabilities" in deterministic_error or "没有调用这些 required_capabilities" in deterministic_error:
            return (
                "按当前脚本的 SkillPlan role + 有效 required_capabilities 补齐真实能力调用：包含 image_generation 必须调用 generate_stable_diffusion_image；"
                "包含 text_generation 必须调用 generate_text_with_llm 或平台 LLM；pdf_builder/exporter 默认只需真实创建文件，并在 stdout JSON 任意业务字段中返回路径，不要因 SKILL.md 全局模型说明而补模型调用。"
                "禁止返回固定 f-string/template-only 文本或 placeholder；蓝图和 SKILL.md 确定后只能修当前脚本。"
            )
        if "试运行" in deterministic_error or "JSON 参数" in deterministic_error or "合法 Python" in deterministic_error:
            return (
                "按脚本合同修复：保持单文件源码，修正语法/参数解析/运行错误；"
                "如果 SKILL.md 命令传 JSON，脚本必须读取 sys.argv[1] 并 json.loads，stdout 输出结构化 JSON，至少包含一个非空字段；字段名由现有 SKILL.md 变量消费关系决定，不要修改蓝图或 SKILL.md。"
            )

    if file_path.startswith("assets/"):
        if "contract 未通过" in deterministic_error or "asset" in deterministic_error or "JSON" in deterministic_error:
            return (
                "按 asset 合同修复：只输出当前资源文件内容；确保非空、JSON 可解析，"
                "删除 Creator 流程、多文件包和运行时代码。"
            )

    if file_path.startswith("references/"):
        if "contract 未通过" in deterministic_error or "多文件包" in deterministic_error or "Creator" in deterministic_error:
            return (
                "按 reference 合同修复：只输出当前参考资料 Markdown 正文；"
                "删除 Creator 流程、写入文件标签、多文件包和其它文件路径标题。"
            )

    return ""

async def _run_generated_file_validator_round(
    *,
    file_path: str,
    content: str,
    deterministic_error: str,
    requested_model: str,
    targeted_repair: str = "",
    contract_text: str = "",
    passed_checks_text: str = "",
    failed_checks_text: str = "",
    repair_mode: str = "minimal_edit",
) -> dict:
    """Ask the validator model for actionable repair feedback for the coder."""
    route = route_model(
        VALIDATOR_TASK,
        requested_model=requested_model,
        reason=f"creator generated file validation: {file_path}",
    )
    _log_creator_model_usage(
        phase="validator.route",
        file_path=file_path,
        route=route,
        extra=f"requested_generation_model={requested_model} repair_mode={repair_mode} content_chars={len(content)}",
    )
    messages = [
        {
            "role": "system",
            "content": (
                "你是 Creator 生成文件校验模型，只输出严格 JSON object。"
                "你不改代码，只给 coder 可执行的局部修复意见。"
            ),
        },
        {
            "role": "user",
            "content": (
                f"目标文件：{file_path}\n"
                "后端确定性校验/试运行错误：\n"
                f"{deterministic_error}\n\n"
                + (f"完整 contract：\n{contract_text}\n\n" if contract_text else "")
                + (f"已通过检查（repair 时必须保留）：\n{passed_checks_text}\n\n" if passed_checks_text else "")
                + (f"未通过检查（repair 只修这些项）：\n{failed_checks_text}\n\n" if failed_checks_text else "")
                + (f"本轮修复模式：{repair_mode}\n\n" if repair_mode else "")
                + (f"后端根据该错误生成的必做修复步骤：\n{targeted_repair}\n\n" if targeted_repair else "")
                + "候选内容：\n"
                "```text\n"
                f"{content[-12000:]}\n"
                "```\n\n"
                "请输出 JSON："
                "{\"passed\": false, \"issues\": [\"...\"], "
                "\"failed_checks\": [{\"id\": \"...\", \"target\": \"...\", \"expected\": \"...\", \"minimal_edit\": \"...\"}], "
                "\"preserve\": [\"已通过检查对应内容\"], "
                "\"repair_instructions\": \"给 coder 的局部修改指令；如果上方有必做修复步骤，必须复述并细化这些步骤，不得改用其它占位符或省略必需的 fenced block。若有 Markdown fence/多文件包，明确要求只返回目标脚本源码本身，不要 fence/说明/写入文件标签。\"}"
            ),
        },
    ]
    try:
        text = await complete_chat_once(messages, route.model)
        logger.info(
            "[Creator][model] phase=validator.response file=%s model=%s chars=%d repair_mode=%s",
            file_path,
            route.model,
            len(text or ""),
            repair_mode,
        )
    except Exception as exc:  # pragma: no cover - defensive fallback
        logger.warning("Creator file validator failed; using deterministic feedback: %s", exc)
        return {
            "passed": False,
            "issues": [deterministic_error],
            "repair_instructions": deterministic_error,
            "model": route.model,
        }

    data = _parse_validator_json_object(text)
    if not data:
        return {
            "passed": False,
            "issues": [deterministic_error],
            "repair_instructions": deterministic_error,
            "model": route.model,
        }

    issues = data.get("issues") if isinstance(data.get("issues"), list) else []
    raw_failed_checks = data.get("failed_checks") if isinstance(data.get("failed_checks"), list) else []
    failed_checks = _filter_validator_failed_checks(raw_failed_checks, failed_checks_text)
    preserve = data.get("preserve") if isinstance(data.get("preserve"), list) else []
    instructions = str(data.get("repair_instructions") or data.get("feedback") or deterministic_error)
    filtered_issues, instructions = _filter_validator_model_call_misjudgements(
        file_path=file_path,
        deterministic_error=deterministic_error,
        failed_checks_text=failed_checks_text,
        issues=issues,
        instructions=instructions,
    )
    return {
        "passed": bool(data.get("passed", data.get("valid", False))) and not filtered_issues and not failed_checks,
        "issues": filtered_issues,
        "failed_checks": failed_checks,
        "preserve": [str(item) for item in preserve],
        "repair_instructions": instructions,
        "model": route.model,
    }



def _filter_validator_model_call_misjudgements(
    *,
    file_path: str,
    deterministic_error: str,
    failed_checks_text: str,
    issues: list[Any],
    instructions: str,
) -> tuple[list[str], str]:
    """Remove validator feedback that invents model-call requirements.

    The validator model may over-generalize from SKILL.md prose. Only the
    deterministic contract for this script may introduce required model calls.
    """
    issue_strings = [str(item) for item in issues]
    if not file_path.startswith("scripts/"):
        return issue_strings, instructions
    deterministic_scope = f"{deterministic_error}\n{failed_checks_text}"
    deterministic_mentions_model_requirement = bool(re.search(
        r"text_generation|image_generation|generate_text_with_llm|generate_stable_diffusion_image|LLM|TEXT_MODEL|IMAGE_MODEL|模型",
        deterministic_scope,
        re.IGNORECASE,
    ))
    if deterministic_mentions_model_requirement:
        return issue_strings, instructions
    model_error_re = re.compile(r"必须.*(?:模型|LLM|TEXT_MODEL|IMAGE_MODEL|generate_text_with_llm|generate_stable_diffusion_image)|(?:未|没有)调用.*(?:模型|LLM|TEXT_MODEL|IMAGE_MODEL)", re.IGNORECASE)
    filtered_issues = [item for item in issue_strings if not model_error_re.search(item)]
    if model_error_re.search(instructions or ""):
        instructions = deterministic_error
    return filtered_issues, instructions

def _format_file_validator_feedback(deterministic_error: str, validator_report: dict, targeted_repair: str = "") -> str:
    issues = validator_report.get("issues") or []
    issue_text = "\n".join(f"- {issue}" for issue in issues) if issues else "- （校验模型未返回额外问题）"
    failed_checks = validator_report.get("failed_checks") or []
    failed_check_text = (
        "\n".join(f"- {json.dumps(check, ensure_ascii=False)}" for check in failed_checks)
        if failed_checks else "- （校验模型未返回结构化 failed_checks）"
    )
    return (
        "后端确定性校验/试运行错误：\n"
        f"{deterministic_error}\n\n"
        + (f"后端确定性修复指令：\n{targeted_repair}\n\n" if targeted_repair else "")
        + f"校验模型：{validator_report.get('model', '')}\n"
        "校验模型问题列表：\n"
        f"{issue_text}\n\n"
        "校验模型结构化 failed_checks：\n"
        f"{failed_check_text}\n\n"
        "校验模型给 coder 的修复意见：\n"
        f"{validator_report.get('repair_instructions') or deterministic_error}"
    )


def _is_valid_normalized_script_source(file_path: str, content: str) -> bool:
    """Return whether content is safe to accept as the requested raw script.

    This helper is intentionally narrower than the full script validator: it only
    checks that deterministic Markdown/bundle cleanup produced one raw source
    file.  Full fake-implementation, contract, dependency and trial-run checks
    still happen in the normal validation pipeline.
    """
    stripped = content.strip()
    if not stripped or "```" in stripped or "~~~" in stripped or _MULTI_FILE_MARKER_RE.search(stripped):
        return False

    if Path(file_path).suffix.lower() == ".py":
        try:
            ast.parse(stripped)
        except SyntaxError:
            return False

    return True


def _extract_single_wrapping_fence(content: str) -> str | None:
    """Extract a code block only when it wraps the entire model response.

    Models often wrap repaired scripts in ```text fences, sometimes with CRLF
    line endings or a closing fence that is longer than the opener.  This parser
    intentionally accepts only a whole-response fence: any prose before/after the
    block, or any non-fence trailing line, returns None and lets validation reject
    the ambiguous output.
    """
    stripped = content.strip().lstrip("\ufeff")
    lines = stripped.splitlines()
    if len(lines) < 2:
        return None

    opening = lines[0].strip()
    opening_match = re.match(r"^(`{3,}|~{3,})[^`~]*$", opening)
    if not opening_match:
        return None

    fence = opening_match.group(1)
    fence_char = fence[0]
    min_fence_len = len(fence)
    closing = lines[-1].strip()
    if not re.fullmatch(rf"{re.escape(fence_char)}{{{min_fence_len},}}", closing):
        return None

    return "\n".join(lines[1:-1]).strip()


def _extract_only_fenced_block(content: str) -> str | None:
    """Extract the body only when exactly one fenced block appears in content."""
    lines = content.strip().lstrip("\ufeff").splitlines()
    blocks: list[str] = []
    idx = 0
    while idx < len(lines):
        opening = lines[idx].strip()
        opening_match = re.match(r"^(`{3,}|~{3,})[^`~]*$", opening)
        if not opening_match:
            idx += 1
            continue

        fence = opening_match.group(1)
        fence_char = fence[0]
        min_fence_len = len(fence)
        body: list[str] = []
        idx += 1
        while idx < len(lines):
            closing = lines[idx].strip()
            if re.fullmatch(rf"{re.escape(fence_char)}{{{min_fence_len},}}", closing):
                blocks.append("\n".join(body).strip())
                break
            body.append(lines[idx])
            idx += 1
        else:
            return None

        if len(blocks) > 1:
            return None
        idx += 1

    if len(blocks) != 1:
        return None
    return blocks[0]


def _extract_first_fenced_block(content: str) -> str | None:
    """Extract the first fenced block body, or None if no complete block exists."""
    lines = content.strip().lstrip("\ufeff").splitlines()
    idx = 0
    while idx < len(lines):
        opening = lines[idx].strip()
        opening_match = re.match(r"^(`{3,}|~{3,})[^`~]*$", opening)
        if not opening_match:
            idx += 1
            continue
        fence = opening_match.group(1)
        fence_char = fence[0]
        min_fence_len = len(fence)
        body: list[str] = []
        idx += 1
        while idx < len(lines):
            closing = lines[idx].strip()
            if re.fullmatch(rf"{re.escape(fence_char)}{{{min_fence_len},}}", closing):
                return "\n".join(body).strip()
            body.append(lines[idx])
            idx += 1
        return None
    return None


def _drop_common_non_code_lines(text: str) -> str:
    """Remove common chat/file-label prose that models place around scripts."""
    drop_patterns = [
        r"^\s*下面是",
        r"^\s*以下是",
        r"^\s*(?:文件|路径)[:：]",
        r"^\s*写入文件[:：]",
        r"^\s*#+\s*scripts/",
        r"^\s*`?scripts/[^`]+`?\s*$",
    ]
    cleaned: list[str] = []
    for line in text.strip().splitlines():
        if any(re.search(pattern, line, flags=re.IGNORECASE) for pattern in drop_patterns):
            continue
        cleaned.append(line)
    return "\n".join(cleaned).strip()


def _looks_like_python_source(text: str) -> bool:
    """Return True for probable Python source without requiring valid syntax yet."""
    stripped = text.strip()
    if not stripped or "```" in stripped or "~~~" in stripped or _MULTI_FILE_MARKER_RE.search(stripped):
        return False
    return bool(re.search(
        r"(?m)^\s*(?:import\s+|from\s+|def\s+|class\s+|if __name__\s*==\s*['\"]__main__['\"]|#!/usr/bin/env python|#)",
        stripped,
    ))


def _extract_probable_python_source(content: str) -> str | None:
    """Extract a raw Python source candidate before syntax validation."""
    stripped = content.strip().lstrip("\ufeff")
    candidates: list[str] = []

    normalized = stripped
    for _ in range(3):
        wrapping = _extract_single_wrapping_fence(normalized)
        if wrapping is None:
            break
        normalized = wrapping.strip()
        candidates.append(normalized)

    only_block = _extract_only_fenced_block(stripped)
    if only_block is not None:
        candidates.append(only_block.strip())

    # Use the first block only when the response does not look like an explicit
    # multi-file bundle.  Multi-file bundles are rejected rather than guessed.
    if not _MULTI_FILE_MARKER_RE.search(stripped) and not re.search(r"(?im)^\s*写入文件[:：]", stripped):
        first_block = _extract_first_fenced_block(stripped)
        if first_block is not None:
            candidates.append(first_block.strip())

    candidates.append(stripped)

    seen: set[str] = set()
    for candidate in candidates:
        cleaned = _drop_common_non_code_lines(candidate)
        if cleaned in seen:
            continue
        seen.add(cleaned)
        if _looks_like_python_source(cleaned):
            return cleaned
    return None


def _normalize_generated_file_content(file_path: str, content: str) -> str:
    """Normalize model output while keeping script extraction conservative."""
    if file_path.startswith("scripts/"):
        stripped = content.strip()
        if Path(file_path).suffix.lower() == ".py":
            extracted = _extract_probable_python_source(stripped)
            if extracted:
                return extracted

        normalized = stripped
        for _ in range(3):
            wrapping_fence = _extract_single_wrapping_fence(normalized)
            if wrapping_fence is None:
                break
            normalized = wrapping_fence.strip()
            if _is_valid_normalized_script_source(file_path, normalized):
                return normalized

        only_block = _extract_only_fenced_block(stripped)
        if only_block is not None:
            normalized = only_block.strip()
            if _is_valid_normalized_script_source(file_path, normalized):
                return normalized

        candidate = _strip_orphan_trailing_fence(stripped)
        if _is_valid_normalized_script_source(file_path, candidate):
            return candidate
        return candidate

    extracted = _extract_target_file_from_bundle(content, file_path)
    return _strip_code_fence(extracted if extracted is not None else content)



def _trim_source_to_runtime_entrypoint(file_path: str, content: str, skill_plan_entry: dict[str, Any] | None = None) -> str:
    """Drop leading/trailing prose around a probable runtime source entrypoint."""
    if not file_path.startswith("scripts/"):
        return content.strip()
    plan_entry = _skill_plan_entry_for_file(file_path=file_path, skill_plan_entry=skill_plan_entry)
    text = _strip_orphan_trailing_fence(content.strip().lstrip("\ufeff"))
    lines = text.splitlines()

    start_patterns: list[str]
    end_patterns: list[str]
    if plan_entry.runtime == "node":
        start_patterns = [r"^\s*(?:const|let|var)\s+", r"^\s*function\s+", r"^\s*#!/usr/bin/env\s+node"]
        end_patterns = [r"console\.log\s*\("]
    elif plan_entry.runtime in {"bash", "shell"}:
        start_patterns = [r"^\s*#!/", r"^\s*set\s+-", r"^\s*payload_json=", r"^\s*[A-Za-z_][A-Za-z0-9_]*="]
        end_patterns = [r"^\s*(?:echo|printf|python\s+-c)\b"]
    else:
        start_patterns = [r"^\s*(?:import\s+|from\s+|def\s+|class\s+|#!/usr/bin/env\s+python)"]
        end_patterns = [r"main\s*\(\s*\)"]

    start_idx = 0
    for idx, line in enumerate(lines):
        if any(re.search(pattern, line) for pattern in start_patterns):
            start_idx = idx
            break
    trimmed = lines[start_idx:]

    end_idx = len(trimmed)
    for idx in range(len(trimmed) - 1, -1, -1):
        line = trimmed[idx]
        if any(re.search(pattern, line) for pattern in end_patterns):
            end_idx = idx + 1
            break
    return "\n".join(trimmed[:end_idx]).strip()

def _sanitize_generated_file_content(file_path: str, content: str, role: str | None = None, skill_plan_entry: dict[str, Any] | None = None) -> str:
    """Normalize model output into exactly the requested file content."""
    if file_path.startswith("scripts/") and _MULTI_FILE_MARKER_RE.search(content) and _extract_only_fenced_block(content) is None:
        sanitized = content.strip()
    else:
        sanitized = _normalize_generated_file_content(file_path, content)
        sanitized = _trim_source_to_runtime_entrypoint(file_path, sanitized, skill_plan_entry=skill_plan_entry)
    _validate_generated_file_content(file_path, sanitized, role=role, skill_plan_entry=skill_plan_entry)
    return sanitized


def _strip_code_fence(content: str) -> str:
    """Strip wrapping code-fence markers that a model may output despite instructions.

    Handles:
      ```python\\n<code>\\n```
      ```\\n<code>\\n```
      ~~~~\\n<code>\\n~~~~
    """
    stripped = content.strip()
    # Opening fence with optional language tag
    m_open = re.match(r"^(`{3,}|~{3,})[^\n]*\n", stripped)
    if m_open:
        fence_char = stripped[0]
        rest = stripped[m_open.end():]
        close_pat = re.compile(r"\n" + re.escape(fence_char) + r"{3,}\s*$")
        m_close = close_pat.search(rest)
        if m_close:
            return rest[: m_close.start()].strip()
        return rest.strip()
    return stripped




def _script_generation_skeleton(
    file_path: str,
    purpose: str,
    blueprint_text: str,
    *,
    role: str | None = None,
    skill_plan_entry: dict[str, Any] | None = None,
) -> str:
    """Return a runtime-aware scaffold selected by SkillPlan role/runtime."""
    plan_entry = _skill_plan_entry_for_file(
        file_path=file_path,
        purpose=purpose,
        blueprint_text=blueprint_text,
        role=role,
        skill_plan_entry=skill_plan_entry,
    )
    input_keys = list(plan_entry.inputs or ["payload"])
    effective_required_capabilities = _effective_required_capabilities_for_script(plan_entry)
    py_value_expr = " or ".join(f"payload.get({key!r})" for key in input_keys) + " or ''"
    js_value_expr = " || ".join(f"payload[{json.dumps(key)}]" for key in input_keys) + " || ''"
    bash_py_expr = " or ".join(f"p.get({key!r})" for key in input_keys) + " or ''"

    if plan_entry.runtime == "node":
        if {"text_generation", "image_generation"} <= set(effective_required_capabilities):
            return (
                "必须使用下面的 node composite_generator skeleton；先调用平台 generate_text_with_llm，再调用 generate_stable_diffusion_image，stdout 只能 console.log JSON 字符串：\n"
                "const { spawnSync } = require('child_process');\n"
                "const payload = process.argv[2] ? JSON.parse(process.argv[2]) : {};\n"
                "function pyEval(code, arg) {\n"
                "  const proc = spawnSync(process.env.PYTHON || 'python', ['-c', code, arg], { encoding: 'utf8' });\n"
                "  if (proc.status !== 0) throw new Error(proc.stderr || 'platform helper failed');\n"
                "  return JSON.parse(proc.stdout);\n"
                "}\n"
                "function run(payload) {\n"
                f"  const prompt = String({js_value_expr}).trim();\n"
                "  const textCode = `from backend.services.skill_runtime import generate_text_with_llm\\nimport json,sys\\nprint(json.dumps({'text': generate_text_with_llm(sys.argv[1])}, ensure_ascii=False))`;\n"
                "  const textResult = pyEval(textCode, prompt);\n"
                "  const imagePrompt = textResult.text || prompt;\n"
                "  const imageCode = `from backend.services.skill_runtime import generate_stable_diffusion_image\\nimport json,sys\\nresult = generate_stable_diffusion_image(sys.argv[1], filename_prefix='generated')\\nprint(json.dumps(result, ensure_ascii=False))`;\n"
                "  const imageResult = pyEval(imageCode, imagePrompt);\n"
                "  return { text: textResult.text, text_with_image_prompts: [{ text: textResult.text, image_prompt: imagePrompt }], image_paths: [imageResult.image_path].filter(Boolean), images: [imageResult] };\n"
                "}\n"
                "console.log(JSON.stringify(run(payload)));"
            )
        if plan_entry.role == "image_generator":
            return (
                "必须使用下面的 node image_generator skeleton；通过 Python 平台 helper 生成图片，stdout 只能 console.log JSON 字符串：\n"
                "const { spawnSync } = require('child_process');\n"
                "const payload = process.argv[2] ? JSON.parse(process.argv[2]) : {};\n"
                "function run(payload) {\n"
                f"  const desc = String({js_value_expr}).trim();\n"
                "  const helper = `from backend.services.skill_runtime import generate_stable_diffusion_image\\nimport json,sys\\nresult = generate_stable_diffusion_image(sys.argv[1], filename_prefix='generated')\\nprint(json.dumps(result, ensure_ascii=False))`;\n"
                "  const proc = spawnSync(process.env.PYTHON || 'python', ['-c', helper, desc], { encoding: 'utf8' });\n"
                "  if (proc.status !== 0) throw new Error(proc.stderr || 'generate_stable_diffusion_image failed');\n"
                "  const result = JSON.parse(proc.stdout);\n"
                "  const image_paths = [];\n"
                "  const images = [];\n"
                "  image_paths.push(result.image_path);\n"
                "  images.push(result);\n"
                "  return { image_paths: image_paths.filter(Boolean), images };\n"
                "}\n"
                "console.log(JSON.stringify(run(payload)));"
            )
        if plan_entry.role == "pdf_builder" or "pdf_generation" in set(effective_required_capabilities):
            return (
                "必须使用下面的 node pdf_builder skeleton；消费 text/image_paths/template/assets 构建 PDF，不生成图片：\n"
                "const fs = require('fs');\n"
                "const path = require('path');\n"
                "const payload = process.argv[2] ? JSON.parse(process.argv[2]) : {};\n"
                "function escapePdfText(value) { return String(value).replace(/[\\\\()]/g, '\\\\$&').slice(0, 1800); }\n"
                "function run(payload) {\n"
                f"  const text = String({js_value_expr}).trim() || 'Generated PDF';\n"
                "  const outputDir = path.resolve(payload.output_dir || 'assets/generated');\n"
                "  fs.mkdirSync(outputDir, { recursive: true });\n"
                "  const pdfPath = path.join(outputDir, 'output.pdf');\n"
                "  const body = escapePdfText(text);\n"
                "  const pdf = `%PDF-1.4\\n1 0 obj<<>>endobj\\n2 0 obj<< /Length 44 >>stream\\nBT /F1 12 Tf 50 760 Td (${body}) Tj ET\\nendstream endobj\\n3 0 obj<< /Type /Page /Parent 4 0 R /Contents 2 0 R /Resources << /Font << /F1 5 0 R >> >> >>endobj\\n4 0 obj<< /Type /Pages /Kids [3 0 R] /Count 1 >>endobj\\n5 0 obj<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>endobj\\n6 0 obj<< /Type /Catalog /Pages 4 0 R >>endobj\\ntrailer<< /Root 6 0 R >>\\n%%EOF\\n`;\n"
                "  fs.writeFileSync(pdfPath, pdf);\n"
                "  return { pdf_path: pdfPath, file_paths: [pdfPath] };\n"
                "}\n"
                "console.log(JSON.stringify(run(payload)));"
            )
        if plan_entry.role == "text_generator":
            return (
                "必须使用下面的 node text_generator skeleton；调用平台 generate_text_with_llm helper，stdout JSON 包含非空 text：\n"
                "const { spawnSync } = require('child_process');\n"
                "const payload = process.argv[2] ? JSON.parse(process.argv[2]) : {};\n"
                "function run(payload) {\n"
                f"  const prompt = String({js_value_expr}).trim();\n"
                "  const helper = `from backend.services.skill_runtime import generate_text_with_llm\\nimport json,sys\\nprint(json.dumps({'text': generate_text_with_llm(sys.argv[1])}, ensure_ascii=False))`;\n"
                "  const proc = spawnSync(process.env.PYTHON || 'python', ['-c', helper, prompt], { encoding: 'utf8' });\n"
                "  if (proc.status !== 0) throw new Error(proc.stderr || 'generate_text_with_llm failed');\n"
                "  return JSON.parse(proc.stdout);\n"
                "}\n"
                "console.log(JSON.stringify(run(payload)));"
            )
        return (
            "必须使用下面的 node_skeleton；解析 process.argv[2] JSON，stdout 只能 console.log JSON 字符串：\n"
            "const payload = process.argv[2] ? JSON.parse(process.argv[2]) : {};\n"
            "function run(payload) {\n"
            f"  const text = String({js_value_expr}).trim();\n"
            "  return { text, file_paths: [] };\n"
            "}\n"
            "console.log(JSON.stringify(run(payload)));"
        )

    if plan_entry.runtime in {"bash", "shell"}:
        if {"text_generation", "image_generation"} <= set(effective_required_capabilities):
            helper = "from backend.services.skill_runtime import generate_text_with_llm, generate_stable_diffusion_image; import json,sys; p=json.loads(sys.argv[1]); prompt=str(" + bash_py_expr + "); text=generate_text_with_llm(prompt); result=generate_stable_diffusion_image(text or prompt, filename_prefix='generated'); print(json.dumps({'story_text': text, 'text': text, 'text_with_image_prompts': [{'text': text, 'image_prompt': text or prompt}], 'image_paths':[result.get('image_path')], 'images':[result]}, ensure_ascii=False))"
        elif plan_entry.role == "image_generator":
            helper = "from backend.services.skill_runtime import generate_stable_diffusion_image; import json,sys; result=generate_stable_diffusion_image(sys.argv[1], filename_prefix='generated'); print(json.dumps({'image_paths':[result.get('image_path')], 'images':[result]}, ensure_ascii=False))"
        elif plan_entry.role == "pdf_builder" or "pdf_generation" in set(effective_required_capabilities):
            helper = "import json,sys; from pathlib import Path; p=json.loads(sys.argv[1]); text=str(" + bash_py_expr + " or 'Generated PDF'); out=Path(p.get('output_dir') or 'assets/generated').resolve(); out.mkdir(parents=True, exist_ok=True); pdf=out/'output.pdf'; pdf.write_text('%PDF-1.4\\nBT ('+text[:1000].replace('(',' ').replace(')',' ') +') Tj ET\\n%%EOF\\n', encoding='latin1'); print(json.dumps({'pdf_path':str(pdf),'file_paths':[str(pdf)]}, ensure_ascii=False))"
        elif plan_entry.role == "text_generator":
            helper = "from backend.services.skill_runtime import generate_text_with_llm; import json,sys; p=json.loads(sys.argv[1]); prompt=str(" + bash_py_expr + "); print(json.dumps({'text': generate_text_with_llm(prompt)}, ensure_ascii=False))"
        else:
            helper = "import json,sys; p=json.loads(sys.argv[1]); text=str(" + bash_py_expr + "); print(json.dumps({'text': text, 'file_paths': []}, ensure_ascii=False))"
        return (
            "必须使用下面的 shell_skeleton；从 $1 读取 JSON argv，并向 stdout 输出 JSON：\n"
            "#!/usr/bin/env bash\n"
            "set -euo pipefail\n"
            "payload_json=${1:-'{}'}\n"
            f"python -c {shlex.quote(helper)} \"$payload_json\""
        )


    if plan_entry.role == "docx_builder" or "docx_generation" in set(effective_required_capabilities):
        return (
            "必须使用下面的 docx_builder 脚本骨架；默认只消费已有 stdout JSON/text/image_paths 并生成 Word；仅当当前脚本 capabilities 显式声明 text/image generation 时才调用模型：\n"
            "import json\n"
            "import sys\n"
            "from pathlib import Path\n"
            "from zipfile import ZipFile, ZIP_DEFLATED\n\n"
            "def parse_args() -> dict:\n"
            "    if len(sys.argv) < 2:\n"
            "        return {}\n"
            "    return json.loads(sys.argv[1])\n\n"
            "def previous_payload(payload: dict) -> dict:\n"
            "    raw = payload.get('previous_stdout') or payload.get('stdout_json') or '{}'\n"
            "    if isinstance(raw, dict):\n"
            "        return raw\n"
            "    try:\n"
            "        data = json.loads(str(raw))\n"
            "        return data if isinstance(data, dict) else {}\n"
            "    except json.JSONDecodeError:\n"
            "        return {}\n\n"
            "def build_docx(payload: dict) -> dict:\n"
            "    prev = previous_payload(payload)\n"
            f"    text = str(payload.get('story_text') or payload.get('text') or prev.get('story_text') or prev.get('text') or {py_value_expr} or 'Generated document').strip()\n"
            "    out_dir = Path(payload.get('output_dir') or 'assets/generated').resolve()\n"
            "    out_dir.mkdir(parents=True, exist_ok=True)\n"
            "    docx_path = out_dir / 'output.docx'\n"
            "    body = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')\n"
            "    document_xml = '<?xml version=\"1.0\" encoding=\"UTF-8\"?><w:document xmlns:w=\"http://schemas.openxmlformats.org/wordprocessingml/2006/main\"><w:body><w:p><w:r><w:t>' + body + '</w:t></w:r></w:p></w:body></w:document>'\n"
            "    with ZipFile(docx_path, 'w', ZIP_DEFLATED) as zf:\n"
            "        zf.writestr('[Content_Types].xml', '<?xml version=\"1.0\" encoding=\"UTF-8\"?><Types xmlns=\"http://schemas.openxmlformats.org/package/2006/content-types\"><Default Extension=\"rels\" ContentType=\"application/vnd.openxmlformats-package.relationships+xml\"/><Default Extension=\"xml\" ContentType=\"application/xml\"/><Override PartName=\"/word/document.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml\"/></Types>')\n"
            "        zf.writestr('_rels/.rels', '<?xml version=\"1.0\" encoding=\"UTF-8\"?><Relationships xmlns=\"http://schemas.openxmlformats.org/package/2006/relationships\"><Relationship Id=\"rId1\" Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument\" Target=\"word/document.xml\"/></Relationships>')\n"
            "        zf.writestr('word/document.xml', document_xml)\n"
            "    return {'docx_path': str(docx_path), 'file_paths': [str(docx_path)]}\n\n"
            "def main() -> None:\n"
            "    print(json.dumps(build_docx(parse_args()), ensure_ascii=False))\n\n"
            "if __name__ == '__main__':\n"
            "    main()"
        )

    if plan_entry.role == "pptx_builder" or "pptx_generation" in set(effective_required_capabilities):
        return (
            "必须使用下面的 pptx_builder 脚本骨架；默认只消费已有 stdout JSON/text/image_paths 并生成 PPT；仅当当前脚本 capabilities 显式声明 text/image generation 时才调用模型：\n"
            "import json\n"
            "import sys\n"
            "from pathlib import Path\n\n"
            "def parse_args() -> dict:\n"
            "    if len(sys.argv) < 2:\n"
            "        return {}\n"
            "    return json.loads(sys.argv[1])\n\n"
            "def build_pptx(payload: dict) -> dict:\n"
            f"    text = str(payload.get('story_text') or payload.get('text') or {py_value_expr} or 'Generated presentation').strip()\n"
            "    out_dir = Path(payload.get('output_dir') or 'assets/generated').resolve()\n"
            "    out_dir.mkdir(parents=True, exist_ok=True)\n"
            "    pptx_path = out_dir / 'output.pptx'\n"
            "    from pptx import Presentation\n"
            "    prs = Presentation()\n"
            "    slide = prs.slides.add_slide(prs.slide_layouts[1])\n"
            "    slide.shapes.title.text = 'Generated'\n"
            "    slide.placeholders[1].text = text[:2000]\n"
            "    prs.save(pptx_path)\n"
            "    return {'pptx_path': str(pptx_path), 'file_paths': [str(pptx_path)]}\n\n"
            "def main() -> None:\n"
            "    print(json.dumps(build_pptx(parse_args()), ensure_ascii=False))\n\n"
            "if __name__ == '__main__':\n"
            "    main()"
        )

    if plan_entry.role in {"html_asset_builder", "asset_builder"} or ({"html_generation", "html_asset_generation"} & set(effective_required_capabilities)):
        return (
            "必须使用下面的 html_asset_builder Python 脚本骨架；只能在当前 Skill 的 assets/generated/ 下写入 HTML，并在 stdout JSON 返回 html_path 与 asset_paths：\n"
            "import html\n"
            "import json\n"
            "import re\n"
            "import sys\n"
            "from pathlib import Path\n\n"
            "def parse_args() -> dict:\n"
            "    if len(sys.argv) < 2:\n"
            "        return {}\n"
            "    return json.loads(sys.argv[1])\n\n"
            "def slugify(value: str) -> str:\n"
            "    slug = re.sub(r'[^A-Za-z0-9_-]+', '-', value).strip('-').lower()\n"
            "    return slug or 'generated'\n\n"
            "def build_html(payload: dict) -> str:\n"
            f"    text = str({py_value_expr} or 'Generated HTML').strip()\n"
            "    safe = html.escape(text)\n"
            "    return '<!doctype html><html><head><meta charset=\"utf-8\"><title>Generated</title></head><body><main><h1>Generated Asset</h1><p>' + safe + '</p></main></body></html>'\n\n"
            "def run(payload: dict) -> dict:\n"
            f"    title = str({py_value_expr} or 'generated').strip()\n"
            "    skill_root = Path(__file__).resolve().parents[1]\n"
            "    out_dir = (skill_root / 'assets' / 'generated').resolve()\n"
            "    required_root = (skill_root / 'assets' / 'generated').resolve()\n"
            "    out_dir.mkdir(parents=True, exist_ok=True)\n"
            "    html_path = (out_dir / (slugify(title) + '.html')).resolve()\n"
            "    html_path.relative_to(required_root)\n"
            "    html_path.write_text(build_html(payload), encoding='utf-8')\n"
            "    rel = html_path.relative_to(skill_root).as_posix()\n"
            "    return {'html_path': rel, 'asset_paths': [rel]}\n\n"
            "def main() -> None:\n"
            "    payload = parse_args()\n"
            "    print(json.dumps(run(payload), ensure_ascii=False))\n\n"
            "if __name__ == '__main__':\n"
            "    main()"
        )

    if {"text_generation", "image_generation"} <= set(effective_required_capabilities):
        return (
            "必须使用下面的 composite_generator 脚本骨架；先调用平台 generate_text_with_llm，再调用 generate_stable_diffusion_image，stdout JSON 包含 text 与 image_paths/images：\n"
            "import json\n"
            "import sys\n"
            "from backend.services.skill_runtime import generate_text_with_llm, generate_stable_diffusion_image\n\n"
            "def parse_args() -> dict:\n"
            "    if len(sys.argv) < 2:\n"
            "        return {}\n"
            "    return json.loads(sys.argv[1])\n\n"
            "def build_prompt(payload: dict) -> str:\n"
            f"    return str({py_value_expr}).strip()\n\n"
            "def generate_text(prompt: str) -> str:\n"
            "    return generate_text_with_llm(prompt).strip()\n\n"
            "def generate_images(text: str, prompt: str) -> tuple[list[str], list[dict]]:\n"
            "    image_prompt = text or prompt\n"
            "    result = generate_stable_diffusion_image(image_prompt, filename_prefix='generated')\n"
            "    image_paths = [result.get('image_path')]\n"
            "    image_paths = [p for p in image_paths if isinstance(p, str) and p]\n"
            "    return image_paths, [result]\n\n"
            "def run(payload: dict) -> dict:\n"
            "    prompt = build_prompt(payload)\n"
            "    text = generate_text(prompt)\n"
            "    image_paths, images = generate_images(text, prompt)\n"
            "    return {'story_text': text, 'text': text, 'text_with_image_prompts': [{'text': text, 'image_prompt': text or prompt}], 'image_paths': image_paths, 'images': images}\n\n"
            "def main() -> None:\n"
            "    payload = parse_args()\n"
            "    print(json.dumps(run(payload), ensure_ascii=False))\n\n"
            "if __name__ == '__main__':\n"
            "    main()"
        )

    if plan_entry.role == "image_generator":
        return (
            "必须使用下面的 image_generator 脚本骨架；不要改变 import/helper/main/JSON stdout 结构，只填充 build_image_prompt() 中的业务 prompt 组装逻辑，必要时补充返回字段：\n"
            "import json\n"
            "import sys\n"
            "from backend.services.skill_runtime import generate_stable_diffusion_image\n\n"
            "def parse_args() -> dict:\n"
            "    if len(sys.argv) < 2:\n"
            "        return {}\n"
            "    return json.loads(sys.argv[1])\n\n"
            "def build_image_prompt(payload: dict) -> str:\n"
            f"    topic = str({py_value_expr}).strip()\n"
            "    return topic\n\n"
            "def run(payload: dict) -> dict:\n"
            "    desc = build_image_prompt(payload)\n"
            "    image_paths = []\n"
            "    images = []\n"
            "    result = generate_stable_diffusion_image(desc, filename_prefix='generated')\n"
            "    image_paths.append(result.get('image_path'))\n"
            "    images.append(result)\n"
            "    image_paths = [p for p in image_paths if isinstance(p, str) and p]\n"
            "    return {'image_paths': image_paths, 'images': images}\n\n"
            "def main() -> None:\n"
            "    payload = parse_args()\n"
            "    print(json.dumps(run(payload), ensure_ascii=False))\n\n"
            "if __name__ == '__main__':\n"
            "    main()"
        )

    if plan_entry.role == "pdf_builder" or "pdf_generation" in set(effective_required_capabilities):
        return (
            "必须使用下面的 pdf_builder 脚本骨架；默认只负责读取已有内容并构建 PDF/Word/PPT/HTML 文件；仅当当前脚本 capabilities 显式声明 text/image generation 时才调用模型：\n"
            "import json\n"
            "import sys\n"
            "from pathlib import Path\n\n"
            "def parse_args() -> dict:\n"
            "    if len(sys.argv) < 2:\n"
            "        return {}\n"
            "    return json.loads(sys.argv[1])\n\n"
            "def build_pdf(payload: dict) -> dict:\n"
            "    output_dir = Path(payload.get('output_dir') or 'assets/generated').resolve()\n"
            "    output_dir.mkdir(parents=True, exist_ok=True)\n"
            "    pdf_path = output_dir / 'output.pdf'\n"
            "    from fpdf import FPDF\n"
            f"    text = str({py_value_expr} or 'Generated PDF').strip()\n"
            "    pdf = FPDF()\n"
            "    pdf.add_page()\n"
            "    pdf.set_font('Helvetica', size=14)\n"
            "    pdf.multi_cell(0, 10, text[:2000])\n"
            "    pdf.output(str(pdf_path))\n"
            "    return {'pdf_path': str(pdf_path), 'file_paths': [str(pdf_path)]}\n\n"
            "def run(payload: dict) -> dict:\n"
            "    return build_pdf(payload)\n\n"
            "def main() -> None:\n"
            "    payload = parse_args()\n"
            "    print(json.dumps(run(payload), ensure_ascii=False))\n\n"
            "if __name__ == '__main__':\n"
            "    main()"
        )

    if plan_entry.role == "text_generator":
        return (
            "必须使用下面的 text_generator 脚本骨架；调用平台 generate_text_with_llm，stdout JSON 必须包含非空 text，不要生成图片或 PDF：\n"
            "import json\n"
            "import sys\n"
            "from backend.services.skill_runtime import generate_text_with_llm\n\n"
            "def parse_args() -> dict:\n"
            "    if len(sys.argv) < 2:\n"
            "        return {}\n"
            "    return json.loads(sys.argv[1])\n\n"
            "def generate_text(payload: dict) -> str:\n"
            f"    prompt = str({py_value_expr}).strip()\n"
            "    return generate_text_with_llm(prompt)\n\n"
            "def run(payload: dict) -> dict:\n"
            "    text = generate_text(payload).strip()\n"
            "    return {'story_text': text, 'text': text}\n\n"
            "def main() -> None:\n"
            "    payload = parse_args()\n"
            "    print(json.dumps(run(payload), ensure_ascii=False))\n\n"
            "if __name__ == '__main__':\n"
            "    main()"
        )

    return (
        "必须使用下面的 generic_script 脚本骨架；不要改变 import/parse_args/main/JSON stdout 结构，只填充 run() 中的真实业务逻辑并按 SkillPlan 使用 payload 字段：\n"
        "import json\n"
        "import sys\n\n"
        "def parse_args() -> dict:\n"
        "    if len(sys.argv) < 2:\n"
        "        return {}\n"
        "    return json.loads(sys.argv[1])\n\n"
        "def run(payload: dict) -> dict:\n"
        f"    text = str({py_value_expr}).strip()\n"
        "    return {'text': text, 'file_paths': []}\n\n"
        "def main() -> None:\n"
        "    payload = parse_args()\n"
        "    print(json.dumps(run(payload), ensure_ascii=False))\n\n"
        "if __name__ == '__main__':\n"
        "    main()"
    )


def _creator_kernel_reference_context() -> str:
    """Load small Creator prompt context from kernel references.

    These references are advisory generation context only; SKILL.md protocol and
    generated file contracts remain unchanged.
    """
    kernel_dir = Path(__file__).resolve().parents[2] / "kernel"
    candidates = [
        kernel_dir / "references" / "best-practices.md",
        kernel_dir / "references" / "workflows.md",
        kernel_dir / "references" / "output-patterns.md",
        kernel_dir / "SKILL.md",
    ]
    chunks: list[str] = []
    for path in candidates:
        if not path.is_file():
            continue
        try:
            text = path.read_text(encoding="utf-8").strip()
        except OSError:
            continue
        if text:
            rel = path.relative_to(kernel_dir.parent)
            chunks.append(f"### INTERNAL-ONLY {rel}\n{text[:1800]}")
    return "\n\n".join(chunks)

def _build_generate_file_prompt(
    file_path: str,
    skill_name: str,
    purpose: str,
    blueprint_text: str,
    conversation_history: list[dict],
    role: str | None = None,
    skill_plan_entry: dict[str, Any] | None = None,
) -> list[dict]:
    """Build a minimal generation prompt for a single Skill file.

    The model is asked to output *only* raw file content — no fences, no JSON,
    no explanations.  This maximises reliability for small or unstable models.
    """
    ext = Path(file_path).suffix.lower()
    lang = _LANG_LABELS.get(ext, "文本")

    clean_blueprint_text = _clean_blueprint_for_file_prompt(blueprint_text)
    declared_paths = _extract_declared_skill_paths(blueprint_text)
    declared_paths_text = "\n".join(f"- {path}" for path in declared_paths) or "- （蓝图未显式列出资源文件）"
    plan_entry = _skill_plan_entry_for_file(
        file_path=file_path, purpose=purpose, blueprint_text=blueprint_text, role=role, skill_plan_entry=skill_plan_entry
    )
    generated_file_contract_text = _build_generated_file_contract_text(
        file_path, blueprint_text, purpose, role=role, skill_plan_entry=skill_plan_entry
    )
    skill_md_contract_text = generated_file_contract_text if file_path == "SKILL.md" else ""
    script_skeleton_text = _script_generation_skeleton(
        file_path,
        purpose,
        blueprint_text,
        role=plan_entry.role,
        skill_plan_entry=skill_plan_entry,
    ) if file_path.startswith("scripts/") else ""
    kernel_reference_context = _creator_kernel_reference_context()
    plan_summary = (
        f"SkillPlan role：{plan_entry.role}；"
        f"inputs：{', '.join(plan_entry.inputs)}；"
        f"outputs：{', '.join(plan_entry.outputs)}；"
        f"language：{plan_entry.language}；"
        f"runtime：{plan_entry.runtime}；"
        f"command_template：{_script_command_template(file_path, blueprint_text, plan_entry)}；"
        f"forbidden_capabilities：{', '.join(plan_entry.forbidden_capabilities)}"
    )

    if file_path == "SKILL.md":
        instruction = (
            f'你正在为 Skill 包 "{skill_name}" 生成 SKILL.md 文件。\n\n'
            "要求：\n"
            "1. 只输出 SKILL.md 的文件内容，不要任何解释，不要 Markdown 代码块包裹。\n"
            "2. 文件必须以 YAML frontmatter 开始，格式严格如下（冒号后有一个空格）：\n"
            "---\n"
            f"name: {skill_name}\n"
            "description: <一句话说明本 Skill 的用途>\n"
            "---\n"
            "3. frontmatter 闭合后，输出 Skill 的核心执行说明（普通 Markdown 正文）。\n"
            "4. SKILL.md 必须作为复合任务 orchestrator：描述执行顺序、上一步 outputs 如何传给下一步 inputs、每步 expected outputs、失败/跳过条件，以及何时读取 references/assets。\n"
            "5. SKILL.md 只写总流程和调用顺序；子任务详细规则必须引用 references/*.md，不要把 text/image/pdf/docx/pptx/html 子任务规范混在 SKILL.md 中。\n"
            "6. 如果蓝图包含 scripts/ 资源，SKILL.md 正文必须为每个 scripts/ 路径提供一个可执行的 ```bash fenced code block，命令参数必须与脚本接口一致。\n"
            "7. 如果蓝图包含 references/ 资源，SKILL.md 正文必须在“参考资料/资源”小节明确引用每个 references/ 路径，并说明何时读取。\n"
            "8. 不要在输出内容的外侧套 ``` 代码块，但 SKILL.md 正文内部必须按需包含示例 ```bash fenced code block。\n"
            "9. 禁止只写‘立即调用 `scripts/...`’这种隐式执行描述；必须写明 assistant 应输出可执行 fenced block。\n"
            "10. 禁止复制 Creator 界面流程、确认清单、‘点击开始创建/开始生成’、系统将自动创建文件等平台创建流程文案。\n"
            "11. 以下宿主 Markdown 执行说明是内部写作约束，只能转化为面向使用者的 Skill 说明，不要逐字复制这些约束或标题。\n"
            f"{_SKILL_MD_MARKDOWN_EXECUTION_GUIDE}\n"
            "生成前请先隐式检查以下合同，最终输出必须逐项满足；如果合同要求内部 ```bash block，必须在 SKILL.md 正文中写出该 block：\n"
            f"{skill_md_contract_text}\n\n"
            f"蓝图声明的文件路径（必须覆盖对应 scripts/references 要求）：\n{declared_paths_text}\n\n"
            f"以下是已确认的蓝图（已移除 Creator UI 确认文案），你的内容必须与此一致：\n\n{clean_blueprint_text}"
        )
    elif file_path.startswith("scripts/"):
        instruction = (
            f'你正在为 Skill 包 "{skill_name}" 生成 {file_path} 文件。\n\n'
            f"职责说明：{purpose}\n"
            f"{plan_summary}\n\n"
            "你是文件内容生成器，不是聊天助手；当前输出会被直接写入目标文件。\n"
            "要求：\n"
            f"1. 只输出完整可运行的 {lang} 文件字节内容本身，不要任何说明文字。\n"
            "2. 禁止 Markdown，禁止 ```，禁止‘下面是代码’，禁止文件名标题，禁止多文件输出；如果输出包含 ```，系统会判定失败。\n"
            "3. 脚本的命令行参数、stdin/stdout 接口必须与蓝图和 SKILL.md 里的 Markdown 命令示例一致。\n"
            "4. 如果命令示例传入 JSON 字符串参数，脚本必须按 SkillPlan.runtime 解析；Python 默认读取 sys.argv[1] 并 json.loads 解析，Node 使用 process.argv[2]+JSON.parse，Bash 使用 $1 JSON。\n"
            "5. 必须实际使用用户可变参数生成结果；禁止把示例结果、示例标题、示例图片路径硬编码成固定输出。\n"
            "6. 文本/代码/视觉理解与图片生成的模型来源必须区分：text_generation 使用 generate_text_with_llm 或 LLM_BASE_URL + TEXT_MODEL；看图/OCR/多模态理解使用 LLM_BASE_URL + VISION_MODEL；image_generation 使用平台 Stable Diffusion 图片运行时（IMAGE_BASE_URL + IMAGE_MODEL），不要把 VISION_MODEL 用于图片生成。\n"
            "7. 是否必须调用文本/图片模型只由当前脚本 SkillPlan.required_capabilities 决定：包含 image_generation 时必须调用 `from backend.services.skill_runtime import generate_stable_diffusion_image`；builder/exporter 默认是确定性文件构建脚本，不要因为整个 Skill.md 提到模型就强制 builder 调模型；若 builder 需要模型辅助，用 optional_capabilities/allowed_capabilities 或显式 required_capabilities 表达。\n"
            "8. image_generation stdout 输出结构化 JSON，并返回 helper 结果里的 image_paths/images；必须使用 result = generate_stable_diffusion_image(desc)、image_paths.append(result.get(\"image_path\"))、images.append(result) 的骨架，禁止 image_path = generate_stable_diffusion_image(...)；禁止输出 base64 data URI，禁止假设接口只返回 url；可按需读取平台注入的 IMAGE_MODEL / IMAGE_BASE_URL / IMAGE_SIZE / IMAGE_API_KEY 等环境变量，但不要硬编码，也不需要额外校验它们是否存在。\n"
            "9. 如果脚本只做确定性计算、转换、文件处理或格式化，必须实现真实算法并使用用户输入；禁止假 API、placeholder 文件、纯色/空白图片或 ASCII 图冒充输出。\n"
            "10. stdout 输出结构化 JSON，字段必须以 SkillPlan.outputs 为准；允许极少数 legacy alias 仅用于兼容旧 Skill，新生成 Skill 必须统一字段命名，不要混用 text/text_content/story_text 等别名；composite_generator 只表示多能力组合，具体输出由 outputs/required_capabilities 决定。\n"
            "11. 所有导入的第三方库必须真实存在且常见；Creator 保存前会先扫描 Python import 并安装缺失依赖，再按“生成→测试→修复生成→再测试”的闭环试运行；脚本仍必须包含必要的错误处理逻辑（如参数校验、文件不存在提示等）。\n"
            "12. 必须基于下方固定骨架生成：默认优先 Python；若 SkillPlan.runtime 为 node/bash，则使用对应骨架并保留入口、参数解析和 JSON stdout。\n"
            f"13. 最终响应必须是单个 {plan_entry.language} 源码文件；去掉 Markdown fence、说明文字、文件路径标题和多文件包。\n"
            "生成前请先隐式检查以下脚本合同，最终输出必须逐项满足：\n"
            f"{generated_file_contract_text}\n\n"
            f"固定脚本骨架（仅用于约束生成结构；输出时应是补全后的源码，不要保留空实现）：\n{script_skeleton_text}\n\n"
            f"Creator internal-only kernel guidance（只可指导生成，禁止把以下 kernel 路径、文件名或组织结构复制到最终业务 SKILL.md / references / assets）：\n{kernel_reference_context}\n\n"
            f"蓝图声明的文件路径：\n{declared_paths_text}\n\n"
            f"以下是已确认的蓝图（scripts/ 生成不会追加聊天历史，只使用本蓝图）：\n\n{clean_blueprint_text}"
        )
    elif file_path.startswith("references/"):
        instruction = (
            f'你正在为 Skill 包 "{skill_name}" 生成 {file_path} 参考资料文件。\n\n'
            f"职责说明：{purpose}\n\n"
            "要求：\n"
            "1. 只输出 Markdown 文档内容，不要额外的说明文字。\n"
            "2. 不要在文档外套 ``` 代码块。\n"
            "3. 内容应是有实际指导价值的参考资料，不是对参考资料的再描述。\n"
            "生成前请先隐式检查以下 reference 合同，最终输出必须逐项满足：\n"
            f"{generated_file_contract_text}\n\n"
            f"蓝图声明的文件路径：\n{declared_paths_text}\n\n"
            f"以下是已确认的蓝图（参考资料职责说明见 references/ 部分）：\n\n{clean_blueprint_text}"
        )
    elif file_path.startswith("assets/"):
        instruction = (
            f'你正在为 Skill 包 "{skill_name}" 生成 {file_path} 资源文件。\n\n'
            f"职责说明：{purpose}\n\n"
            "要求：\n"
            f"1. 只输出 {lang} 格式的文件内容，不要任何说明文字。\n"
            "2. 不要用 ``` 代码块包裹输出。\n"
            "生成前请先隐式检查以下 asset 合同，最终输出必须逐项满足：\n"
            f"{generated_file_contract_text}\n\n"
            f"蓝图声明的文件路径：\n{declared_paths_text}\n\n"
            f"以下是已确认的蓝图：\n\n{clean_blueprint_text}"
        )
    else:
        instruction = (
            f'你正在为 Skill 包 "{skill_name}" 生成 {file_path} 文件。\n\n'
            f"职责说明：{purpose}\n\n"
            "要求：直接输出文件内容，不要任何解释，不要 Markdown 代码块包裹。\n\n"
            f"蓝图声明的文件路径：\n{declared_paths_text}\n\n"
            f"蓝图：\n\n{clean_blueprint_text}"
        )

    messages: list[dict] = [{"role": "system", "content": instruction}]

    if file_path.startswith("scripts/"):
        # Scripts are generated from the system instruction plus the confirmed
        # blueprint only.  Do not append conversation history: recent Creator UI
        # copy (file-list previews, confirmation instructions, panel messages)
        # has repeatedly polluted first-pass script output.
        return messages

    # Include recent user context but skip Creator UI confirmation text. Assistant
    # blueprint confirmations often contain "click Start" operational prose that
    # must never be copied into generated files.
    for msg in conversation_history[-_MAX_HISTORY_TURNS:]:
        if not isinstance(msg, dict) or msg.get("role") not in {"user", "assistant"}:
            continue
        content = str(msg.get("content") or "")
        if msg.get("role") == "assistant" and _CREATOR_FLOW_LEAK_RE.search(content):
            continue
        messages.append({**msg, "content": _clean_blueprint_for_file_prompt(content)})

    return messages


def _sse(data: dict) -> str:
    return f"data: {json.dumps(data, ensure_ascii=False)}\n\n"


# ---------------------------------------------------------------------------
# Endpoints
# ---------------------------------------------------------------------------

@router.post("/analyze-blueprint", response_model=AnalyzeBlueprintResponse)
async def analyze_blueprint(request: AnalyzeBlueprintRequest):
    """Extract a file-creation plan from the conversation blueprint.

    Pure rule-based extraction — no LLM call is made.
    """
    plan: BlueprintPlan = parse_blueprint(request.messages)
    entries_by_path = {entry.path: entry for entry in (plan.skill_plan.files if plan.skill_plan else [])}
    return AnalyzeBlueprintResponse(
        skill_name=plan.skill_name,
        files=[
            FileSpecOut(
                path=f.path,
                purpose=f.purpose,
                required=f.required,
                can_skip=f.can_skip,
                file_type=entries_by_path[f.path].file_type if f.path in entries_by_path else None,
                role=entries_by_path[f.path].role if f.path in entries_by_path else None,
                inputs=entries_by_path[f.path].inputs if f.path in entries_by_path else [],
                outputs=entries_by_path[f.path].outputs if f.path in entries_by_path else [],
                dependencies=entries_by_path[f.path].dependencies if f.path in entries_by_path else [],
                required_capabilities=entries_by_path[f.path].required_capabilities if f.path in entries_by_path else [],
                forbidden_capabilities=entries_by_path[f.path].forbidden_capabilities if f.path in entries_by_path else [],
                reference_files=entries_by_path[f.path].reference_files if f.path in entries_by_path else [],
                skill_local_references=entries_by_path[f.path].skill_local_references if f.path in entries_by_path else [],
                creator_internal_references=entries_by_path[f.path].creator_internal_references if f.path in entries_by_path else [],
                language=entries_by_path[f.path].language if f.path in entries_by_path else "text",
                runtime=entries_by_path[f.path].runtime if f.path in entries_by_path else "none",
                entrypoint=entries_by_path[f.path].entrypoint if f.path in entries_by_path else "",
                command_template=entries_by_path[f.path].command_template if f.path in entries_by_path else "",
                references=entries_by_path[f.path].reference_files if f.path in entries_by_path else [],
                low_confidence=(entries_by_path[f.path].confidence < 0.7) if f.path in entries_by_path else False,
                confidence=entries_by_path[f.path].confidence if f.path in entries_by_path else 0.0,
                reason=entries_by_path[f.path].reason if f.path in entries_by_path else "",
                heuristic_signals=entries_by_path[f.path].heuristic_signals if f.path in entries_by_path else [],
            )
            for f in plan.files
        ],
        warnings=plan.warnings,
    )


@router.post("/init-skill", response_model=InitSkillResponse)
async def init_skill(request: InitSkillRequest):
    """Initialise a new Skill directory structure."""
    skill_name = _validate_skill_name(request.skill_name)
    result = run_action({"action": "init", "name": skill_name})
    return InitSkillResponse(
        success=result["success"],
        path=result.get("path"),
        message=result["message"],
    )


@router.post("/generate-file")
async def generate_file(request: GenerateFileRequest):
    """Stream the generated content for a single Skill file.

    SSE event shapes:
      ``{ "content": "<chunk>" }``  — content chunk
      ``{ "done": true }``           — generation complete
      ``{ "error": "<message>" }``   — generation failed
    """
    _validate_file_path(request.file_path)
    skill_name = _validate_skill_name(request.skill_name)
    route = route_creator_file_model(
        file_path=request.file_path,
        purpose=request.purpose,
        requested_model=request.model,
    )
    model = route.model
    _log_creator_model_usage(
        phase="generate.route",
        skill_name=skill_name,
        file_path=request.file_path,
        route=route,
        extra=f"purpose_chars={len(request.purpose or '')} requested_ui_model={request.model or ''}",
    )

    prompt_messages = _build_generate_file_prompt(
        file_path=request.file_path,
        skill_name=skill_name,
        purpose=request.purpose,
        blueprint_text=request.blueprint_text,
        conversation_history=request.conversation_history,
        role=request.role,
        skill_plan_entry=request.skill_plan_entry,
    )

    async def event_stream():
        last_validation_error = ""
        repair_attempts = 0
        try:
            yield _sse({"model_ack": route.ack()})
            ack_payload = {}

            def _capture_ack(payload: dict) -> None:
                ack_payload.update(payload)
                _log_creator_model_usage(
                    phase="generate.provider_ack",
                    skill_name=skill_name,
                    file_path=request.file_path,
                    route=route,
                    actual_model=payload.get("actual_model"),
                    provider=payload,
                )

            generated_chunks: list[str] = []
            async for chunk in stream_chat(prompt_messages, model, model_ack_callback=_capture_ack):
                if ack_payload:
                    yield _sse({"model_ack": {**route.ack(actual_model=ack_payload.get("actual_model")), "provider": ack_payload}})
                    ack_payload.clear()
                generated_chunks.append(chunk)

            raw_content = "".join(generated_chunks)
            logger.info(
                "[Creator][model] phase=generate.response skill=%s file=%s model=%s raw_chars=%d chunks=%d",
                skill_name,
                request.file_path,
                model,
                len(raw_content),
                len(generated_chunks),
            )

            should_repair = (
                request.file_path.startswith("scripts/")
                or request.file_path.startswith("references/")
                or request.file_path.startswith("assets/")
                or request.file_path == "SKILL.md"
            )
            if should_repair:
                content = raw_content
                last_error = ""
                repeated_error_counts: dict[str, int] = {}
                contract_text = _build_generated_file_contract_text(
                    request.file_path,
                    request.blueprint_text,
                    request.purpose,
                    role=request.role,
                    skill_plan_entry=request.skill_plan_entry,
                )
                for attempt in range(_MAX_FILE_REPAIR_ATTEMPTS + 1):
                    try:
                        content = _sanitize_generated_file_content(request.file_path, content, role=request.role, skill_plan_entry=request.skill_plan_entry)
                        if request.file_path == "SKILL.md":
                            _validate_skill_md_contract(content, request.blueprint_text)
                        elif request.file_path.startswith("references/"):
                            _validate_reference_file_contract(request.file_path, content, request.purpose)
                        elif request.file_path.startswith("assets/"):
                            _validate_asset_file_contract(request.file_path, content)
                        else:
                            _trial_run_generated_script_with_plan(
                                skill_name,
                                request.file_path,
                                content,
                                role=request.role,
                                skill_plan_entry=request.skill_plan_entry,
                            )
                        last_error = ""
                        break
                    except ValueError as validation_exc:
                        deterministic_error = str(validation_exc)
                        repeated_error_counts[deterministic_error] = repeated_error_counts.get(deterministic_error, 0) + 1
                        first_attempt_failed = attempt == 0
                        if request.file_path.startswith("scripts/") and first_attempt_failed:
                            repair_mode = "strict_contract_rewrite"
                        else:
                            repair_mode = "strict_contract_rewrite" if repeated_error_counts[deterministic_error] >= 2 else "minimal_edit"
                        contract_results = (
                            validation_exc.results
                            if isinstance(validation_exc, ContractValidationError)
                            else []
                        )
                        passed_checks_text = _format_contract_checks(contract_results, passed=True) if contract_results else ""
                        failed_checks_text = _format_contract_checks(contract_results, passed=False) if contract_results else ""
                        targeted_repair = _targeted_generated_file_repair_instructions(
                            file_path=request.file_path,
                            deterministic_error=deterministic_error,
                        )
                        validator_report = await _run_generated_file_validator_round(
                            file_path=request.file_path,
                            content=content,
                            deterministic_error=deterministic_error,
                            requested_model=model,
                            targeted_repair=targeted_repair,
                            contract_text=contract_text,
                            passed_checks_text=passed_checks_text,
                            failed_checks_text=failed_checks_text,
                            repair_mode=repair_mode,
                        )
                        last_error = _format_file_validator_feedback(deterministic_error, validator_report, targeted_repair)
                        last_validation_error = last_error
                        if attempt >= _MAX_FILE_REPAIR_ATTEMPTS:
                            repair_attempts = attempt
                            yield _sse({
                                "validation": {
                                    "status": "failed",
                                    "attempt": attempt,
                                    "error": deterministic_error,
                                    "validator": validator_report,
                                }
                            })
                            raise
                        repair_attempts = attempt + 1
                        logger.info(
                            "repairing generated file skill=%s file=%s attempt=%s error=%s validator_issues=%s",
                            skill_name,
                            request.file_path,
                            repair_attempts,
                            deterministic_error,
                            validator_report.get("issues"),
                        )
                        yield _sse({
                            "validation": {
                                "status": "repairing",
                                "attempt": repair_attempts,
                                "error": deterministic_error,
                                "validator": validator_report,
                            }
                        })
                        content = await _repair_generated_file_with_feedback(
                            prompt_messages=prompt_messages,
                            model=model,
                            file_path=request.file_path,
                            previous_content=content,
                            validation_error=last_error,
                            targeted_repair=targeted_repair,
                            contract_text=contract_text,
                            passed_checks_text=passed_checks_text,
                            failed_checks_text=failed_checks_text,
                            repair_mode=repair_mode,
                            skill_plan_entry=request.skill_plan_entry,
                        )
            else:
                content = _sanitize_generated_file_content(request.file_path, raw_content, role=request.role, skill_plan_entry=request.skill_plan_entry)

            yield _sse({"content": content})
            yield _sse({"done": True})
        except Exception as exc:
            logger.exception(
                "generate-file stream error skill=%s file=%s",
                skill_name,
                request.file_path,
            )
            # Return a safe user-facing message; full stack trace is in server logs.
            if last_validation_error:
                yield _sse({
                    "error": (
                        f"文件内容生成失败：已自动修复 {repair_attempts} 次仍未通过。"
                        f"最后错误：{last_validation_error}"
                    )
                })
            else:
                yield _sse({"error": "文件内容生成失败，请重试。详情已记录在服务器日志中。"})
        yield "data: [DONE]\n\n"

    return StreamingResponse(event_stream(), media_type="text/event-stream")


@router.post("/write-file", response_model=WriteFileResponse)
async def write_file(request: WriteFileRequest):
    """Write generated content to a Skill file on disk.

    Automatically strips spurious code-fence wrappers the LLM may add.
    """
    _validate_file_path(request.file_path)
    skill_name = _validate_skill_name(request.skill_name)

    try:
        content = _sanitize_generated_file_content(request.file_path, request.content, role=request.role, skill_plan_entry=request.skill_plan_entry)
        if request.file_path == "SKILL.md":
            _validate_skill_md_against_existing_files(skill_name, content)
        _validate_script_against_existing_skill_contract(skill_name, request.file_path, content)
        _trial_run_generated_script(
            skill_name,
            request.file_path,
            content,
            request.role,
            request.skill_plan_entry,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    if request.file_path == "SKILL.md":
        result = run_action({"action": "write", "name": skill_name, "content": content})
    else:
        p = Path(request.file_path)
        folder = p.parts[0]   # "scripts" / "references" / "assets"
        filename = p.name
        result = run_action(
            {
                "action": "write_file",
                "name": skill_name,
                "folder": folder,
                "filename": filename,
                "content": content,
            }
        )

    return WriteFileResponse(
        success=result["success"],
        path=result.get("path"),
        bytes=len(content.encode("utf-8")) if result["success"] else 0,
        message=result["message"],
    )



def _validate_skill_package_smoke(skill_name: str, *, mode: str = "trial") -> list[str]:
    """End-to-end dry-run for the final Skill package across docs and scripts.

    ``mode=trial`` uses deterministic helper behavior through the existing
    script trial runner so validation does not spend expensive model/image
    calls.  ``mode=sandbox`` currently exercises the same sandbox execution
    path with trial-run helpers enabled, while still validating the real docs,
    command templates, script paths, JSON argv, stdout outputs, and helper call
    wiring.
    """
    skill_dir = settings.skills_path / skill_name
    errors: list[str] = []
    skill_md_path = skill_dir / "SKILL.md"
    if not skill_md_path.is_file():
        return ["SKILL.md 合同错误：无法加载 SKILL.md。"]
    skill_md = skill_md_path.read_text(encoding="utf-8")
    try:
        _validate_skill_md_contract(skill_md, skill_md)
    except ValueError as exc:
        errors.append(f"SKILL.md 合同错误：{exc}")

    script_paths = sorted(set(_skill_local_paths_in_markdown(skill_md)) | {f"scripts/{p.name}" for p in (skill_dir / "scripts").glob("*.py")})
    script_paths = [path for path in script_paths if path.startswith("scripts/")]
    reference_paths = sorted(path for path in _skill_local_paths_in_markdown(skill_md) if path.startswith("references/"))

    for ref_path in reference_paths:
        ref_file = skill_dir / ref_path
        if not ref_file.is_file():
            errors.append(f"resource 不存在：{ref_path}")
            continue
        try:
            _validate_reference_file_contract(ref_path, ref_file.read_text(encoding="utf-8"), skill_md)
        except ValueError as exc:
            errors.append(f"reference 合同冲突：{ref_path}: {exc}")

    for script_path in script_paths:
        source_path = skill_dir / script_path
        if not source_path.is_file():
            errors.append(f"script 参数不匹配：找不到脚本 {script_path}")
            continue
        commands = _extract_script_command_templates(skill_md, script_path)
        for ref_path in reference_paths:
            ref_file = skill_dir / ref_path
            if ref_file.is_file():
                commands.extend(cmd for ref_script, cmd in _reference_script_commands(ref_file.read_text(encoding="utf-8")) if ref_script == script_path)
        if not commands:
            errors.append(f"SKILL.md 合同错误：无法识别 {script_path} 的命令块。")
        entry = _skill_plan_entry_for_file(file_path=script_path, blueprint_text=skill_md)
        failed_command_checks = [r for r in _check_command_block_contract(script_path, commands, entry) if not r.passed]
        if failed_command_checks:
            errors.append(f"command_template 不一致：{_format_contract_checks(failed_command_checks, passed=False)}")
        content = source_path.read_text(encoding="utf-8")
        failed_script_checks = [
            r for r in _check_script_file_contract(script_path, content, role=entry.role, skill_plan_entry=entry.__dict__)
            if not r.passed
        ]
        if failed_script_checks:
            errors.append(f"helper 未调用或调用失败：{script_path}: {_format_contract_checks(failed_script_checks, passed=False)}")
        try:
            _trial_run_generated_script(skill_name, script_path, content, entry.role, entry.__dict__)
        except ValueError as exc:
            msg = str(exc)
            if "stdout" in msg and "JSON" in msg:
                layer = "stdout JSON 不符合 outputs"
            elif "required_capabilities" in msg or "helper" in msg:
                layer = "helper 未调用或调用失败"
            else:
                layer = "sandbox 无法执行"
            errors.append(f"{layer}：{script_path}: {exc}")
    return errors

@router.post("/validate-skill", response_model=SkillActionResponse)
async def validate_skill(request: SkillActionRequest):
    """Validate SKILL.md format for a Skill package."""
    skill_name = _validate_skill_name(request.skill_name)
    result = run_action({"action": "validate", "name": skill_name})
    if result["success"]:
        skill_dir = settings.skills_path / skill_name
        trial_errors: list[str] = []
        for script_path in sorted((skill_dir / "scripts").glob("*.py")) if (skill_dir / "scripts").is_dir() else []:
            rel_path = f"scripts/{script_path.name}"
            try:
                _trial_run_generated_script(skill_name, rel_path, script_path.read_text(encoding="utf-8"))
            except ValueError as exc:
                trial_errors.append(f"{rel_path}: {exc}")
        if trial_errors:
            return SkillActionResponse(
                success=False,
                path=None,
                message="SKILL.md 格式校验通过，但脚本试运行失败：\n" + "\n\n".join(trial_errors),
            )
        smoke_errors = _validate_skill_package_smoke(skill_name, mode="trial")
        if smoke_errors:
            return SkillActionResponse(
                success=False,
                path=None,
                message="端到端 sandbox dry-run / smoke test 失败：\n" + "\n\n".join(smoke_errors),
            )
    return SkillActionResponse(
        success=result["success"],
        path=result.get("path"),
        message=result["message"],
    )


@router.post("/package-skill", response_model=SkillActionResponse)
async def package_skill(request: SkillActionRequest):
    """Package a Skill directory into a distributable .skill archive."""
    skill_name = _validate_skill_name(request.skill_name)
    result = run_action({"action": "package", "name": skill_name})
    return SkillActionResponse(
        success=result["success"],
        path=result.get("path"),
        message=result["message"],
    )


@router.post("/init-from-blueprint", response_model=InitFromBlueprintResponse)
async def init_from_blueprint(request: InitFromBlueprintRequest):
    """Initialize Skill directory structure from blueprint file list.
    
    Creates empty files based on the blueprint analysis result.
    This ensures the file structure matches exactly what the user confirmed.
    
    Workflow:
    1. Create main skill directory
    2. Create required subdirectories (scripts/, references/, assets/)
    3. Create empty files based on the blueprint file list
    """
    skill_name = _validate_skill_name(request.skill_name)
    skill_root = settings.skill_public_dir / skill_name
    
    try:
        # Create main skill directory
        skill_root.mkdir(parents=True, exist_ok=True)
        
        files_created = 0
        
        for file_spec in request.files:
            file_path = skill_root / file_spec.path
            
            # Ensure parent directory exists
            file_path.parent.mkdir(parents=True, exist_ok=True)
            
            # Create empty file if it doesn't exist
            if not file_path.exists():
                file_path.touch()
                files_created += 1
        
        return InitFromBlueprintResponse(
            success=True,
            path=str(skill_root),
            files_created=files_created,
            message=f"已创建 {files_created} 个文件",
        )
    
    except Exception as exc:
        logger.exception("init-from-blueprint error")
        return InitFromBlueprintResponse(
            success=False,
            path=None,
            files_created=0,
            message=f"初始化失败：{exc}",
        )


@router.post("/list-files", response_model=ListFilesResponse)
async def list_files(request: ListFilesRequest):
    """List all files in a Skill directory.
    
    Returns the actual file structure on disk, useful for displaying
    to the user after initializing the Skill directory structure.
    """
    skill_name = _validate_skill_name(request.skill_name)
    
    skill_root = settings.skill_public_dir / skill_name
    if not skill_root.exists():
        return ListFilesResponse(
            success=False,
            files=[],
            message=f"Skill '{skill_name}' 不存在",
        )
    
    files: list[FileInfo] = []
    
    def scan_dir(base: Path, rel_path: Path = Path("")):
        for entry in sorted(base.iterdir()):
            entry_rel = rel_path / entry.name
            if entry.is_dir():
                files.append(FileInfo(
                    path=str(entry_rel),
                    is_directory=True,
                ))
                scan_dir(entry, entry_rel)
            else:
                files.append(FileInfo(
                    path=str(entry_rel),
                    is_directory=False,
                    size=entry.stat().st_size,
                ))
    
    scan_dir(skill_root)
    
    return ListFilesResponse(
        success=True,
        files=files,
        message=f"已列出 {len(files)} 个文件",
    )
