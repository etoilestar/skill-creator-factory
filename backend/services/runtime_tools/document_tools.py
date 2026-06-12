"""Document runtime helpers for generated Skill scripts.

The helpers in this module intentionally import optional document libraries
inside each function.  That keeps ``backend.services.skill_runtime`` importable
in lightweight environments while still making Creator tool status aware that
these helpers exist and are platform-owned.
"""

from __future__ import annotations

import os
import re
from pathlib import Path
from typing import Any, Iterable

_SAFE_FILENAME_RE = re.compile(r"[^a-zA-Z0-9_.-]+")


def _safe_filename(filename: str, default: str) -> str:
    candidate = _SAFE_FILENAME_RE.sub("-", str(filename or default).strip()).strip("-._")
    return candidate or default


def _output_path(
    *,
    output_path: str | os.PathLike[str] | None,
    output_dir: str | os.PathLike[str] | None,
    filename: str,
) -> Path:
    """Resolve a document output path constrained to OUTPUT_DIR.

    Platform helpers must not write arbitrary host paths before the later
    artifact validator has a chance to reject unsafe stdout declarations.  The
    caller may provide a relative subpath, or an absolute path that is already
    inside OUTPUT_DIR, but traversal/absolute escapes are rejected up front.
    """
    out_dir = Path(output_dir or os.environ.get("OUTPUT_DIR") or "outputs").expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    if output_path:
        raw_path = Path(output_path).expanduser()
        candidate = raw_path.resolve() if raw_path.is_absolute() else (out_dir / raw_path).resolve()
        try:
            candidate.relative_to(out_dir)
        except ValueError as exc:
            raise ValueError("output_path must stay under OUTPUT_DIR") from exc
        candidate.parent.mkdir(parents=True, exist_ok=True)
        return candidate

    return out_dir / _safe_filename(filename, filename)


def _coerce_lines(text: str | Iterable[Any]) -> list[str]:
    if isinstance(text, str):
        lines = text.splitlines()
    else:
        lines = [str(item) for item in text]
    return [line if line else " " for line in lines] or ["Generated document"]



def _write_minimal_trial_pdf(path: Path) -> dict[str, Any]:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(b"%PDF-1.4\n1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n2 0 obj<</Type/Pages/Count 0>>endobj\ntrailer<</Root 1 0 R>>\n%%EOF\n")
    return {"pdf_path": str(path), "file_paths": [str(path)], "file_outputs": [str(path)]}

def create_pdf(
    text: str | Iterable[Any],
    *,
    output_path: str | os.PathLike[str] | None = None,
    output_dir: str | os.PathLike[str] | None = None,
    filename: str = "output.pdf",
    title: str | None = None,
    font_name: str = "STSong-Light",
) -> dict[str, Any]:
    """Create a Unicode-capable PDF and return JSON-serializable paths.

    ``STSong-Light`` is a built-in ReportLab CID font, so generated Chinese text
    works without bundling a TTF file.  Generated scripts should print the
    returned dict (or include its paths) as stdout JSON.
    """
    pdf_path = _output_path(output_path=output_path, output_dir=output_dir, filename=filename)
    if _trial():
        return _write_minimal_trial_pdf(pdf_path)

    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.pdfgen import canvas

    pdfmetrics.registerFont(UnicodeCIDFont(font_name))
    c = canvas.Canvas(str(pdf_path), pagesize=A4)
    _, height = A4
    c.setTitle(title or pdf_path.stem)
    c.setFont(font_name, 14)
    y = height - 72
    for raw_line in _coerce_lines(text):
        line = str(raw_line)
        while line:
            chunk, line = line[:42], line[42:]
            c.drawString(72, y, chunk)
            y -= 22
            if y < 72:
                c.showPage()
                c.setFont(font_name, 14)
                y = height - 72
    c.save()
    return {"pdf_path": str(pdf_path), "file_paths": [str(pdf_path)], "file_outputs": [str(pdf_path)]}


def create_docx(
    text: str | Iterable[Any],
    *,
    output_path: str | os.PathLike[str] | None = None,
    output_dir: str | os.PathLike[str] | None = None,
    filename: str = "output.docx",
    title: str | None = None,
) -> dict[str, Any]:
    """Create a Word document and return JSON-serializable paths."""
    docx_path = _output_path(output_path=output_path, output_dir=output_dir, filename=filename)

    from docx import Document

    document = Document()
    if title:
        document.add_heading(str(title), level=1)
    for line in _coerce_lines(text):
        document.add_paragraph(str(line))
    document.save(str(docx_path))
    return {"docx_path": str(docx_path), "file_paths": [str(docx_path)], "file_outputs": [str(docx_path)]}


def create_pptx(
    slides: str | Iterable[Any],
    *,
    output_path: str | os.PathLike[str] | None = None,
    output_dir: str | os.PathLike[str] | None = None,
    filename: str = "output.pptx",
    title: str = "Generated Presentation",
) -> dict[str, Any]:
    """Create a simple PowerPoint deck and return JSON-serializable paths."""
    pptx_path = _output_path(output_path=output_path, output_dir=output_dir, filename=filename)

    from pptx import Presentation

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = str(title or "Generated Presentation")
    title_slide.placeholders[1].text = "Created by Superskills runtime tools"

    for index, line in enumerate(_coerce_lines(slides), start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Slide {index}"
        slide.placeholders[1].text = str(line)
    prs.save(str(pptx_path))
    return {"pptx_path": str(pptx_path), "file_paths": [str(pptx_path)], "file_outputs": [str(pptx_path)]}


def extract_pdf_text(path: str | os.PathLike[str], *, max_pages: int | None = None) -> dict[str, Any]:
    """Extract text from a PDF with pypdf and return page-level text."""
    from pypdf import PdfReader

    pdf_path = Path(path).expanduser().resolve()
    reader = PdfReader(str(pdf_path))
    pages = []
    for index, page in enumerate(reader.pages):
        if max_pages is not None and index >= max_pages:
            break
        pages.append(page.extract_text() or "")
    return {
        "text": "\n".join(pages).strip(),
        "pages": pages,
        "page_count": len(reader.pages),
        "pdf_path": str(pdf_path),
    }

_MAX_INPUT_BYTES = 25 * 1024 * 1024
_MAX_INPUT_FILES = 50


def _trial() -> bool:
    return os.environ.get("SKILL_TRIAL_RUN") == "1"


def _allowed_input_roots() -> list[Path]:
    roots = [Path.cwd()]
    for name in ("SKILL_WORKDIR", "SKILL_DIR", "INPUT_DIR", "UPLOAD_DIR", "OUTPUT_DIR"):
        if os.environ.get(name):
            roots.append(Path(os.environ[name]))
    roots.extend([Path.cwd() / "inputs", Path.cwd() / "assets", Path.cwd() / "uploads"])
    return [root.expanduser().resolve() for root in roots]


def _safe_input_path(path: str | os.PathLike[str], suffixes: set[str]) -> Path:
    resolved = Path(path).expanduser().resolve()
    if resolved.suffix.lower() not in suffixes:
        raise ValueError(f"unsupported file type: {resolved.suffix}")
    if not resolved.is_file():
        raise FileNotFoundError("input file does not exist")
    if resolved.stat().st_size > _MAX_INPUT_BYTES:
        raise ValueError("input file is too large")
    if not any(resolved == root or resolved.is_relative_to(root) for root in _allowed_input_roots()):
        raise ValueError("input path must stay under the skill workdir, inputs, assets, uploads, or OUTPUT_DIR")
    return resolved


def read_docx_text(docx_path: str | os.PathLike[str]) -> dict[str, Any]:
    """Read text from a Word document."""
    if _trial():
        return {"text": "Mock DOCX text during SKILL_TRIAL_RUN.", "paragraphs": ["Mock DOCX text during SKILL_TRIAL_RUN."], "source_path": str(docx_path)}
    path = _safe_input_path(docx_path, {".docx"})
    from docx import Document

    document = Document(str(path))
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    return {"text": "\n".join(paragraphs), "paragraphs": paragraphs, "source_path": str(path)}


def read_pptx_text(pptx_path: str | os.PathLike[str]) -> dict[str, Any]:
    """Read text from a PowerPoint deck."""
    if _trial():
        return {"text": "Mock PPTX text during SKILL_TRIAL_RUN.", "paragraphs": ["Mock PPTX text during SKILL_TRIAL_RUN."], "source_path": str(pptx_path)}
    path = _safe_input_path(pptx_path, {".pptx"})
    from pptx import Presentation

    paragraphs: list[str] = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                paragraphs.append(str(shape.text))
    return {"text": "\n".join(paragraphs), "paragraphs": paragraphs, "source_path": str(path)}


def read_spreadsheet(path: str | os.PathLike[str], sheet_name: str | None = None, max_rows: int = 500) -> dict[str, Any]:
    """Read rows from an Excel spreadsheet."""
    max_rows = max(1, min(int(max_rows or 500), 5000))
    if _trial():
        return {"sheets": [sheet_name or "Sheet1"], "columns": ["A", "B"], "rows": [{"A": "mock", "B": "value"}], "row_count": 1, "truncated": False}
    safe_path = _safe_input_path(path, {".xlsx", ".xlsm"})
    from openpyxl import load_workbook

    workbook = load_workbook(str(safe_path), read_only=True, data_only=True)
    worksheet = workbook[sheet_name] if sheet_name else workbook[workbook.sheetnames[0]]
    rows_iter = worksheet.iter_rows(values_only=True)
    header_values = next(rows_iter, None) or []
    columns = [str(value) if value not in (None, "") else f"Column{index}" for index, value in enumerate(header_values, start=1)]
    rows: list[dict[str, Any]] = []
    truncated = False
    for index, values in enumerate(rows_iter, start=1):
        if index > max_rows:
            truncated = True
            break
        rows.append({columns[i] if i < len(columns) else f"Column{i+1}": value for i, value in enumerate(values)})
    return {"sheets": workbook.sheetnames, "columns": columns, "rows": rows, "row_count": len(rows), "truncated": truncated, "source_path": str(safe_path)}


def merge_pdfs(pdf_paths: list[str], output_path: str | os.PathLike[str] | None = None) -> dict[str, Any]:
    """Merge PDFs into a PDF under OUTPUT_DIR."""
    if not pdf_paths or len(pdf_paths) > _MAX_INPUT_FILES:
        raise ValueError("pdf_paths must contain 1 to 50 files")
    out_path = _output_path(output_path=output_path, output_dir=None, filename="merged.pdf")
    if _trial():
        return _write_minimal_trial_pdf(out_path)
    from pypdf import PdfReader, PdfWriter

    writer = PdfWriter()
    for raw_path in pdf_paths:
        reader = PdfReader(str(_safe_input_path(raw_path, {".pdf"})))
        for page in reader.pages:
            writer.add_page(page)
    with out_path.open("wb") as file_obj:
        writer.write(file_obj)
    return {"pdf_path": str(out_path), "file_paths": [str(out_path)], "file_outputs": [str(out_path)]}


def images_to_pdf(image_paths: list[str], output_path: str | os.PathLike[str] | None = None) -> dict[str, Any]:
    """Convert images to one PDF under OUTPUT_DIR."""
    if not image_paths or len(image_paths) > _MAX_INPUT_FILES:
        raise ValueError("image_paths must contain 1 to 50 files")
    out_path = _output_path(output_path=output_path, output_dir=None, filename="images.pdf")
    if _trial():
        return _write_minimal_trial_pdf(out_path)
    from PIL import Image

    images = []
    for raw_path in image_paths:
        image = Image.open(_safe_input_path(raw_path, {".png", ".jpg", ".jpeg", ".webp"})).convert("RGB")
        images.append(image)
    first, rest = images[0], images[1:]
    first.save(str(out_path), save_all=True, append_images=rest)
    return {"pdf_path": str(out_path), "file_paths": [str(out_path)], "file_outputs": [str(out_path)]}


def build_pdf_report(title: str, sections: list[dict], image_paths: list[str] | None = None, *, filename: str = "report.pdf") -> dict[str, Any]:
    """Build a simple PDF report from text sections and optional image references."""
    lines = [str(title or "Report"), ""]
    for section in sections or []:
        lines.append(str(section.get("title") or "Section"))
        lines.extend(_coerce_lines(str(section.get("text") or section.get("content") or "")))
        lines.append("")
    if image_paths:
        lines.append("Images:")
        lines.extend(str(_safe_input_path(path, {".png", ".jpg", ".jpeg", ".webp"})) for path in image_paths[:_MAX_INPUT_FILES])
    return create_pdf(lines, filename=filename or "report.pdf", title=title or "Report")
